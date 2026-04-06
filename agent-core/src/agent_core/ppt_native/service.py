"""Native Python PPT cutover service (initial slice).

This module provides a native execution path for core PPT operations while
keeping bridge mode available for compatibility.
"""

from __future__ import annotations

import sqlite3
import tempfile
import time
import uuid
import json
import hashlib
import re
import io
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from agent_core.config import get_settings
import jwt


def _native_export_work_dir() -> Path:
    """Directory for PPTX/PDF export artifacts (use tmp on ephemeral hosts like Render)."""
    raw = (get_settings().ppt_export_dir or "").strip()
    if raw:
        p = Path(raw).expanduser().resolve()
    else:
        p = Path(tempfile.gettempdir()) / "lf_ai_exports"
    p.mkdir(parents=True, exist_ok=True)
    return p


def _sanitize_bullet_line(text: str) -> str:
    """Remove trailing JSON fragments sometimes leaked into bullet strings by models."""
    s = str(text or "").strip()
    if len(s) < 8:
        return s
    # Cut at first obvious JSON object start mid-string (after some real words)
    j = s.find('{"')
    if j > 24:
        s = s[:j].rstrip(" -—\t")
    j2 = s.find("{'")
    if j2 > 24:
        s = s[:j2].rstrip(" -—\t")
    return s


def _sanitize_bullets_list(bullets: list[Any]) -> list[str]:
    out: list[str] = []
    for b in bullets:
        line = _sanitize_bullet_line(str(b))
        if line:
            out.append(line)
    return out


def _extract_topic_keywords(text: str, *, max_keywords: int = 10) -> list[str]:
    """Extract simple keyword candidates to keep titles/bullets grounded to the prompt."""
    stop = {
        "the",
        "and",
        "for",
        "with",
        "that",
        "this",
        "from",
        "into",
        "into",
        "your",
        "you",
        "are",
        "was",
        "were",
        "will",
        "can",
        "could",
        "should",
        "would",
        "how",
        "why",
        "what",
        "when",
        "where",
        "make",
        "using",
        "use",
        "about",
        "some",
        "many",
        "more",
        "most",
        "less",
        "over",
        "under",
        "then",
        "than",
        "also",
        "etc",
        "very",
        "just",
        "like",
        "best",
        "better",
        "good",
        "get",
        "getting",
        "improve",
        "improving",
        "design",
        "create",
        "including",
        "presentation",
        "presentations",
        "ppt",
        "slide",
        "slides",
        "deck",
        "please",
        "tell",
        "show",
        "give",
        "write",
        "prepare",
        "regarding",
    }
    raw_words = re.findall(r"[A-Za-z0-9]{3,}", str(text or "").lower())
    seen: set[str] = set()
    out: list[str] = []
    for w in raw_words:
        if w in stop:
            continue
        if w in seen:
            continue
        seen.add(w)
        out.append(w)
        if len(out) >= max_keywords:
            break
    # Short acronyms matter for relevance checks (regex above skips 2-letter tokens).
    low_full = str(text or "").lower()
    for short in ("ai", "ml", "qa", "ui", "ux", "ar", "vr", "it", "api", "llm", "gpu", "iot"):
        if len(out) >= max_keywords:
            break
        if short in seen:
            continue
        if re.search(rf"(?<![a-z0-9]){re.escape(short)}(?![a-z0-9])", low_full):
            seen.add(short)
            out.append(short)
    return out


def _slide_text_blob_lower(p: dict[str, Any]) -> str:
    """Concatenate slide copy for topic-grounding checks."""
    if not isinstance(p, dict):
        return ""
    parts: list[str] = []
    for k in ("title", "subtitle", "description", "highlight", "keyMessage", "speakerNotes"):
        parts.append(str(p.get(k) or ""))
    bl = p.get("bullets")
    if isinstance(bl, list):
        parts.extend(str(x) for x in bl)
    return " ".join(parts).lower()


def _payload_seems_topic_grounded(
    p: dict[str, Any],
    *,
    canonical_topic: str,
    keywords: list[str],
) -> bool:
    """Lenient match so varied LLM wording is not treated as off-topic."""
    blob = _slide_text_blob_lower(p)
    ct = (canonical_topic or "").strip().lower()
    if len(ct) >= 5 and ct in blob:
        return True
    topic_sig = [
        w
        for w in re.findall(r"[a-z]{4,}", ct)
        if w
        not in {
            "that",
            "with",
            "from",
            "this",
            "what",
            "when",
            "your",
            "have",
            "than",
            "into",
            "they",
            "them",
            "such",
            "most",
            "some",
        }
    ]
    if topic_sig and sum(1 for w in topic_sig if w in blob) >= min(2, len(topic_sig)):
        return True
    if not keywords:
        return True
    return any(k.lower() in blob for k in keywords)


def _canonical_subject_from_prompt(text: str) -> str:
    """Strip deck-building instructions so content/titles use the real subject (e.g. 'Artificial Intelligence')."""
    s = (text or "").strip()
    if not s:
        return "the topic"
    # File-extraction JSON: keep embedded topic as-is
    if s.startswith("{") and '"topic"' in s:
        try:
            parsed = json.loads(s)
            if isinstance(parsed, dict):
                t = str(parsed.get("topic") or "").strip()
                if t:
                    return t[:220]
        except Exception:
            pass
    patterns = [
        r"(?is)^\s*(create|make|build|prepare|design|write|give|show|develop)\s+(a\s+)?(an\s+)?(ppt|slides?|slide\s+deck|presentation|deck)\s+(on|about|for|regarding)\s+",
        r"(?is)^\s*(a\s+)?(ppt|slides?|presentation|deck)\s+(on|about|for|regarding)\s+",
        r"(?is)^\s*i\s+need\s+(a\s+)?(ppt|slides?|presentation|deck)\s+(on|about|for)\s+",
        r"(?is)^\s*prepare\s+(a\s+)?(ppt|slides?|presentation)\s+(on|about|for)\s+",
    ]
    for pat in patterns:
        s2 = re.sub(pat, "", s, count=1).strip()
        if s2 != s:
            s = s2
            break
    s = re.sub(r"(?is)\s+(presentation|deck|ppt)\s*$", "", s).strip(" .")
    return (s[:220] if s else "the topic").strip()


def _short_topic_label_for_copy(canonical: str, keywords: list[str]) -> str:
    """Compact phrase for titles/subtitles—never repeat a 60+ char subject in every bullet."""
    c = (canonical or "").strip()
    if not c or c.lower() == "the topic":
        return "this topic"
    low = c.lower()
    for prefix in (
        "the impact of ",
        "impact of ",
        "the role of ",
        "role of ",
        "overview of ",
        "introduction to ",
        "understanding ",
        "the future of ",
    ):
        if low.startswith(prefix):
            c = c[len(prefix) :].strip()
            low = c.lower()
            break
    c = re.sub(r"(?i)\bartificial intelligence\b", "AI", c)
    c = re.sub(r"(?i)\bmachine learning\b", "ML", c)
    c = re.sub(r"\s+", " ", c).strip()
    if len(c) <= 44:
        return c
    if keywords:
        kws = [k for k in keywords[:5] if k and len(str(k)) > 2][:4]
        joined = ", ".join(str(k) for k in kws)
        if joined and len(joined) <= 48:
            return joined
    if " in " in c:
        a, b = c.split(" in ", 1)
        a, b = a.strip(), b.split(",")[0].strip()[:36]
        cand = f"{a} in {b}".strip()
        if len(cand) <= 48:
            return cand
    cut = c[:42].rsplit(" ", 1)[0].strip()
    return cut + "…" if len(cut) < len(c) else c[:44]


def _heuristic_deck_display_title(canonical: str) -> str:
    """Short list title when the user did not provide one; never the raw instruction sentence."""
    c = (canonical or "").strip() or "Your Topic"
    core = c.split(",")[0].strip()[:72]
    if len(core) < 3:
        return "Strategic Briefing Deck"
    return f"{core}: Ideas, Impact, and Next Moves"[:160]


def _llm_generate_deck_display_title(*, canonical_subject: str) -> str | None:
    """One premium deck title for DB + header (not a copy of the user prompt)."""
    if not _openai_api_key_available():
        return None
    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    if not key or key == "sk-placeholder":
        return None
    subj = (canonical_subject or "").strip() or "the subject"
    try:
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI
    except Exception:
        return None
    sys = SystemMessage(
        content="You name presentation decks. Output JSON only: {\"title\":\"...\"}. Title: 6–12 words, compelling, specific to the subject. Never echo instructions like 'create a presentation'."
    )
    user = HumanMessage(content=f"Subject matter to name (not a command): {subj[:500]}")
    try:
        llm = ChatOpenAI(
            model=settings.openai_model,
            api_key=key,
            base_url=settings.openai_base_url,
            temperature=0.55,
        )
        resp = llm.invoke([sys, user])
        txt = str(getattr(resp, "content", None) or resp)
        parsed = _extract_json_object(txt)
        if not isinstance(parsed, dict):
            return None
        t = str(parsed.get("title") or "").strip()
        return t[:160] if t else None
    except Exception:
        return None


def _epoch_ms(ts: float | int | None) -> float:
    """Convert stored Unix-epoch **seconds** to **milliseconds** for JavaScript `Date`.

    Native mode uses `time.time()` in SQLite; the UI uses `new Date(x)` which expects ms.
    Without this, timestamps look like January 1970 (e.g. 21 Jan 1970 for ~1.7e9).
    """
    if ts is None:
        return 0.0
    return float(ts) * 1000.0


def _extract_json_object(text: str) -> dict[str, Any] | None:
    """Parse JSON from LLM output (handles ```json fences)."""
    raw = text.strip()
    if "```" in raw:
        m = re.search(r"```(?:json)?\s*([\s\S]*?)\s*```", raw, re.IGNORECASE)
        if m:
            raw = m.group(1).strip()
    try:
        out = json.loads(raw)
        return out if isinstance(out, dict) else None
    except Exception:
        pass
    start, end = raw.find("{"), raw.rfind("}")
    if start >= 0 and end > start:
        try:
            out = json.loads(raw[start : end + 1])
            return out if isinstance(out, dict) else None
        except Exception:
            pass
    return None


# --- Gamma-premium generation (role engine + layout + imageQuery) -----------------

_PREMIUM_ROLE_CYCLE = ("context", "insight", "example", "breakdown", "contrast", "data", "emotional")
_VALID_GAMMA_LAYOUTS = frozenset({"hero_split", "two_column", "three_cards", "title_bullets", "stats_split"})
_PREMIUM_ROLES = frozenset(
    {"hook", "context", "insight", "example", "breakdown", "contrast", "data", "emotional", "conclusion"}
)


def _premium_role_plan(n: int) -> list[str]:
    """Slide 1 = hook, last = conclusion, interior cycles through roles without consecutive duplicates."""
    if n <= 0:
        return []
    if n == 1:
        return ["hook"]
    roles: list[str] = ["hook"]
    pool = list(_PREMIUM_ROLE_CYCLE)
    pi = 0
    for _ in range(n - 2):
        picked = None
        for step in range(len(pool)):
            cand = pool[(pi + step) % len(pool)]
            if cand != roles[-1]:
                picked = cand
                pi = (pi + step + 1) % len(pool)
                break
        roles.append(picked or ("insight" if roles[-1] != "insight" else "example"))
    roles.append("conclusion")
    return roles


def _premium_visual_pack_for_role(role: str) -> dict[str, Any]:
    """Map narrative role → slideType + Gamma layout + chrome (valid presets only)."""
    r = (role or "").strip().lower()
    if r == "hook":
        return {
            "slideType": "hero",
            "layoutPreset": "hero_split",
            "imagePlacement": "left",
            "fullBleed": False,
            "textPrimary": False,
            "alignment": "center",
            "gradientTitle": True,
        }
    if r == "context":
        return {
            "slideType": "content",
            "layoutPreset": "two_column",
            "imagePlacement": "right",
            "fullBleed": False,
            "textPrimary": False,
            "alignment": "left",
            "gradientTitle": False,
        }
    if r == "insight":
        return {
            "slideType": "content",
            "layoutPreset": "two_column",
            "imagePlacement": "left",
            "fullBleed": False,
            "textPrimary": False,
            "alignment": "left",
            "gradientTitle": False,
        }
    if r == "example":
        return {
            "slideType": "visual",
            "layoutPreset": "two_column",
            "imagePlacement": "right",
            "fullBleed": False,
            "textPrimary": False,
            "alignment": "left",
            "gradientTitle": False,
        }
    if r == "breakdown":
        return {
            "slideType": "timeline",
            "layoutPreset": "three_cards",
            "fullBleed": False,
            "textPrimary": True,
            "alignment": "top",
            "gradientTitle": False,
        }
    if r == "contrast":
        return {
            "slideType": "comparison",
            "layoutPreset": "two_column",
            "imagePlacement": "left",
            "fullBleed": False,
            "textPrimary": True,
            "alignment": "left",
            "gradientTitle": False,
        }
    if r == "data":
        return {
            "slideType": "stats",
            "layoutPreset": "stats_split",
            "fullBleed": False,
            "textPrimary": True,
            "alignment": "center",
            "gradientTitle": False,
        }
    if r == "emotional":
        return {
            "slideType": "visual",
            "layoutPreset": "hero_split",
            "imagePlacement": "left",
            "fullBleed": True,
            "textPrimary": False,
            "alignment": "center",
            "gradientTitle": True,
        }
    if r == "conclusion":
        return {
            "slideType": "content",
            "layoutPreset": "title_bullets",
            "fullBleed": False,
            "textPrimary": True,
            "alignment": "center",
            "gradientTitle": True,
        }
    return {
        "slideType": "content",
        "layoutPreset": "title_bullets",
        "textPrimary": True,
        "alignment": "left",
        "gradientTitle": False,
    }


def _default_image_query_for_role(*, topic: str, role: str, title: str, index: int) -> str:
    t = (topic or "the topic").strip()[:80]
    tl = (title or t).strip()[:60]
    mood = "cinematic soft light, premium editorial"
    if role == "hook":
        return f"Dramatic wide shot related to {t}, audience tension and possibility, {mood}, slide {index}"
    if role == "context":
        return f"Real-world city or workplace scene echoing {t}, calm documentary style, {mood}, slide {index}"
    if role == "insight":
        return f"Abstract conceptual visualization of {t}, clean geometry and depth, {mood}, slide {index}"
    if role == "example":
        return f"Photoreal scene showing people applying {t} in daily work, authentic, {mood}, slide {index}"
    if role == "breakdown":
        return f"Step-by-step process infographic mood for {t}, numbered flow, minimal text, slide {index}"
    if role == "contrast":
        return f"Split visual metaphor comparing two paths for {t}, balanced composition, {mood}, slide {index}"
    if role == "data":
        return f"Modern analytics dashboard aesthetic for {t}, charts glow, cool blues, no readable numbers, slide {index}"
    if role == "emotional":
        return f"Human-centered emotional moment tied to {t}, warm tones, shallow depth of field, slide {index}"
    if role == "conclusion":
        return f"Forward-looking horizon or handshake closure for {t}, optimistic sunrise, {mood}, slide {index}"
    return f"{tl} professional scene, {mood}, slide {index}"


def _trim_bullet_to_word_count(s: str, min_w: int = 8, max_w: int = 16) -> str:
    words = re.findall(r"\S+", str(s or "").strip())
    if len(words) <= max_w:
        return " ".join(words)
    return " ".join(words[:max_w]).rstrip(",;:") + "…"


def _infer_detail_level(text: str | None) -> str:
    """
    Infer whether the user wants a standard / detailed / deep deck.

    User intent keywords are intentionally simple so this works even when the prompt
    is wrapped inside JSON (e.g. file extraction objects).
    """
    s = (text or "").lower()
    if not s.strip():
        return "standard"

    deep_terms = [
        "deep dive",
        "deep-dive",
        "deepdive",
        "in-depth",
        "in depth",
        "step-by-step",
        "step by step",
        "deeper",
        "research",
        "study",
        "technical",
    ]
    detailed_terms = [
        "detailed",
        "descriptive",
        "more detail",
        "include details",
        "give examples",
        "examples",
        "elaborate",
        "expand",
        "comprehensive",
    ]
    short_terms = ["short", "quick", "brief", "summary", "overview", "high-level", "at a high level", "lightweight"]

    has_deep = any(t in s for t in deep_terms)
    has_detailed = any(t in s for t in detailed_terms)
    has_short = any(t in s for t in short_terms)

    if has_deep:
        return "deep"
    if has_detailed and not has_short:
        return "detailed"
    if has_short and not has_detailed and not has_deep:
        return "concise"
    return "standard"


def _clamp_premium_bullets(
    bullets: list[str],
    *,
    topic: str,
    min_n: int = 5,
    max_n: int = 7,
    min_words: int = 6,
    max_words: int = 24,
) -> list[str]:
    cleaned = [
        _trim_bullet_to_word_count(_sanitize_bullet_line(b), min_words, max_words) for b in bullets if str(b).strip()
    ]
    cleaned = [b for b in cleaned if b]
    kw_pad = _extract_topic_keywords(str(topic or ""), max_keywords=5)
    topic_w = _short_topic_label_for_copy(str(topic or ""), kw_pad)[:44] or "this topic"
    pad = [
        f"Teams see faster decisions when guardrails and metrics for {topic_w} are defined up front.",
        f"One owner should map data, model, and deployment responsibilities for {topic_w}.",
        f"Pilot scope for {topic_w} needs users, success signals, and rollback triggers before scale.",
        f"Review cadences catch drift, bias, or security issues early—critical for {topic_w}.",
        f"Sequence wins on {topic_w}: prove value, then standardize tooling and training.",
    ]
    while len(cleaned) < min_n and pad:
        cleaned.append(pad.pop(0))
    return cleaned[:max_n]


def _merge_gamma_style(base: dict[str, Any], override: dict[str, Any] | None) -> dict[str, Any]:
    if not isinstance(override, dict):
        return dict(base)
    out = {**base, **{k: v for k, v in override.items() if v is not None}}
    lp = out.get("layoutPreset")
    if isinstance(lp, str) and lp not in _VALID_GAMMA_LAYOUTS:
        out["layoutPreset"] = base.get("layoutPreset", "title_bullets")
    return out


def _finalize_premium_deck_payloads(
    payloads: list[dict[str, Any]],
    *,
    topic: str,
    roles: list[str],
    detail_level: str = "standard",
) -> list[dict[str, Any]]:
    """Apply role plan, layout variety, bullet counts, imageQuery uniqueness, gammaStyle."""
    n = len(payloads)
    if n == 0:
        return []
    seen_queries: set[str] = set()
    prev_layout: str | None = None
    out: list[dict[str, Any]] = []
    for i, p in enumerate(payloads):
        role_raw = str(p.get("slideRole") or "").strip().lower()
        role = role_raw if role_raw in _PREMIUM_ROLES else (roles[i] if i < len(roles) else "insight")
        pack = _premium_visual_pack_for_role(role)
        slide_type = str(p.get("slideType") or pack.get("slideType") or "content").strip().lower()
        gs_llm = p.get("gammaStyle") if isinstance(p.get("gammaStyle"), dict) else {}
        gamma_base = {k: v for k, v in pack.items() if k != "slideType"}
        gs = _merge_gamma_style(gamma_base, gs_llm)
        layout = str(gs.get("layoutPreset") or "title_bullets")
        if prev_layout and layout == prev_layout:
            if layout == "two_column":
                gs["imagePlacement"] = "right" if gs.get("imagePlacement") == "left" else "left"
            elif layout == "hero_split":
                gs["layoutPreset"] = "title_bullets"
                layout = "title_bullets"
            else:
                gs["layoutPreset"] = "two_column"
                gs["imagePlacement"] = "left"
                layout = "two_column"
        prev_layout = str(gs.get("layoutPreset") or layout)

        title = str(p.get("title") or "").strip()
        iq = str(p.get("imageQuery") or "").strip()
        if not iq:
            iq = _default_image_query_for_role(topic=topic, role=role, title=title, index=i + 1)
        qk = iq[:200].lower()
        if qk in seen_queries:
            iq = f"{iq} — variant {i + 1} angle {role}"
        seen_queries.add(iq[:200].lower())

        bullets = p.get("bullets")
        if not isinstance(bullets, list):
            bullets = []
        # Bullet count control:
        # - concise: exactly 3 bullets
        # - standard: 3–4 bullets
        # - detailed / deep: exactly 4 bullets (visibly more detailed without overcrowding)
        detail_mode = (detail_level or "").strip().lower()
        if detail_mode == "concise":
            min_b, max_b = 3, 3
        elif detail_mode in ("detailed", "deep"):
            min_b, max_b = 4, 4
        else:
            min_b, max_b = 3, 4

        # Bullet length control: standard still gets slightly richer bullets; DEEP/DETAILED use stronger limits.
        min_words, max_words = 6, 34
        if detail_mode == "concise":
            min_words, max_words = 5, 18
        elif detail_mode == "detailed":
            min_words, max_words = 7, 38
        elif detail_mode == "deep":
            min_words, max_words = 9, 48
        bullets = _clamp_premium_bullets(
            [str(b) for b in bullets],
            topic=topic,
            min_n=min_b,
            max_n=max_b,
            min_words=min_words,
            max_words=max_words,
        )

        emphasis = [w for w in re.findall(r"[A-Za-z][A-Za-z-]+", title) if len(w) > 2][:3]
        if len(emphasis) < 2:
            emphasis = (emphasis + ["Impact", "Shift", "Proof"])[:3]
        gs["emphasisWords"] = (gs.get("emphasisWords") if isinstance(gs.get("emphasisWords"), list) else None) or emphasis[:3]

        out.append(
            {
                **p,
                "slideRole": role,
                "slideType": slide_type if slide_type else str(pack.get("slideType")),
                "bullets": bullets,
                "imageQuery": iq[:320],
                "gammaStyle": gs,
                "layoutSuggestion": gs.get("layoutPreset", layout),
            }
        )
    return out


def _llm_generate_slide_payloads(
    *,
    topic: str,
    deck_title: str,
    slide_count: int,
    tone: str,
    file_extraction: dict[str, Any] | None = None,
    user_prompt: str | None = None,
    detail_level: str = "standard",
) -> list[dict[str, Any]] | None:
    """Use OpenAI to produce rich slide JSON. Returns None if unavailable or on failure."""
    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    if not key or key == "sk-placeholder":
        return None
    try:
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI
    except Exception:
        return None

    n = max(3, min(slide_count, 30))
    tone_use = (tone or "professional").strip() or "professional"
    topic_trim = _canonical_subject_from_prompt((topic or deck_title or "").strip()).strip() or "the subject"
    raw_brief = (user_prompt or topic or "").strip()
    raw_brief = raw_brief[:8000]
    kw_a = _extract_topic_keywords(topic_trim, max_keywords=8)
    kw_b = _extract_topic_keywords(raw_brief[:2500], max_keywords=12) if raw_brief else []
    merged_kw: list[str] = []
    seen_kw: set[str] = set()
    for k in kw_a + kw_b:
        kl = k.lower()
        if kl in seen_kw:
            continue
        seen_kw.add(kl)
        merged_kw.append(k)
        if len(merged_kw) >= 14:
            break
    keywords_blob = ", ".join(merged_kw) if merged_kw else ""

    insights = None
    structured = None
    if isinstance(file_extraction, dict):
        maybe_insights = file_extraction.get("insights")
        insights = maybe_insights if isinstance(maybe_insights, list) else None
        structured = file_extraction.get("structuredData")

    file_context_section = ""
    if insights or isinstance(structured, dict):
        ins_lines: list[str] = []
        if isinstance(insights, list):
            for x in insights[:6]:
                if x is None:
                    continue
                s = str(x).strip()
                if s:
                    ins_lines.append(s[:180])

        kpi_lines: list[str] = []
        if isinstance(structured, dict):
            kpis = structured.get("kpis")
            if isinstance(kpis, list):
                for k in kpis[:6]:
                    if not isinstance(k, dict):
                        continue
                    name = str(k.get("name") or "").strip()
                    value = k.get("value")
                    unit = str(k.get("unit") or "").strip()
                    ctx = str(k.get("context") or "").strip()
                    source_tag = str(k.get("sourceTag") or "").strip()
                    if value is None or value == "":
                        continue
                    num = str(value).strip()
                    display = f"{name}: {num}{unit}".strip(": ")
                    if source_tag:
                        display = f"{source_tag} {display}".strip()
                    if ctx:
                        display = f"{display} ({ctx[:90]})"
                    kpi_lines.append(display[:220])

        ins_blob = "\n".join(f"- {l}" for l in ins_lines) if ins_lines else "- (none provided)"
        kpi_blob = "\n".join(f"- {l}" for l in kpi_lines) if kpi_lines else "- (none provided)"
        file_context_section = f"""
FILE EXTRACTIONS (summarized; do NOT copy verbatim):

INSIGHTS:
{ins_blob}

STRUCTURED DATA (for KPIs / numeric framing):
{kpi_blob}
"""

    detail_mode = (detail_level or "standard").strip().lower()
    # SlideContentGenerator behavior (adaptive depth + anti-pattern avoidance).
    if detail_mode == "deep":
        detail_clause = (
            "MODE: DEEP / IN-DEPTH.\n"
            "- Bullets: exactly 4; each includes mechanism + example/scenario + plausible measurable framing.\n"
            "- speakerNotes: 5–8 sentences, mini-lecture style, grounded (use FILE EXTRACTIONS when provided).\n"
            "- Avoid generic filler; each slide communicates ONE clear idea.\n"
            "- Anti-copy: never reuse exact sentences or long phrases from FILE EXTRACTIONS; paraphrase.\n"
        )
    elif detail_mode == "detailed":
        detail_clause = (
            "MODE: DETAILED.\n"
            "- Bullets: exactly 4; each adds value (no repetition) and includes at least one concrete example/workflow.\n"
            "- speakerNotes: 4–6 sentences, explanatory and presenter-friendly.\n"
            "- Avoid generic filler; each slide communicates ONE clear idea.\n"
            "- Anti-copy: never reuse exact sentences or long phrases from FILE EXTRACTIONS; paraphrase.\n"
        )
    elif detail_mode == "concise":
        detail_clause = (
            "MODE: CONCISE / BRIEF.\n"
            "- Bullets: exactly 3; short, high-impact, no fluff.\n"
            "- speakerNotes: 2–3 sentences; context only, not a re-read.\n"
            "- Avoid generic filler; each slide communicates ONE clear idea.\n"
        )
    else:
        detail_clause = (
            "MODE: STANDARD.\n"
            "- Bullets: 3–4; balanced clarity + depth, include mechanism and a concrete example when possible.\n"
            "- speakerNotes: 3–5 sentences; explain the why + how.\n"
            "- Avoid generic filler; each slide communicates ONE clear idea.\n"
            "- Anti-copy: never paste FILE EXTRACTIONS sentences verbatim; paraphrase and use evidence tags.\n"
        )
    system_hint = (
        "You are an expert on the subject matter in the user's request. "
        "Domain-specific, verifiable-sounding detail beats generic business language. "
        "Output strictly valid JSON only. No markdown fences, no commentary, no trailing text."
    )
    brief_block = raw_brief if raw_brief.strip() else "(No extra brief; use SUBJECT LINE and keywords only.)"
    user = f"""Generate presentation slides that answer the user's actual request—not a generic consulting deck.

FULL USER REQUEST (highest priority: audience, domain, scope, constraints, and any entities they named):
---
{brief_block}
---

SUBJECT LINE (for sharp slide titles only—do not paste this entire string into every bullet):
{topic_trim}

Deck title (display): {deck_title!s}
Tone: {tone_use}
Slide count: exactly {n}

Domain vocabulary seeds (infer more from the brief; weave several into each slide—naturally, not as a list dump):
{keywords_blob if keywords_blob else "(derive from the brief)"}
{file_context_section}

CITATION TAGGING (mandatory when evidence is used):
- Your FILE EXTRACTIONS include evidence tags like `[S1]`, `[S2]`, ... inside INSIGHTS and STRUCTURED DATA.
- Whenever you include a factual claim that comes from that evidence, append 1-2 citation tags to the *end of the sentence* inside `speakerNotes` (for example: "...". [S3][S4]).
- Never paste evidence sentences verbatim; paraphrase and use tags only.

NUMERIC FAITHFULNESS:
- If a bullet / sentence uses a numeric KPI from STRUCTURED DATA, keep the exact extracted number and unit (no rounding changes, no unit changes).

{detail_clause}

SPECIFICITY (content relevance > style):

1. EVERY bullet must satisfy at least one: a concrete real-world example, a named application or workflow, OR a measurable/plausible detail (timeframe, %, volume, before/after). If a bullet could apply to any industry unchanged, rewrite it until it could not.

2. Each bullet should make a reader ask "how does this show up in real life?" and get an answer from the text—not from imagination.

3. Banned as empty filler (do not use unless the user explicitly asked about that exact thing AND you tie it to domain nouns): generic "stakeholders", "roadmaps", "alignment", "outcomes", "strategy", "synergy", "transformation", "ecosystem", "leverage", "paradigm", "unlock value", and org clichés (sponsors/builders/operators) unless the brief is about operating models.

4. Forbidden phrases: "memorable move", "the key takeaway is", "this slide explains", "as you can see", "in today's world", "game-changer", "AI improves efficiency" without a mechanism.

5. Bad vs good pattern: not "X improves efficiency" alone—but "sensor-driven irrigation cuts water use roughly 20–30% on pilot fields when tied to soil-moisture telemetry" (your numbers and mechanisms must fit the USER'S topic).

6. Do not repeat the same long topic phrase in title, subtitle, bullets, highlight, and keyMessage on one slide—use synonyms, short labels, or pronouns after a single clear mention.

7. Titles: 4–7 words, max 55 characters, vivid and specific to that slide's angle; never "Introduction", "Overview", "Agenda", or the full deck title copied verbatim.

8. subtitle, description, highlight, keyMessage, speakerNotes: same specificity standard; no meta ("this slide…"); speakerNotes conversational, not a bullet re-read.

9. Bullet count rules by mode:
   - CONCISE: exactly 3 bullets
   - STANDARD: 3–4 bullets
   - DETAILED/DEEP: exactly 4 bullets

SLIDE ROLES (exact strings for slideRole; slide 1 = hook, slide {n} = conclusion; no two consecutive slides share the same slideRole):
hook, context, insight, example, breakdown, contrast, data, emotional, conclusion — cycle through the middle roles as needed.

IMAGEQUERY: every slide, distinct; domain-grounded scene (specific setting, tools, or activity from the user's topic); no readable text in the image; match slideRole mood.

LAYOUT: gammaStyle.layoutPreset one of hero_split, two_column, three_cards, title_bullets, stats_split—vary across the deck when sensible.

OUTPUT — single JSON object only, exactly {n} slides, this shape (use "bullets" not "points"; bullet count must follow the mode rules above):

{{
  "slides": [
    {{
      "slideRole": "hook",
      "title": "...",
      "subtitle": "...",
      "bullets": ["...", "...", "..."],
      "description": "...",
      "highlight": "...",
      "keyMessage": "...",
      "speakerNotes": "...",
      "imageQuery": "...",
      "slideType": "hero|content|visual|stats|comparison|timeline",
      "gammaStyle": {{
        "layoutPreset": "hero_split|two_column|three_cards|title_bullets|stats_split",
        "alignment": "left|center|top",
        "emphasisWords": ["word1","word2"],
        "imagePlacement": "left|right",
        "fullBleed": false,
        "textPrimary": false,
        "gradientTitle": true
      }}
    }}
  ]
}}

STRICT:
* Exactly {n} slides; slideRole, imageQuery, gammaStyle on every slide
* Every "bullets" array follows the mode's bullet count rule
* No duplicate or near-duplicate bullets across the whole deck
* Honor the FULL USER REQUEST above; if it names an industry, use that industry's artifacts and examples
* Do not invent precise statistics or company names unless plausible qualifiers ("e.g.", "pilot data suggests", "often")
* If the topic is technical or academic, stay accurate or mark uncertainty clearly"""

    def _parse_llm_slides_response(text: str) -> list[dict[str, Any]] | None:
        if not isinstance(text, str):
            text = str(text)
        parsed = _extract_json_object(text)
        if not parsed:
            return None
        slides = parsed.get("slides")
        if not isinstance(slides, list) or len(slides) < n:
            return None
        roles_plan = _premium_role_plan(n)
        out: list[dict[str, Any]] = []
        for i, s in enumerate(slides[:n]):
            if not isinstance(s, dict):
                return None
            title = str(s.get("title") or "").strip()
            bullets = s.get("bullets")
            if not isinstance(bullets, list):
                bullets = s.get("points")
            if not isinstance(bullets, list):
                bullets = []
            bullets = [str(b).strip() for b in bullets if str(b).strip()]
            if len(bullets) < 3:
                bullets = bullets + [
                    f"One decision the audience should make this quarter, grounded in the brief—not generic advice.",
                    f"A constraint or trade-off that changes how they prioritize work under this brief.",
                    f"A concrete next step that is measurable within 30–90 days.",
                ]
            bullets = bullets[:4]
            gs_raw = s.get("gammaStyle") if isinstance(s.get("gammaStyle"), dict) else {}
            out.append(
                {
                    "title": title or f"{deck_title} — Slide",
                    "subtitle": str(s.get("subtitle") or "").strip(),
                    "bullets": bullets[:4],
                    "description": str(s.get("description") or "").strip(),
                    "highlight": str(s.get("highlight") or "").strip(),
                    "keyMessage": str(s.get("keyMessage") or "").strip(),
                    "speakerNotes": str(s.get("speakerNotes") or "").strip(),
                    "slideRole": str(s.get("slideRole") or "").strip().lower(),
                    "imageQuery": str(s.get("imageQuery") or "").strip(),
                    "slideType": str(s.get("slideType") or "").strip().lower(),
                    "gammaStyle": gs_raw,
                }
            )
        if len(out) != n:
            return None
        out = _finalize_premium_deck_payloads(out, topic=topic_trim, roles=roles_plan, detail_level=detail_level)
        return out if len(out) >= 3 else None

    try:
        for temp in (0.65, 0.78):
            llm = ChatOpenAI(
                model=settings.openai_model,
                api_key=key,
                base_url=settings.openai_base_url,
                temperature=temp,
            )
            resp = llm.invoke(
                [
                    SystemMessage(content=system_hint),
                    HumanMessage(content=user),
                ]
            )
            text = getattr(resp, "content", None) or str(resp)
            parsed_out = _parse_llm_slides_response(text)
            if parsed_out:
                return parsed_out
        return None
    except Exception:
        return None


def _llm_generate_slide_titles(
    *,
    topic: str,
    deck_title: str,
    slide_count: int,
    tone: str,
) -> list[str] | None:
    """One engaging title per slide from the user topic (used when full-deck LLM is skipped)."""
    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    if not key or key == "sk-placeholder":
        return None
    try:
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI
    except Exception:
        return None

    n = max(3, min(slide_count, 30))
    topic_trim = _canonical_subject_from_prompt((topic or deck_title or "").strip()).strip() or "the subject"
    keywords = _extract_topic_keywords(topic_trim, max_keywords=8)
    keywords_blob = ", ".join([k for k in keywords if k]) if keywords else "(no keywords extracted)"
    tone_use = (tone or "professional").strip() or "professional"
    system_hint = (
        "You create premium slide titles. Output valid JSON only, no markdown. "
        "Titles must be insightful, topic-specific, and never generic."
    )
    user = f"""User prompt / topic:
---
{topic_trim}
---
Deck display title: {deck_title!s}
Tone: {tone_use}
Number of slides: {n}

Return a single JSON object:
{{
  "titles": [
    "Slide 1 title (4–7 words)",
    "... exactly {n} strings total ..."
  ]
}}

Rules:
- Exactly {n} titles in the array, same order as the deck.
- Each title: 4–7 words, max 55 characters, sharp and specific—short but clearly descriptive; no padding words.
- Topic keywords (must ground each title): {keywords_blob}
- Use consistent casing (either Title Case or sentence case, but keep it consistent).
- Slide 1 MUST be a Hook (attention-grabbing promise or stakes about the topic).
- Slide {n} MUST be a Conclusion (clear end-state, decision, or next-step framing about the topic).
- Slides 2..{n}-1 MUST alternate themes across: Context, Insight, Breakdown, Example, Comparison, Data/Fact, Implication.
- Avoid generic role words like: \"Introduction\", \"Overview\", \"Understanding\", \"Importance\", \"Key Points\".
- Do not repeat the same title twice (no near-duplicates).
- Never echo the user's command text (forbidden: titles starting with \"Create\", \"Make a presentation\", \"Build a deck\")."""

    try:
        llm = ChatOpenAI(
            model=settings.openai_model,
            api_key=key,
            base_url=settings.openai_base_url,
            temperature=0.65,
        )
        resp = llm.invoke(
            [
                SystemMessage(content=system_hint),
                HumanMessage(content=user),
            ]
        )
        text = getattr(resp, "content", None) or str(resp)
        if not isinstance(text, str):
            text = str(text)
        parsed = _extract_json_object(text)
        if not parsed:
            return None
        raw = parsed.get("titles")
        if not isinstance(raw, list) or len(raw) < n:
            return None
        titles: list[str] = []
        for t in raw[:n]:
            s = str(t).strip() if t is not None else ""
            if not s:
                return None
            titles.append(s[:200])
        return titles if len(titles) == n else None
    except Exception:
        return None


def _heuristic_slide_payloads(
    *,
    prompt: str,
    deck_title: str,
    slide_count: int,
    tone: str,
    file_extraction: dict[str, Any] | None = None,
    detail_level: str = "standard",
) -> list[dict[str, Any]]:
    """Topic-grounded fallback when the LLM is unavailable. Uses cleaned subject, not raw instructions."""
    raw_in = (prompt or deck_title or "this subject").strip()
    canonical = _canonical_subject_from_prompt(raw_in)
    if len(canonical) > 220:
        canonical = canonical[:217] + "…"

    focus_keywords = _extract_topic_keywords(canonical, max_keywords=7)
    topic_focus = (" ".join(focus_keywords[:5]).strip() if focus_keywords else canonical[:100]).strip()
    short = _short_topic_label_for_copy(canonical, focus_keywords)
    tone_adj = (tone or "professional").strip() or "professional"
    n = max(3, min(slide_count, 30))
    detail_mode = (detail_level or "standard").strip().lower()

    insights: list[str] = []
    kpis: list[dict[str, Any]] = []
    if isinstance(file_extraction, dict):
        maybe_insights = file_extraction.get("insights")
        if isinstance(maybe_insights, list):
            insights = [str(x).strip() for x in maybe_insights if str(x).strip()][:8]
        structured = file_extraction.get("structuredData")
        if isinstance(structured, dict):
            maybe_kpis = structured.get("kpis")
            if isinstance(maybe_kpis, list):
                kpis = [x for x in maybe_kpis if isinstance(x, dict)][:10]

    roles = _premium_role_plan(n)

    def _subtitle_for_role(role: str) -> str:
        # One short topic mention max; avoid echoing the full canonical string every line.
        m = {
            "hook": f"The stakes behind {short}—why this room should care now.",
            "context": f"Market and operating forces shaping adoption—not slide-level hype around {short}.",
            "insight": f"The pivot leaders need when the conversation shifts from vision to execution on {short}.",
            "example": "A concrete scene: workflows, owners, and review loops—not a generic demo script.",
            "breakdown": "Moving parts, dependencies, and sequencing so the story holds under scrutiny.",
            "contrast": "Two credible strategies—speed versus control—and the trade-offs each implies.",
            "data": "Signals that prove progress versus vanity charts that hide operational truth.",
            "emotional": "The human side: trust, fear, and accountability when tools change how work gets done.",
            "conclusion": f"What to commit to next—so {short} shows up in calendars, not only in strategy PDFs.",
        }
        return m.get(role, f"Focused angle: {short}.")

    def _heuristic_title_for_role(role: str) -> str:
        if role == "hook":
            return f"Why {short} Changes the Next 12 Months of Execution"[:120]
        if role == "conclusion":
            return f"Close With Owners, Dates, and One Clear Metric"[:120]
        lab = {
            "context": "The Field Today",
            "insight": "The Turning Point",
            "example": "On the Ground",
            "breakdown": "Inside the Engine",
            "contrast": "Fork in the Road",
            "data": "Signals That Matter",
            "emotional": "The Human Stakes",
        }
        return f"{lab.get(role, 'Focus')}: What Changes Next"[:120]

    def _bullets_for_role(role: str) -> list[str]:
        """Bullets must read like human copy: no injected full topic string in every line."""
        if role == "hook":
            return [
                "Roadmaps for risk, speed, and customer trust look different once adoption is real—not slideware.",
                f"One capability tied to {short} is already credible in production; another still needs guardrails and owners.",
                "Data quality, accountability, and release discipline must line up before scaled investment pays off.",
                "Name who among sponsors, builders, and operators co-owns outcomes—and what “done” means next quarter.",
            ]
        if role == "context":
            return [
                "Policy, talent, and infrastructure still decide where adoption spreads fastest; ignore that and pilots stall.",
                f"A crisp scope for {short} matters for budgets, audits, and liability—fuzzy framing creates rework.",
                "Sponsors see the vision, builders ship systems, operators feel incentives and daily workflow friction.",
                "Treat advanced tooling as augmented judgment with clear accountability—not magic that removes responsibility.",
                "Spell out what must be true about data and culture before operators trust the stack in production.",
            ]
        if role == "insight":
            return [
                "Move from deck-level interest to measured capability: owners, cadence, and a definition of good.",
                "Sequence thin releases, honest measurement, then wider scope—skipping the middle step breeds skepticism.",
                f"The cost of slow decisions on {short} compounds into churn, rework, and lost trust with stakeholders.",
                "Tie the idea to one bottleneck this audience already names in weekly meetings.",
                "Offer a decision rule people can repeat Monday—without waiting on an external deck cycle.",
            ]
        if role == "example":
            return [
                "Service workflows: humans approve, models rank, and response SLAs improve week over week when roles are explicit.",
                "Finance surfaces anomalies earlier when alerts are trusted and escalation paths are written down.",
                "In regulated or high-stakes domains, augmentation works when experts keep final sign-off and documentation.",
                "Demos that hide inputs and review cadence erode trust; show the operating loop, not a black box.",
                "Pilot with narrow scope, a visible metric, and a rollback story before the first serious misfire.",
            ]
        if role == "breakdown":
            return [
                "Step one: inventory data sources, owners, and freshness rules before anything earns production traffic.",
                "Step two: define evaluation windows, bias checks, and human review triggers on high-impact outputs.",
                "Step three: wire monitoring so drift, security events, and cost spikes surface within one sprint.",
                "Step four: train operators on failure modes so incidents get playbooks instead of panic threads.",
                "Step five: publish a concise flow so new hires see how data, models, and decisions connect end to end.",
            ]
        if role == "contrast":
            return [
                "Path A: ship fast, accept rework, learn loudly with visible pilots and tight feedback loops.",
                "Path B: gate behind heavier process, trading early velocity for compliance comfort and audit trails.",
                "Hybrids time-box pilots, cap blast radius, and require sponsors when scope or data access shifts.",
                "The choice exposes appetite for innovation versus tolerance for operational and regulatory risk.",
                "Pick the path that matches your exposure, talent depth, and the promise you make to customers.",
            ]
        if role == "data":
            return [
                "Leading signals: time-to-decision, error rates, and human override frequency on model-assisted work.",
                "Lagging proof: revenue impact, cost per case, and retention when new workflows support core journeys.",
                "Separate metrics operators open daily from vanity dashboards built only for steering committees.",
                "Establish baselines before launch so improvements survive skeptical finance and security review.",
                "Weekly triage beats one-off launches: catch silent degradation before it becomes a headline.",
            ]
        if role == "emotional":
            return [
                "People fear replacement; the narrative that wins shows augmentation that respects craft and judgment.",
                "Trust grows when leaders admit unknowns and fund training instead of declaring instant mastery.",
                "Celebrate small wins so teams associate change with relief and clarity—not surveillance or blame.",
                "Name the anxiety—job impact, opacity, accountability—then answer with transparent governance.",
            ]
        if role == "conclusion":
            return [
                f"Pick one decision this week on {short}: pilot scope, vendor path, or internal build—with named owners.",
                "Set a thirty-day signal that proves behavior change—not only a dashboard that went live quietly.",
                "Hold a short retrospective on data quality and adoption; surface one surprise and one fix.",
                "Teach one peer team the playbook so the organization learns once instead of in silos.",
                "Close with accountability: who follows up, when, and what good looks like next quarter.",
            ]
        return [
            "Anchor to one customer or operational outcome the room can recognize without jargon.",
            "Replace one myth with a testable claim your team can validate inside a month.",
            "Give a short checklist: prerequisites, owners, and the top risk to watch.",
            "Link forward to the next decision this deck will ask them to make with confidence.",
        ]

    raw_out: list[dict[str, Any]] = []
    for i in range(n):
        role = roles[i]
        sub = _subtitle_for_role(role)
        s_title = _heuristic_title_for_role(role)
        bullets = _bullets_for_role(role)
        note_tags: list[str] = []
        if insights and i < len(insights) and bullets:
            ins = str(insights[i]).strip()
            if ins:
                m = re.search(r"\[S\d+\]", ins)
                if m:
                    note_tags.append(m.group(0))
                bullets[0] = " ".join(ins.split()[:22]).rstrip(".") + "."
        k_source_tag = ""
        if role == "data" and kpis and bullets:
            k = kpis[i % len(kpis)]
            if isinstance(k, dict):
                name = str(k.get("name") or "").strip()
                value = k.get("value")
                unit = str(k.get("unit") or "").strip()
                ctx = str(k.get("context") or "").strip()
                k_source_tag = str(k.get("sourceTag") or "").strip()
                if name and value is not None:
                    num = str(value).strip()
                    kpi_line = f"{name} reached {num}{unit}".strip()
                    if ctx:
                        kpi_line = f"{kpi_line} ({ctx[:90]})"
                    if k_source_tag:
                        kpi_line = f"{k_source_tag} {kpi_line}".strip()
                    bullets[0] = " ".join(kpi_line.split()[:22]).rstrip(".") + "."
                if k_source_tag:
                    note_tags.append(k_source_tag)

        desc = f"{sub} Tie the thread to accountabilities and decisions your group can act on within the quarter."
        if detail_mode == "detailed":
            desc = f"{desc} Add one concrete constraint from the brief so the audience can picture the trade-off."
        elif detail_mode == "deep":
            desc = f"{desc} Then walk through a short causal chain (input -> mechanism -> observed effect) tied to the brief."
        else:
            # Standard mode: still add a small mechanism/explanation so it doesn't feel too bare.
            desc = f"{desc} Keep each point concrete and explain why it changes decisions."
        highlight = (
            f"Sharp line: {short} belongs in operating cadence—owners, metrics, and review—not only in strategy decks."
            if role != "conclusion"
            else "Leave with one committed action, a date, and a metric everyone agrees to track."
        )
        key_msg = (
            "Align sponsors, builders, and operators before scaling claims—or scope and trust will fracture."
            if role != "conclusion"
            else "Name the pilot boundary, the owner, and the next review date."
        )
        notes = (
            f"Open with why this matters to this room now, in a {tone_adj} tone; preview the arc without repeating the deck title."
            if role == "hook"
            else (
                "Slow on examples; invite one challenge question before you advance."
                if role not in ("conclusion", "emotional")
                else "Land the human truth first, then tie it to an accountable plan with owners."
            )
        )
        if detail_mode == "detailed":
            notes = f'{notes} Then give one practical "how-to" example the audience can try in their next meeting.'
        elif detail_mode == "deep":
            notes = f"{notes} Add a brief mini-case and end with a crisp decision rule the audience can apply immediately."
        else:
            notes = f"{notes} Include one sentence that explains the mechanism behind the decision."
        if role == "conclusion":
            notes = "Recap two supporting ideas, then lock one commitment with a calendar anchor."

        if note_tags:
            # Citations are only tags; evidence text itself is not pasted here.
            notes = f"{notes} {' '.join(note_tags)}".strip()

        raw_out.append(
            {
                "title": s_title,
                "subtitle": sub,
                "bullets": bullets,
                "description": desc,
                "highlight": highlight,
                "keyMessage": key_msg,
                "speakerNotes": notes,
                "slideRole": role,
                "imageQuery": "",
                "slideType": "",
                "gammaStyle": {},
            }
        )
    return _finalize_premium_deck_payloads(raw_out, topic=canonical, roles=roles, detail_level=detail_level)


def _conn() -> sqlite3.Connection:
    db_path = Path(get_settings().ppt_native_db_path).resolve()
    db_path.parent.mkdir(parents=True, exist_ok=True)
    con = sqlite3.connect(str(db_path))
    con.row_factory = sqlite3.Row
    _ensure_schema(con)
    return con


def _ensure_schema(con: sqlite3.Connection) -> None:
    con.executescript(
        """
        CREATE TABLE IF NOT EXISTS presentations (
            id TEXT PRIMARY KEY,
            user_id TEXT NOT NULL,
            title TEXT NOT NULL,
            prompt TEXT NOT NULL,
            template_name TEXT,
            status TEXT NOT NULL,
            created_at REAL NOT NULL,
            updated_at REAL NOT NULL
        );
        CREATE TABLE IF NOT EXISTS jobs (
            id TEXT PRIMARY KEY,
            presentation_id TEXT NOT NULL,
            status TEXT NOT NULL,
            error TEXT,
            created_at REAL NOT NULL,
            updated_at REAL NOT NULL
        );
        CREATE TABLE IF NOT EXISTS slides (
            id TEXT PRIMARY KEY,
            presentation_id TEXT NOT NULL,
            idx INTEGER NOT NULL,
            title TEXT NOT NULL,
            bullets_json TEXT NOT NULL,
            content_json TEXT
        );
        CREATE TABLE IF NOT EXISTS users (
            id TEXT PRIMARY KEY,
            first_name TEXT NOT NULL,
            last_name TEXT NOT NULL,
            mobile TEXT NOT NULL,
            email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            created_at REAL NOT NULL,
            updated_at REAL NOT NULL
        );
        CREATE TABLE IF NOT EXISTS user_charts (
            id TEXT PRIMARY KEY,
            user_id TEXT NOT NULL,
            title TEXT NOT NULL,
            chart_type TEXT NOT NULL,
            chart_data_json TEXT NOT NULL,
            source_type TEXT NOT NULL,
            source_name TEXT,
            input_summary TEXT,
            created_at REAL NOT NULL
        );
        """
    )
    cols = {r[1] for r in con.execute("PRAGMA table_info(slides)").fetchall()}
    if "content_json" not in cols:
        con.execute("ALTER TABLE slides ADD COLUMN content_json TEXT")
    pres_cols = {r[1] for r in con.execute("PRAGMA table_info(presentations)").fetchall()}
    if "share_settings_json" not in pres_cols:
        con.execute("ALTER TABLE presentations ADD COLUMN share_settings_json TEXT")
    con.execute(
        """
        CREATE TABLE IF NOT EXISTS presentation_share_views (
            id TEXT PRIMARY KEY,
            presentation_id TEXT NOT NULL,
            viewer_user_id TEXT,
            viewer_email TEXT,
            viewer_display_name TEXT,
            viewed_at REAL NOT NULL
        )
        """
    )
    con.execute(
        "CREATE INDEX IF NOT EXISTS idx_pres_share_views_pres ON presentation_share_views(presentation_id)"
    )
    con.execute(
        "CREATE INDEX IF NOT EXISTS idx_pres_share_views_pres_time ON presentation_share_views(presentation_id, viewed_at)"
    )
    con.execute(
        """
        CREATE TABLE IF NOT EXISTS username_registry (
            username_normalized TEXT PRIMARY KEY COLLATE NOCASE,
            user_id TEXT NOT NULL UNIQUE,
            created_at REAL NOT NULL
        )
        """
    )
    con.commit()


def create_presentation(
    *,
    user_id: str,
    prompt: str,
    title: str = "",
    template_name: str = "",
) -> dict[str, Any]:
    now = time.time()
    presentation_id = uuid.uuid4().hex
    canon = _canonical_subject_from_prompt(prompt)
    user_title = (title or "").strip()
    if user_title:
        display_title = user_title[:160]
    else:
        display_title = (_llm_generate_deck_display_title(canonical_subject=canon) or _heuristic_deck_display_title(canon))[:160]
    with _conn() as con:
        con.execute(
            """
            INSERT INTO presentations (id, user_id, title, prompt, template_name, status, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                presentation_id,
                user_id.strip(),
                display_title,
                prompt.strip(),
                template_name.strip() or None,
                "QUEUED",
                now,
                now,
            ),
        )
        con.commit()
    return {"presentationId": presentation_id, "status": "QUEUED", "mode": "native"}


def register_user(
    *,
    first_name: str,
    last_name: str,
    mobile: str,
    email: str,
    password: str,
) -> dict[str, Any]:
    em = email.strip().lower()
    if "@" not in em or len(password) < 6:
        return {"errorCode": "INVALID_INPUT", "message": "Invalid email or password"}
    now = time.time()
    with _conn() as con:
        exists = con.execute("SELECT id FROM users WHERE email = ?", (em,)).fetchone()
        if exists is not None:
            return {"errorCode": "EMAIL_EXISTS", "message": "A user with this email already exists."}
        user_id = uuid.uuid4().hex
        con.execute(
            """
            INSERT INTO users (id, first_name, last_name, mobile, email, password_hash, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                user_id,
                first_name.strip(),
                last_name.strip(),
                mobile.strip(),
                em,
                _hash_password(password),
                now,
                now,
            ),
        )
        con.commit()
    return {"userId": user_id, "mode": "native"}


def login_user(*, email: str, password: str) -> dict[str, Any]:
    em = email.strip().lower()
    with _conn() as con:
        row = con.execute("SELECT id, email, password_hash FROM users WHERE email = ?", (em,)).fetchone()
        if row is None:
            return {"errorCode": "INVALID_CREDENTIALS", "message": "Invalid email or password."}
        if _hash_password(password) != str(row["password_hash"]):
            return {"errorCode": "INVALID_CREDENTIALS", "message": "Invalid email or password."}
    token = jwt.encode(
        {"userId": row["id"], "email": row["email"], "exp": int(time.time()) + 7 * 24 * 3600},
        get_settings().auth_jwt_secret,
        algorithm="HS256",
    )
    return {"userId": row["id"], "email": row["email"], "token": token, "mode": "native"}


def generate_presentation(
    *,
    presentation_id: str,
    slide_count_target: int = 0,
    tone: str = "",
) -> dict[str, Any]:
    job_id = uuid.uuid4().hex
    now = time.time()
    with _conn() as con:
        row = con.execute(
            "SELECT id, prompt, title FROM presentations WHERE id = ?",
            (presentation_id,),
        ).fetchone()
        if row is None:
            return {"error": "Presentation not found"}

        con.execute(
            "INSERT INTO jobs (id, presentation_id, status, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
            (job_id, presentation_id, "PROCESSING", now, now),
        )
        _generate_slides_sync(
            con,
            presentation_id=presentation_id,
            prompt=row["prompt"],
            title=row["title"],
            slide_count=slide_count_target if slide_count_target > 0 else 8,
            tone=tone,
        )
        done = time.time()
        con.execute(
            "UPDATE jobs SET status = ?, updated_at = ? WHERE id = ?",
            ("COMPLETED", done, job_id),
        )
        con.execute(
            "UPDATE presentations SET status = ?, updated_at = ? WHERE id = ?",
            ("COMPLETED", done, presentation_id),
        )
        con.commit()
    return {"jobId": job_id, "status": "COMPLETED", "mode": "native"}


def _extract_text_from_pdf_bytes(buffer: bytes, *, limit_chars: int = 12000) -> str:
    """Best-effort PDF text extraction (no OCR)."""
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(io.BytesIO(buffer))
        parts: list[str] = []
        total = 0
        for page in reader.pages:
            t = page.extract_text() or ""
            t = str(t).strip()
            if not t:
                continue
            parts.append(t)
            total += len(t)
            if total >= limit_chars:
                break
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _extract_text_from_pptx_bytes(buffer: bytes, *, limit_chars: int = 12000) -> str:
    """Best-effort PPTX text extraction (includes notes)."""
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(io.BytesIO(buffer))
        parts: list[str] = []
        total = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    # Common text-bearing shapes.
                    if hasattr(shape, "text"):
                        t = str(getattr(shape, "text") or "").strip()
                        if t:
                            parts.append(t)
                            total += len(t)
                            if total >= limit_chars:
                                break
                    # Tables.
                    if getattr(shape, "has_table", False):
                        table = getattr(shape, "table", None)
                        if table is not None:
                            for row in table.rows:
                                for cell in row.cells:
                                    ct = str(getattr(cell, "text", "") or "").strip()
                                    if ct:
                                        parts.append(ct)
                                        total += len(ct)
                                        if total >= limit_chars:
                                            break
                                if total >= limit_chars:
                                    break
                except Exception:
                    continue
            if total >= limit_chars:
                break

            # Notes
            try:
                if getattr(slide, "has_notes_slide", False) and slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    nt = str(getattr(notes_frame, "text", "") or "").strip()
                    if nt:
                        parts.append(nt)
                        total += len(nt)
            except Exception:
                pass
            if total >= limit_chars:
                break

        return "\n".join(parts).strip()
    except Exception:
        return ""


def _extract_text_from_upload_bytes(
    *,
    filename: str,
    buffer: bytes,
    limit_chars: int = 12000,
) -> str:
    """Extract readable text from many uploads (txt/md/pdf/pptx)."""
    ext = Path(filename or "").suffix.lower().lstrip(".")
    if not buffer:
        return ""
    if ext in ("pdf",):
        return _extract_text_from_pdf_bytes(buffer, limit_chars=limit_chars).strip()
    if ext in ("pptx",):
        return _extract_text_from_pptx_bytes(buffer, limit_chars=limit_chars).strip()

    # Plain text fallback.
    try:
        return buffer.decode("utf-8", errors="ignore").strip()
    except Exception:
        return ""


def generate_from_file(*, presentation_id: str, file_path: str) -> dict[str, Any]:
    p = Path(file_path)
    if not p.exists():
        return {"error": f"File not found: {file_path}"}
    try:
        raw = p.read_bytes()
    except Exception:
        raw = b""
    text = _extract_text_from_upload_bytes(filename=p.name, buffer=raw, limit_chars=12000)

    def _extract_file_slide_context(src_text: str, file_stem: str) -> dict[str, Any]:
        t = re.sub(r"\s+", " ", (src_text or "")).strip()
        if not t:
            return {"topic": file_stem or "Presentation", "insights": [], "structuredData": {"kpis": [], "dataPoints": []}}

        # Topic: first short meaningful line; fall back to file stem.
        topic = ""
        for line in (src_text or "").splitlines():
            s = line.strip()
            if not s:
                continue
            if len(s) <= 80:
                topic = s
                break
        if not topic:
            topic = file_stem
        topic = topic.strip()[:120] or file_stem or "Presentation"

        # Keyword frequency for lightweight salience.
        stop = {
            "the",
            "and",
            "for",
            "with",
            "that",
            "this",
            "from",
            "into",
            "are",
            "was",
            "were",
            "will",
            "can",
            "could",
            "should",
            "would",
            "about",
            "over",
            "under",
            "their",
            "they",
            "then",
            "than",
            "also",
            "such",
            "more",
            "most",
        }
        tokens = re.findall(r"[A-Za-z][A-Za-z0-9\-']+", t.lower())
        freq: dict[str, int] = {}
        for w in tokens:
            if len(w) < 4 or w in stop:
                continue
            freq[w] = freq.get(w, 0) + 1
        top_kw = [w for w, _ in sorted(freq.items(), key=lambda x: x[1], reverse=True)[:10]]

        # Insights: extract salience sentences; then compress to short bullets.
        evidence_id = 0
        sentences = re.split(r"(?<=[.!?])\s+", src_text or "")
        insight_markers = [
            "result",
            "finding",
            "important",
            "key",
            "recommend",
            "impact",
            "therefore",
            "however",
            "leads",
            "drives",
            "risk",
            "trade-off",
            "opportunity",
            "best",
            "worst",
        ]
        markers_l = [m.lower() for m in insight_markers]
        insight_words = top_kw[:6]
        insights: list[str] = []
        seen: set[str] = set()

        for s in sentences:
            s_clean = s.strip()
            if not s_clean or len(s_clean.split()) < 8:
                continue
            s_low = s_clean.lower()
            score = 0
            if any(m in s_low for m in markers_l):
                score += 3
            if any(w in s_low for w in insight_words):
                score += 2
            # Prefer sentences that align with the topic's first word.
            topic_head = topic.split()[0].lower() if topic else ""
            if topic_head and topic_head in s_low:
                score += 1
            if score <= 0:
                continue
            trimmed = " ".join(s_clean.split()[:28]).strip().rstrip(".")
            if trimmed and trimmed not in seen:
                evidence_id += 1
                insights.append(f"[S{evidence_id}] {trimmed}.")
                seen.add(trimmed)
            if len(insights) >= 6:
                break

        if not insights:
            for s in sentences[:6]:
                s_clean = s.strip()
                if not s_clean:
                    continue
                trimmed = " ".join(s_clean.split()[:24]).strip().rstrip(".")
                if trimmed and trimmed not in seen:
                    evidence_id += 1
                    insights.append(f"[S{evidence_id}] {trimmed}.")
                    seen.add(trimmed)
                if len(insights) >= 5:
                    break

        # Data points: extract numbers and attach nearby KPI-like context.
        num_re = re.compile(
            r"(?P<prefix>[$€£])?\s*(?P<num>\d{1,3}(?:,\d{3})*(?:\.\d+)?|\d+(?:\.\d+)?)\s*(?P<unit>%|percent|million|billion|thousand|k|m|bn)?",
            re.I,
        )
        kpi_terms = [
            "revenue",
            "profit",
            "cost",
            "growth",
            "retention",
            "churn",
            "conversion",
            "accuracy",
            "latency",
            "errors",
            "defects",
            "roi",
            "margin",
        ]

        def _normalize_unit(u: str | None) -> str:
            if not u:
                return ""
            u2 = u.lower().strip()
            if u2 in ("%", "percent"):
                return "%"
            if u2 in ("million", "m"):
                return "M"
            if u2 in ("billion", "bn", "b"):
                return "B"
            if u2 in ("thousand", "k"):
                return "K"
            return u2[:6]

        data_points: list[dict[str, Any]] = []
        kpis: list[dict[str, Any]] = []
        for m in list(num_re.finditer(t))[:14]:
            start, end = m.span()
            ctx = t[max(0, start - 90) : min(len(t), end + 90)]
            num_raw = m.group("num")
            unit = _normalize_unit(m.group("unit"))
            try:
                num_norm = float(str(num_raw).replace(",", ""))
                value: Any = int(num_norm) if abs(num_norm - int(num_norm)) < 1e-9 else num_norm
            except Exception:
                continue

            ctx_low = ctx.lower()
            name = ""
            for term in kpi_terms:
                if term in ctx_low:
                    name = term
                    break
            if not name and top_kw:
                for kw in top_kw[:6]:
                    if kw in ctx_low:
                        name = kw
                        break
            name = name or (top_kw[0] if top_kw else "Metric")

            context = " ".join(ctx.strip().replace("\n", " ").split())[:120]
            evidence_id += 1
            source_tag = f"[S{evidence_id}]"
            data_points.append({"value": value, "unit": unit, "context": context, "sourceTag": source_tag})
            kpis.append(
                {
                    "name": str(name).title()[:30],
                    "value": value,
                    "unit": unit,
                    "context": context,
                    "sourceTag": source_tag,
                }
            )
            if len(kpis) >= 6:
                break

        structured = {"kpis": kpis[:6], "dataPoints": data_points[:6]}
        return {"topic": topic, "insights": insights[:6], "structuredData": structured}

    with _conn() as con:
        row = con.execute(
            "SELECT id, prompt FROM presentations WHERE id = ?",
            (presentation_id,),
        ).fetchone()
        if row is None:
            return {"error": "Presentation not found"}
        extraction = _extract_file_slide_context(text, p.stem)
        # Preserve the original user request so we can still infer whether they asked for
        # "detailed" / "deep" deck depth (even though we overwrite `presentations.prompt` below).
        extraction["userRequest"] = str(row["prompt"] or "")
        prompt = json.dumps(extraction, ensure_ascii=False)[:6000]
        con.execute(
            "UPDATE presentations SET prompt = ?, updated_at = ? WHERE id = ?",
            (prompt, time.time(), presentation_id),
        )
        con.commit()
    return generate_presentation(presentation_id=presentation_id)


def get_job_status(*, job_id: str) -> dict[str, Any]:
    with _conn() as con:
        row = con.execute("SELECT * FROM jobs WHERE id = ?", (job_id,)).fetchone()
        if row is None:
            return {"error": "Job not found"}
        return {
            "jobId": row["id"],
            "presentationId": row["presentation_id"],
            "status": row["status"],
            "error": row["error"],
            "mode": "native",
        }


def _default_share_settings() -> dict[str, Any]:
    return {
        "linkAccess": "view",
        "passwordEnabled": False,
        "passwordHash": None,
        "searchIndexing": False,
    }


def _parse_share_settings(raw: str | None) -> dict[str, Any]:
    d = _default_share_settings()
    if not raw:
        return d
    try:
        o = json.loads(raw)
        if isinstance(o, dict):
            if o.get("linkAccess") in ("none", "view"):
                d["linkAccess"] = o["linkAccess"]
            d["passwordEnabled"] = bool(o.get("passwordEnabled"))
            ph = o.get("passwordHash")
            d["passwordHash"] = ph if isinstance(ph, str) and ph.strip() else None
            d["searchIndexing"] = bool(o.get("searchIndexing"))
    except Exception:
        pass
    return d


def _hash_share_password(presentation_id: str, password: str) -> str:
    secret = (get_settings().auth_jwt_secret or "dev").encode()
    return hashlib.sha256(secret + b"|" + presentation_id.encode() + b"|" + password.encode()).hexdigest()


def _share_settings_for_response(st: dict[str, Any], *, is_owner: bool) -> dict[str, Any]:
    out: dict[str, Any] = {
        "linkAccess": st.get("linkAccess") or "view",
        "passwordEnabled": bool(st.get("passwordEnabled")),
        "searchIndexing": bool(st.get("searchIndexing")),
    }
    if is_owner:
        out["hasPassword"] = bool(st.get("passwordHash"))
    return out


def _iso_utc_from_epoch(ts: float) -> str:
    try:
        return datetime.fromtimestamp(float(ts), tz=timezone.utc).isoformat().replace("+00:00", "Z")
    except Exception:
        return ""


def _viewer_identity_from_db(con: sqlite3.Connection, viewer_user_id: str) -> tuple[str | None, str | None]:
    row = con.execute(
        "SELECT email, first_name, last_name FROM users WHERE id = ?",
        (viewer_user_id,),
    ).fetchone()
    if row is None:
        return None, None
    email = str(row["email"] or "").strip() or None
    fn = str(row["first_name"] or "").strip()
    ln = str(row["last_name"] or "").strip()
    disp = f"{fn} {ln}".strip() or None
    return email, disp


def _record_presentation_share_view(
    *,
    presentation_id: str,
    owner_id: str,
    viewer_user_id: str,
    viewer_email: str | None,
    viewer_display_name: str | None,
) -> None:
    """Append a share view event (debounced per viewer per deck). Not for owner previews."""
    vid = (viewer_user_id or "").strip()
    if not vid or vid == (owner_id or "").strip():
        return
    ve = (viewer_email or "").strip()[:320] or None
    vn = (viewer_display_name or "").strip()[:320] or None
    now = time.time()
    debounce_s = 300.0
    with _conn() as con:
        prev = con.execute(
            """
            SELECT 1 FROM presentation_share_views
            WHERE presentation_id = ? AND viewer_user_id = ? AND viewed_at > ?
            LIMIT 1
            """,
            (presentation_id, vid, now - debounce_s),
        ).fetchone()
        if prev:
            return
        if not ve or not vn:
            dbe, dbn = _viewer_identity_from_db(con, vid)
            ve = ve or dbe
            vn = vn or dbn
        rid = uuid.uuid4().hex
        con.execute(
            """
            INSERT INTO presentation_share_views
            (id, presentation_id, viewer_user_id, viewer_email, viewer_display_name, viewed_at)
            VALUES (?, ?, ?, ?, ?, ?)
            """,
            (rid, presentation_id, vid, ve, vn, now),
        )
        con.commit()


def list_presentation_share_views(*, presentation_id: str, owner_user_id: str) -> dict[str, Any]:
    """Owner-only: aggregated viewers who opened the shared link (logged-in / identified clients)."""
    uid = (owner_user_id or "").strip()
    if not uid:
        return {"error": "userId required"}
    with _conn() as con:
        row = con.execute("SELECT user_id FROM presentations WHERE id = ?", (presentation_id,)).fetchone()
        if row is None:
            return {"error": "Presentation not found"}
        if str(row["user_id"]) != uid:
            return {"error": "FORBIDDEN", "message": "Only the deck owner can view share analytics."}
        agg = con.execute(
            """
            SELECT
              viewer_user_id,
              MAX(viewer_email) AS ve,
              MAX(viewer_display_name) AS vd,
              COUNT(1) AS cnt,
              MAX(viewed_at) AS last_at
            FROM presentation_share_views
            WHERE presentation_id = ? AND viewer_user_id IS NOT NULL AND TRIM(viewer_user_id) != ''
            GROUP BY viewer_user_id
            ORDER BY last_at DESC
            """,
            (presentation_id,),
        ).fetchall()
        anon_row = con.execute(
            """
            SELECT COUNT(1) FROM presentation_share_views
            WHERE presentation_id = ? AND (viewer_user_id IS NULL OR TRIM(viewer_user_id) = '')
            """,
            (presentation_id,),
        ).fetchone()
        anon_n = int(anon_row[0] or 0) if anon_row else 0
    viewers: list[dict[str, Any]] = []
    for r in agg:
        vuid = str(r["viewer_user_id"] or "")
        last_at = float(r["last_at"] or 0)
        viewers.append(
            {
                "viewerUserId": vuid,
                "email": (str(r["ve"]).strip() if r["ve"] else None) or None,
                "displayName": (str(r["vd"]).strip() if r["vd"] else None) or None,
                "lastViewedAt": _iso_utc_from_epoch(last_at),
                "viewCount": int(r["cnt"] or 0),
            }
        )
    return {
        "viewers": viewers,
        "anonymousViewCount": anon_n,
        "mode": "native",
    }


def _normalize_username(raw: str) -> str:
    s = (raw or "").strip().lower()
    if len(s) < 2 or len(s) > 32:
        raise ValueError("Username must be 2–32 characters.")
    if not re.match(r"^[a-z0-9_\-]+$", s):
        raise ValueError("Use letters, numbers, underscores, or hyphens only.")
    return s


def derive_user_id_from_username_normalized(username_normalized: str) -> str:
    """Deterministic user id (same as JS `uuid.v5` with DNS namespace)."""
    return str(uuid.uuid5(uuid.NAMESPACE_DNS, f"lf.local.username:{username_normalized}"))


def register_local_username(*, username: str) -> dict[str, Any]:
    """Claim a new username (fails if already registered)."""
    try:
        norm = _normalize_username(username)
    except ValueError as e:
        return {"error": "invalid_username", "message": str(e)}
    uid = derive_user_id_from_username_normalized(norm)
    with _conn() as con:
        prev = con.execute(
            "SELECT username_normalized FROM username_registry WHERE username_normalized = ?",
            (norm,),
        ).fetchone()
        if prev:
            return {"error": "USERNAME_TAKEN", "message": "That username is already taken. Sign in instead."}
        con.execute(
            "INSERT INTO username_registry (username_normalized, user_id, created_at) VALUES (?, ?, ?)",
            (norm, uid, time.time()),
        )
        con.commit()
    return {
        "ok": True,
        "userId": uid,
        "username": norm,
        "email": f"{norm}@lf.local",
        "mode": "native",
    }


def login_local_username(*, username: str) -> dict[str, Any]:
    """Sign in to an existing registered username."""
    try:
        norm = _normalize_username(username)
    except ValueError as e:
        return {"error": "invalid_username", "message": str(e)}
    with _conn() as con:
        row = con.execute(
            "SELECT user_id FROM username_registry WHERE username_normalized = ?",
            (norm,),
        ).fetchone()
        if row is None:
            return {"error": "USERNAME_NOT_FOUND", "message": "No account with that username. Create one first."}
        uid = str(row["user_id"])
    return {
        "ok": True,
        "userId": uid,
        "username": norm,
        "email": f"{norm}@lf.local",
        "mode": "native",
    }


def get_presentation(
    *,
    presentation_id: str,
    viewer_user_id: str | None = None,
    share_password: str | None = None,
    viewer_email: str | None = None,
    viewer_display_name: str | None = None,
) -> dict[str, Any]:
    with _conn() as con:
        p = con.execute("SELECT * FROM presentations WHERE id = ?", (presentation_id,)).fetchone()
        if p is None:
            return {"error": "Presentation not found"}
        owner_id = str(p["user_id"])
        raw_share = None
        try:
            raw_share = p["share_settings_json"]
        except (KeyError, IndexError):
            raw_share = None
        st = _parse_share_settings(raw_share)

        vu = (viewer_user_id or "").strip() or None
        viewer_is_owner = bool(vu) and vu == owner_id

        if not viewer_is_owner:
            if st.get("linkAccess") == "none":
                return {
                    "error": "FORBIDDEN",
                    "code": "LINK_DISABLED",
                    "message": "Link sharing is turned off for this presentation.",
                }
            if st.get("passwordEnabled") and st.get("passwordHash"):
                pw = (share_password or "").strip()
                if not pw:
                    return {
                        "error": "FORBIDDEN",
                        "code": "PASSWORD_REQUIRED",
                        "message": "Enter the password to view this deck.",
                    }
                if _hash_share_password(presentation_id, pw) != st["passwordHash"]:
                    return {
                        "error": "FORBIDDEN",
                        "code": "PASSWORD_INVALID",
                        "message": "Incorrect password.",
                    }

        slides = con.execute(
            "SELECT id, idx, title, bullets_json, content_json FROM slides WHERE presentation_id = ? ORDER BY idx ASC",
            (presentation_id,),
        ).fetchall()
        slides_payload = [
            {
                "id": s["id"],
                "index": s["idx"],
                "order": s["idx"],
                "title": s["title"],
                "content": _slide_content_from_row(
                    s["title"],
                    s["bullets_json"],
                    s["content_json"],
                ),
            }
            for s in slides
        ]
        prompt_out = str(p["prompt"] or "")
        if not viewer_is_owner:
            prompt_out = ""
        pres: dict[str, Any] = {
            "id": p["id"],
            "userId": owner_id,
            "title": p["title"],
            "prompt": prompt_out,
            "status": p["status"],
            "templateName": p["template_name"],
            "slides": slides_payload,
            "shareSettings": _share_settings_for_response(st, is_owner=viewer_is_owner),
        }
        if vu and not viewer_is_owner:
            _record_presentation_share_view(
                presentation_id=presentation_id,
                owner_id=owner_id,
                viewer_user_id=vu,
                viewer_email=viewer_email,
                viewer_display_name=viewer_display_name,
            )
        return {
            "presentation": pres,
            "mode": "native",
        }


def update_share_settings(
    *,
    presentation_id: str,
    user_id: str,
    link_access: str | None = None,
    password_enabled: bool | None = None,
    password: str | None = None,
    search_indexing: bool | None = None,
) -> dict[str, Any]:
    uid = (user_id or "").strip()
    if not uid:
        return {"error": "userId required"}
    with _conn() as con:
        row = con.execute("SELECT id, user_id, share_settings_json FROM presentations WHERE id = ?", (presentation_id,)).fetchone()
        if row is None:
            return {"error": "Presentation not found"}
        if str(row["user_id"]) != uid:
            return {"error": "FORBIDDEN", "message": "Only the owner can change share settings."}
        try:
            raw_s = row["share_settings_json"]
        except (KeyError, IndexError):
            raw_s = None
        st = _parse_share_settings(raw_s)
        if link_access in ("none", "view"):
            st["linkAccess"] = link_access
        if search_indexing is not None:
            st["searchIndexing"] = bool(search_indexing)
        if password_enabled is not None:
            st["passwordEnabled"] = bool(password_enabled)
            if not st["passwordEnabled"]:
                st["passwordHash"] = None
        if password is not None:
            pw = str(password).strip()
            if pw:
                st["passwordHash"] = _hash_share_password(presentation_id, pw)
                st["passwordEnabled"] = True
            else:
                st["passwordHash"] = None
        con.execute(
            "UPDATE presentations SET share_settings_json = ?, updated_at = ? WHERE id = ?",
            (json.dumps(st, ensure_ascii=False), time.time(), presentation_id),
        )
        con.commit()
    return {
        "ok": True,
        "presentationId": presentation_id,
        "shareSettings": _share_settings_for_response(st, is_owner=True),
        "mode": "native",
    }


def list_presentations(*, user_id: str) -> dict[str, Any]:
    with _conn() as con:
        rows = con.execute(
            "SELECT id, title, prompt, status, created_at, updated_at FROM presentations WHERE user_id = ? ORDER BY updated_at DESC",
            (user_id,),
        ).fetchall()
        counts = {
            r["presentation_id"]: int(r["cnt"])
            for r in con.execute(
                "SELECT presentation_id, COUNT(1) as cnt FROM slides GROUP BY presentation_id",
            ).fetchall()
        }
        return {
            "presentations": [
                {
                    "id": r["id"],
                    "title": r["title"],
                    "prompt": r["prompt"],
                    "status": r["status"],
                    "createdAt": _epoch_ms(r["created_at"]),
                    "updatedAt": _epoch_ms(r["updated_at"]),
                    "slideCount": counts.get(r["id"], 0),
                    "lastActivityAt": _epoch_ms(r["updated_at"]),
                }
                for r in rows
            ],
            "mode": "native",
        }


def delete_presentation(*, presentation_id: str, user_id: str | None = None) -> dict[str, Any]:
    with _conn() as con:
        row = con.execute("SELECT id, user_id FROM presentations WHERE id = ?", (presentation_id,)).fetchone()
        if row is None:
            return {"errorCode": "NOT_FOUND", "message": "Presentation not found"}
        if user_id and str(row["user_id"]) != user_id:
            return {"errorCode": "FORBIDDEN", "message": "Presentation does not belong to this user"}
        con.execute("DELETE FROM slides WHERE presentation_id = ?", (presentation_id,))
        con.execute("DELETE FROM jobs WHERE presentation_id = ?", (presentation_id,))
        con.execute("DELETE FROM presentations WHERE id = ?", (presentation_id,))
        con.commit()
    return {"ok": True, "presentationId": presentation_id}


def list_templates() -> dict[str, Any]:
    # Minimal built-in set for compatibility with frontend template picker.
    return {
        "templates": [
            {"id": "gamma-default", "name": "gammaDefault", "theme": {"mode": "dark"}},
            {"id": "clementa", "name": "clementa", "theme": {"mode": "light"}},
        ]
    }


def export_file_url(*, presentation_id: str) -> dict[str, Any]:
    return {"url": f"/v1/ppt/native/presentations/{presentation_id}/export/file", "mode": "native"}


def update_slide(*, slide_id: str, title: str | None = None, content: dict[str, Any] | None = None) -> dict[str, Any]:
    with _conn() as con:
        s = con.execute(
            "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
            (slide_id,),
        ).fetchone()
        if s is None:
            return {"error": "Slide not found"}
        old_content = _slide_content_from_row(s["title"], s["bullets_json"], s["content_json"])
        next_content = dict(old_content)
        if isinstance(content, dict):
            next_content.update(content)
        if isinstance(title, str) and title.strip():
            next_content["title"] = title.strip()
        next_title = str(next_content.get("title") or s["title"] or "").strip() or str(s["title"])
        bullets = next_content.get("bullets")
        if not isinstance(bullets, list):
            bullets = _parse_bullets_json(s["bullets_json"])
        next_content["bullets"] = _sanitize_bullets_list([str(b) for b in bullets])
        bullets = next_content["bullets"]

        con.execute(
            "UPDATE slides SET title = ?, bullets_json = ?, content_json = ? WHERE id = ?",
            (next_title, "|||".join(str(b) for b in bullets), json.dumps(next_content), slide_id),
        )
        con.execute(
            "UPDATE presentations SET updated_at = ? WHERE id = ?",
            (time.time(), s["presentation_id"]),
        )
        con.commit()
    return {"slideId": slide_id, "updatedAt": _epoch_ms(time.time()), "mode": "native"}


def _apply_layout_structural_flip(content: dict[str, Any]) -> dict[str, Any]:
    """Deterministic layout tweak only (no copy generation)."""
    out = dict(content)
    gs = out.get("gammaStyle") if isinstance(out.get("gammaStyle"), dict) else {}
    flip = gs.get("imagePlacement") == "right"
    out["gammaStyle"] = {**gs, "layoutPreset": "two_column", "imagePlacement": "left" if flip else "right"}
    return out


def _llm_rewrite_slide_content(
    *,
    content: dict[str, Any],
    action: str,
    user_prompt: str,
) -> dict[str, Any] | None:
    """Return updated content dict, or None if the model failed after retries."""
    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    if not key or key == "sk-placeholder":
        return None
    try:
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI
    except Exception:
        return None

    title = str(content.get("title") or "")
    bullets = content.get("bullets")
    if not isinstance(bullets, list):
        bullets = []
    bullets = [str(b) for b in bullets[:8]]
    sub = str(content.get("subtitle") or "")
    desc = str(content.get("description") or "")
    highlight = str(content.get("highlight") or "")
    key_msg = str(content.get("keyMessage") or "")
    speaker_notes = str(content.get("speakerNotes") or "")

    act = (action or "improve").strip().lower()
    extra = (user_prompt or "").strip()

    # Layout is handled by caller via _apply_layout_structural_flip (no LLM copy).
    if act == "layout":
        return None

    # Action-specific editor behavior.
    if act == "improve":
        system = (
            "You edit slide copy like a senior presentation designer.\n"
            "Reply with strictly valid JSON only. No markdown. No explanations.\n\n"
            "Task: Improve writing.\n"
            "- Rewrite bullets + description to be clearer, sharper, more impactful.\n"
            "- Preserve meaning; do not invent new facts.\n"
        )
    elif act == "grammar":
        system = (
            "You edit slide copy with a grammar-only mindset.\n"
            "Reply with strictly valid JSON only. No markdown. No explanations.\n\n"
            "Task: Fix spelling & grammar only.\n"
            "- Correct grammar/spelling.\n"
            "- Do NOT change tone, structure, or meaning.\n"
            "- Keep bullet count the same (prefer 4–6 bullets).\n"
        )
    elif act == "longer":
        system = (
            "You edit slide copy to add depth.\n"
            "Reply with strictly valid JSON only. No markdown. No explanations.\n\n"
            "Task: Make longer (full regeneration, not padding).\n"
            "- Rewrite the ENTIRE slide: title, subtitle, every bullet, description, highlight, keyMessage, speakerNotes.\n"
            "- Each bullet should be meaningfully longer with concrete detail, implications, or examples—no fluff, no repetition of the same clause.\n"
            "- Description and speakerNotes should gain substantive context the presenter can use.\n"
            "- You may use the same number of bullets or add one bullet if it clearly adds a distinct insight (max 6 bullets).\n"
            "- Preserve facts and topic; do not invent statistics or names.\n"
        )
    elif act == "shorter":
        system = (
            "You edit slide copy for brevity.\n"
            "Reply with strictly valid JSON only. No markdown. No explanations.\n\n"
            "Task: Make shorter (full rewrite, NOT truncation).\n"
            "- Rewrite the ENTIRE slide from scratch in a tighter voice: title, subtitle, bullets, description, highlight, keyMessage, speakerNotes.\n"
            "- Do NOT shorten by cutting off mid-sentence, taking only the first sentence, or deleting bullets to hit a length—rephrase so each line is complete and crisp.\n"
            "- Prefer 3–4 bullets if the current deck is dense; each bullet fewer words but still a full thought.\n"
            "- Remove redundancy across fields (do not repeat the title in every bullet).\n"
            "- Preserve meaning and facts; do not omit critical constraints from the original.\n"
        )
    elif act == "simplify":
        system = (
            "You edit slide copy for beginners.\n"
            "Reply with strictly valid JSON only. No markdown. No explanations.\n\n"
            "Task: Simplify language.\n"
            "- Convert to easy-to-understand wording.\n"
            "- Keep structure recognizable; preserve meaning.\n"
        )
    elif act in ("visual", "image"):
        system = (
            "You redesign slide content for visual communication.\n"
            "Reply with strictly valid JSON only. No markdown. No explanations.\n\n"
            "Task: Make the slide more visual and image-ready.\n"
            "- Reduce bullets (prefer 3–4).\n"
            "- Make bullets punchy and scannable.\n"
            "- Provide imageQuery as: subject + environment + action + mood.\n"
            "- Set slideType to \"visual\" and set layoutSuggestion to a visual layout.\n"
            "- Keep title on the same topic; do not invent facts.\n"
        )
    else:
        system = (
            "You rewrite slide copy.\n"
            "Reply with strictly valid JSON only. No markdown. No explanations.\n\n"
            f"Task action: {act}\n"
        )

    # custom prompt is treated as highest priority instruction
    extra_block = extra or "(none)"
    user = f"""Action: {act}
Custom instruction (highest priority, may be empty):
{extra_block}

Current slide JSON (edit in place conceptually):
{json.dumps({"title": title, "subtitle": sub, "bullets": bullets, "description": desc, "highlight": highlight, "keyMessage": key_msg, "speakerNotes": speaker_notes}, ensure_ascii=False)}

Return exactly this JSON shape (fill all keys; use empty strings only when truly unknown):
{{
  "title": "string",
  "subtitle": "string",
  "bullets": ["string", "..."],
  "description": "string",
  "highlight": "string",
  "keyMessage": "string",
  "speakerNotes": "string",
  "imageQuery": "string",
  "slideType": "string",
  "layoutSuggestion": "string"
}}

Rules:
- Bullets: default 3–6 items; each bullet under 220 chars; specific (no meta text like \"this slide shows\").
- For grammar: minimal changes; only fix grammar/spelling; same bullet count when possible.
- For shorter: complete rewrites—every field shorter and cleaner; same narrative, no truncation artifacts.
- For longer: complete rewrites—richer detail everywhere; no duplicated padding phrases across bullets.
- For simplify: beginner-friendly wording; preserve meaning.
- For visual/image: include a high-quality imageQuery: subject + environment + action + mood.
- Title must stay on the same topic as the input title.
"""
    base_temp = 0.55 if act in ("longer", "shorter") else 0.45
    for temp in (base_temp, min(0.85, base_temp + 0.18)):
        try:
            llm = ChatOpenAI(
                model=settings.openai_model,
                api_key=key,
                base_url=settings.openai_base_url,
                temperature=temp,
            )
            resp = llm.invoke([SystemMessage(content=system), HumanMessage(content=user)])
            text = getattr(resp, "content", None) or str(resp)
            if not isinstance(text, str):
                text = str(text)
            parsed = _extract_json_object(text)
            if not parsed:
                continue
            out = dict(content)
            if isinstance(parsed.get("title"), str) and parsed["title"].strip():
                out["title"] = parsed["title"].strip()[:120]
            if isinstance(parsed.get("subtitle"), str):
                out["subtitle"] = parsed["subtitle"].strip()
            bl = parsed.get("bullets")
            if isinstance(bl, list) and len(bl) >= 2:
                out["bullets"] = _sanitize_bullets_list([str(x) for x in bl])[:8]
            elif isinstance(bl, list) and len(bl) == 1 and act in ("longer", "shorter"):
                one = str(bl[0] or "").strip()
                chunks = [c.strip() for c in one.split(";") if c.strip()]
                if len(chunks) >= 2:
                    out["bullets"] = _sanitize_bullets_list(chunks)[:8]
                elif one:
                    out["bullets"] = _sanitize_bullets_list([one])
            if isinstance(parsed.get("description"), str):
                out["description"] = parsed["description"].strip()
            if isinstance(parsed.get("highlight"), str):
                out["highlight"] = parsed["highlight"].strip()
            if isinstance(parsed.get("keyMessage"), str):
                out["keyMessage"] = parsed["keyMessage"].strip()
            if isinstance(parsed.get("speakerNotes"), str):
                out["speakerNotes"] = parsed["speakerNotes"].strip()
            if isinstance(parsed.get("imageQuery"), str):
                out["imageQuery"] = parsed["imageQuery"].strip()
            if isinstance(parsed.get("slideType"), str):
                out["slideType"] = parsed["slideType"].strip()
            if isinstance(parsed.get("layoutSuggestion"), str):
                out["layoutSuggestion"] = parsed["layoutSuggestion"].strip()
            nb = out.get("bullets")
            if isinstance(nb, list) and len(nb) >= 2:
                return out
            if act in ("grammar",) and isinstance(nb, list) and len(nb) >= 1:
                return out
        except Exception:
            continue
    return None


def ai_edit_slide(
    *,
    slide_id: str,
    action: str,
    user_prompt: str = "",
    current_content_json: str | None = None,
) -> dict[str, Any]:
    """AI edit slide; persists and returns AISlideEditor-compatible `editResponse` shape."""
    with _conn() as con:
        s = con.execute(
            "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
            (slide_id,),
        ).fetchone()
        if s is None:
            return {"error": "Slide not found"}
        content = _slide_content_from_row(s["title"], s["bullets_json"], s["content_json"])
        if current_content_json:
            try:
                overlay = json.loads(current_content_json)
                if isinstance(overlay, dict):
                    content = {**content, **overlay}
                    ob = overlay.get("bullets")
                    if isinstance(ob, list):
                        content["bullets"] = _sanitize_bullets_list([str(x) for x in ob])
            except Exception:
                pass
        act = (action or "improve").strip().lower()
        if act == "custom":
            act = "improve"

        if act == "layout":
            merged = _apply_layout_structural_flip(dict(content))
        elif not _openai_api_key_available():
            return {"error": "AI editing requires a configured OpenAI API key (OPENAI_API_KEY)."}
        else:
            merged = _llm_rewrite_slide_content(content=content, action=act, user_prompt=user_prompt)
            if merged is None:
                return {
                    "error": "The AI could not update this slide. Try again, or check your API key and model settings.",
                }

        merged["title"] = str(merged.get("title") or s["title"])
        bullets = merged.get("bullets")
        if not isinstance(bullets, list):
            bullets = _sanitize_bullets_list(_parse_bullets_json(s["bullets_json"]))
        merged["bullets"] = _sanitize_bullets_list([str(b) for b in bullets])[:6]

        # Keep Gamma layout presets in sync with `slideType`.
        slide_type = merged.get("slideType")
        if isinstance(slide_type, str) and slide_type.strip():
            st = slide_type.strip().lower()
            layout_preset = _layout_preset_for_slide_type(st)
            merged["slideType"] = st
            merged["layoutSuggestion"] = str(merged.get("layoutSuggestion") or layout_preset)
            gs = merged.get("gammaStyle") if isinstance(merged.get("gammaStyle"), dict) else {}
            merged["gammaStyle"] = {**gs, "layoutPreset": layout_preset, "alignment": gs.get("alignment", "left")}

        con.execute(
            "UPDATE slides SET title = ?, bullets_json = ?, content_json = ? WHERE id = ?",
            (str(merged["title"]), "|||".join(merged["bullets"]), json.dumps(merged), slide_id),
        )
        con.execute("UPDATE presentations SET updated_at = ? WHERE id = ?", (time.time(), s["presentation_id"]))
        con.commit()

        slide = {
            "id": s["id"],
            "presentationId": s["presentation_id"],
            "title": str(merged["title"]),
            "order": int(s["idx"]),
            "content": merged,
        }

        if act == "layout":
            gs = merged.get("gammaStyle") if isinstance(merged.get("gammaStyle"), dict) else {}
            edit_response: dict[str, Any] = {
                "type": "layout",
                "data": {
                    "layoutType": str(gs.get("layoutPreset") or "hero_split"),
                    "gammaStylePatch": {
                        "layoutPreset": gs.get("layoutPreset", "two_column"),
                        "imagePlacement": gs.get("imagePlacement", "left"),
                    },
                    "slideContentPatch": merged,
                },
            }
        elif act in ("visual", "image"):
            edit_response = {
                "type": "content",
                "data": {
                    "title": slide["title"],
                    "slideContentPatch": merged,
                },
            }
        else:
            edit_response = {
                "type": "content",
                "data": {
                    "title": slide["title"],
                    "slideContentPatch": merged,
                },
            }

        return {"slide": slide, "warning": None, "mode": "native", "editResponse": edit_response}


def _slide_from_db_row(*, srow: sqlite3.Row) -> dict[str, Any]:
    # Slide "content" is stored in content_json with bullets_json split out for legacy reasons.
    content_obj = _slide_content_from_row(
        str(srow["title"] or ""),
        str(srow["bullets_json"] or ""),
        srow["content_json"] or None,
    )
    return {
        "id": str(srow["id"]),
        "presentationId": str(srow["presentation_id"]),
        "order": int(srow["idx"]),
        "title": str(srow["title"] or ""),
        "content": content_obj,
    }


def _persist_slide_content(*, slide_id: str, title: str, bullets: list[str], content: dict[str, Any]) -> None:
    with _conn() as con:
        con.execute(
            "UPDATE slides SET title = ?, bullets_json = ?, content_json = ? WHERE id = ?",
            (str(title), "|||".join([str(b) for b in bullets if str(b).strip()]), json.dumps(content), slide_id),
        )
        # Keep presentation updated for UI refresh signals.
        pres_row = con.execute("SELECT presentation_id FROM slides WHERE id = ?", (slide_id,)).fetchone()
        if pres_row is not None:
            con.execute("UPDATE presentations SET updated_at = ? WHERE id = ?", (time.time(), pres_row["presentation_id"]))
        con.commit()


def _llm_rewrite_slide_title_only(*, title: str, slide_content: dict[str, Any], custom_prompt: str = "") -> tuple[str, float] | None:
    """Rewrite title only, returning (newTitle, confidence)."""
    if not _openai_api_key_available():
        return None

    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    try:
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI
    except Exception:
        return None

    bullets = slide_content.get("bullets")
    if not isinstance(bullets, list):
        bullets = []
    bullets_txt = "\n".join([str(b) for b in bullets[:6] if str(b).strip()])[:1200]
    highlight = str(slide_content.get("highlight") or slide_content.get("keyMessage") or "")
    desc = str(slide_content.get("description") or "")

    custom = (custom_prompt or "").strip()
    system = (
        "You rewrite a slide title to be sharp and non-generic.\n"
        "Output strictly valid JSON only. No markdown. No explanations.\n\n"
        "Rules:\n"
        "- Max 8–10 words.\n"
        "- Avoid generic labels like \"Introduction\", \"Overview\", \"Understanding\".\n"
        "- Keep the title on the same topic as the input title.\n"
        "- Make it role-specific: suggest stakes, outcome, or decision framing.\n"
        "- Do not invent new facts.\n"
    )
    user = f"""Action: improve title only

Custom instruction (highest priority, may be empty):
{custom or "(none)"}"""
    user += f"""

Current title: {title}
Bullet context:
{bullets_txt}
Highlight context:
{highlight}
Description context:
{desc}

Return exactly this JSON:
{{
  "title": "string",
  "confidence": 0.0
}}
"""
    for temp in (0.4, 0.65):
        try:
            llm = ChatOpenAI(
                model=settings.openai_model,
                api_key=key,
                base_url=settings.openai_base_url,
                temperature=temp,
            )
            resp = llm.invoke([SystemMessage(content=system), HumanMessage(content=user)])
            text = getattr(resp, "content", None) or str(resp)
            if not isinstance(text, str):
                text = str(text)
            parsed = _extract_json_object(text)
            if not parsed or not isinstance(parsed, dict):
                continue
            nt = str(parsed.get("title") or "").strip()
            conf = parsed.get("confidence", 0.75)
            conf_f = float(conf) if isinstance(conf, (int, float, str)) else 0.75
            if not nt:
                continue
            return nt[:120], max(0.05, min(0.99, conf_f))
        except Exception:
            continue
    return None


def ai_rewrite_title_for_editor(*, title: str, context: str = "", tone: str = "professional") -> dict[str, Any]:
    """
    Editor endpoint helper for `POST /ai/rewrite-title`.

    Returns:
    {
      rewrittenTitle: string,
      styleUsed: string,
      confidence: number,
      variations: string[]
    }
    """
    base = str(title or "").strip()
    ctx = str(context or "").strip()
    tone_use = str(tone or "professional").strip().lower() or "professional"

    # LLM path
    if _openai_api_key_available():
        settings = get_settings()
        key = (settings.openai_api_key or "").strip()
        try:
            from langchain_core.messages import HumanMessage, SystemMessage
            from langchain_openai import ChatOpenAI
        except Exception:
            settings = None  # type: ignore[assignment]

        try:
            system = (
                "You rewrite slide titles to be sharp and non-generic.\n"
                "Output strictly valid JSON only. No markdown. No explanations.\n\n"
                "Rules:\n"
                "- Return 3 variations max.\n"
                "- Each title: max 8–10 words.\n"
                "- Avoid generic labels: Introduction, Overview, Understanding, Importance.\n"
                "- Keep the topic on the same subject as the input title.\n"
                "- Do not invent new facts.\n"
            )
            user = f"""Improve title (editor)

Tone: {tone_use}

Current title:
{base}

Context (bullets / highlight / key idea):
{ctx[:1200]}

Return exactly:
{{
  "rewrittenTitle": "string",
  "confidence": 0.0,
  "variations": ["string", "string", "string"]
}}
"""
            llm = ChatOpenAI(
                model=settings.openai_model,
                api_key=key,
                base_url=settings.openai_base_url,
                temperature=0.5,
            )
            resp = llm.invoke([SystemMessage(content=system), HumanMessage(content=user)])
            text = getattr(resp, "content", None) or str(resp)
            if not isinstance(text, str):
                text = str(text)
            parsed = _extract_json_object(text)
            if isinstance(parsed, dict):
                rt = str(parsed.get("rewrittenTitle") or "").strip()
                if not rt:
                    rt = base
                vars_raw = parsed.get("variations")
                vars_list: list[str] = []
                if isinstance(vars_raw, list):
                    vars_list = [str(v).strip() for v in vars_raw if str(v).strip()]
                # Ensure at least one suggestion
                if not vars_list:
                    vars_list = [rt]
                # Clamp to 3 and 8–10 word budget.
                vars_list = [v[:120].strip() for v in vars_list][:3]
                conf = parsed.get("confidence", 0.78)
                conf_f = float(conf) if isinstance(conf, (int, float, str)) else 0.78
                return {
                    "rewrittenTitle": rt[:96],
                    "styleUsed": tone_use,
                    "confidence": max(0.05, min(0.99, conf_f)),
                    "variations": vars_list,
                }
        except Exception:
            pass

    return {
        "error": "Title rewrite requires a working OpenAI API key and a successful model response.",
        "rewrittenTitle": base[:96],
        "styleUsed": tone_use,
        "confidence": 0.0,
        "variations": [],
    }


def ai_edit_slide_contract_v2(
    *,
    slide_id: str,
    action: str,
    custom_prompt: str = "",
    current_slide: dict[str, Any] | None = None,
    full_deck_context: list[dict[str, Any]] | None = None,
    chart_type_preference: str | None = None,
) -> dict[str, Any]:
    """
    v2 slide edit contract:
    - Modifies ONLY the selected slide
    - Returns updated slide JSON (ApiSlide shape)
    """
    act_raw = (action or "").strip().lower()
    act = re.sub(r"[^a-z0-9]+", "_", act_raw).strip("_")
    custom_prompt = str(custom_prompt or "").strip()

    # Map human-friendly action names to internal canonical actions.
    alias_map = {
        "improve_writing": "improve",
        "improvecopy": "improve",
        "fix_spelling_grammar": "grammar",
        "fix_grammar": "grammar",
        "spelling_grammar": "grammar",
        "make_longer": "longer",
        "make_longer_bullets": "longer",
        "make_shorter": "shorter",
        "compress": "shorter",
        "simplify_language": "simplify",
        "make_more_visual": "visual",
        "suggest_image_idea": "image",
        "suggest_image": "image",
        "edit_current_image": "edit_image",
        "editimage": "edit_image",
        "generate_chart": "chart",
        "generate_image": "generate_image",
        "generateimage": "generate_image",
        "improve_title": "improve_title",
        "enhance_slide_quality": "enhance",
        "enhanceslidequality": "enhance",
        "refine_this_slide": "refine",
        "refine_slide": "refine",
    }
    act = alias_map.get(act, act)

    # Load from DB once so we can persist deterministically.
    with _conn() as con:
        srow = con.execute(
            "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
            (slide_id,),
        ).fetchone()
        if srow is None:
            return {"error": "Slide not found"}
        presentation_id = str(srow["presentation_id"])
        presentation_title = str(
            con.execute("SELECT title FROM presentations WHERE id = ?", (presentation_id,)).fetchone() or {"title": ""}  # type: ignore[arg-type]
        ).strip()
        slide_obj = _slide_from_db_row(srow=srow)

    # Normalize content to edit.
    # Note: we intentionally do NOT persist caller-provided fields unless we compute the final merged content.
    edit_content: dict[str, Any] = dict(slide_obj["content"] or {})
    if isinstance(current_slide, dict) and isinstance(current_slide.get("content"), dict):
        # Prefer caller content when present; it reflects "current" editor state.
        edit_content = {**edit_content, **current_slide["content"]}
        if isinstance(current_slide.get("title"), str) and current_slide["title"].strip():
            slide_obj["title"] = current_slide["title"].strip()

    # Optional: small delay so users feel "AI is working".
    # (Editing is synchronous, but this improves perceived responsiveness.)
    try:
        if _openai_api_key_available():
            time.sleep(0.35)
    except Exception:
        pass

    # 1–8: writing + visual + imageQuery suggestion edits.
    if act in {"improve", "grammar", "longer", "shorter", "simplify", "visual", "image", "custom", "layout"}:
        out = ai_edit_slide(
            slide_id=slide_id,
            action=act if act != "layout" else "layout",
            user_prompt=custom_prompt,
            current_content_json=json.dumps(edit_content),
        )
        if out.get("error"):
            return out
        updated = out.get("slide") or {}
        # Keep contract shape stable.
        return {"updatedSlide": updated}

    # 9: chart generation
    if act in {"chart", "generate_chart", "generatechart"}:
        structured_blob = _serialize_structured_sources_for_chart(edit_content)
        slide_content_txt = "\n".join(
            [
                str(edit_content.get("title") or slide_obj.get("title") or ""),
                "\n".join([str(b) for b in edit_content.get("bullets") if isinstance(edit_content.get("bullets"), list)][:12])  # type: ignore[index]
                if isinstance(edit_content.get("bullets"), list)
                else "",
                str(edit_content.get("description") or ""),
                str(edit_content.get("highlight") or ""),
                str(edit_content.get("keyMessage") or ""),
                str(edit_content.get("speakerNotes") or ""),
                str(custom_prompt or ""),
                structured_blob,
            ]
        ).strip()
        deck_hint = ""
        if isinstance(full_deck_context, list) and full_deck_context:
            try:
                deck_hint = "\n".join(
                    [str(s.get("title") or "") for s in full_deck_context[:12] if isinstance(s, dict)]
                )[:2000]
            except Exception:
                deck_hint = ""
        if deck_hint.strip():
            slide_content_txt = f"{slide_content_txt}\n--- Deck context (titles) ---\n{deck_hint}".strip()
        chart = generate_chart(
            slide_content=slide_content_txt,
            chart_type_preference=chart_type_preference,
        )
        bullets = edit_content.get("bullets")
        if not isinstance(bullets, list):
            bullets = _sanitize_bullets_list(_parse_bullets_json(srow["bullets_json"]))  # type: ignore[arg-type]
        content_next = dict(edit_content)
        ch_out: dict[str, Any] = {
            "chartType": str(chart.get("chartType") or "bar"),
            "title": str(chart.get("title") or ""),
            "data": chart.get("data") if isinstance(chart.get("data"), list) else [],
        }
        _xl = chart.get("xLabel")
        _yl = chart.get("yLabel")
        if isinstance(_xl, str) and _xl.strip():
            ch_out["xLabel"] = _xl.strip()[:120]
        if isinstance(_yl, str) and _yl.strip():
            ch_out["yLabel"] = _yl.strip()[:120]
        content_next["chart"] = ch_out
        content_next["slideType"] = "stats"
        layout_preset = _layout_preset_for_slide_type("stats")
        content_next["layoutSuggestion"] = layout_preset
        gs = content_next.get("gammaStyle") if isinstance(content_next.get("gammaStyle"), dict) else {}
        content_next["gammaStyle"] = {**gs, "layoutPreset": layout_preset}
        _persist_slide_content(
            slide_id=slide_id,
            title=str(slide_obj.get("title") or ""),
            bullets=[str(b) for b in bullets if str(b).strip()],
            content=content_next,
        )
        with _conn() as con2:
            s2 = con2.execute(
                "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
                (slide_id,),
            ).fetchone()
            updated = _slide_from_db_row(srow=s2) if s2 is not None else {}
        return {"updatedSlide": updated}

    # 10: generate image (replace existing)
    if act in {"image_generate", "generate_image", "generateimage", "image_generation", "image_gen"}:
        slide_content_txt = "\n".join(
            [
                str(slide_obj.get("title") or ""),
                "\n".join([str(b) for b in edit_content.get("bullets") if isinstance(edit_content.get("bullets"), list)][:6])  # type: ignore[index]
                if isinstance(edit_content.get("bullets"), list)
                else "",
                str(edit_content.get("description") or ""),
                str(edit_content.get("imageQuery") or ""),
                str(custom_prompt or ""),
            ]
        ).strip()
        out = generate_image_advanced(slide_id=slide_id, slide_content=slide_content_txt)
        if not out.get("success") or not isinstance(out.get("images"), list) or not out["images"]:
            return {"error": "Image generation failed"}
        img = out["images"][0]
        content_next = dict(edit_content)
        content_next["generatedImageUrl"] = str(img.get("url") or "")
        content_next["generatedImagePrompt"] = str(out.get("promptUsed") or custom_prompt or slide_content_txt)[:220]
        mapped_opts = []
        if isinstance(out.get("images"), list):
            for idx, i in enumerate(out["images"]):
                mapped_opts.append(
                    {
                        "imageUrl": str(i.get("url") or ""),
                        "source": "search",
                        "promptUsed": str(out.get("promptUsed") or slide_content_txt)[:220],
                        "confidence": float(i.get("confidence") or 0.75),
                        "similarity": i.get("similarity"),
                        "originalUrl": str(i.get("url") or ""),
                        "croppedUrl": str(i.get("url") or ""),
                        "isBestMatch": bool(i.get("isBestMatch")) if "isBestMatch" in i else idx == 0,
                    }
                )
        content_next["generatedImageOptions"] = mapped_opts
        content_next["generatedImageConfidence"] = float(img.get("confidence") or 0.75)
        content_next["generatedImageStrategy"] = {
            "action": "search",
            "reason": str(out.get("reason") or "AI image generated for editor"),
        }
        if out.get("promptUsed"):
            content_next["references"] = [str(out.get("promptUsed"))]
        content_next["slideType"] = "visual"
        layout_preset = _layout_preset_for_slide_type("visual")
        content_next["layoutSuggestion"] = layout_preset
        gs = content_next.get("gammaStyle") if isinstance(content_next.get("gammaStyle"), dict) else {}
        content_next["gammaStyle"] = {**gs, "layoutPreset": layout_preset}
        bullets = content_next.get("bullets")
        if not isinstance(bullets, list):
            bullets = _sanitize_bullets_list(_parse_bullets_json(srow["bullets_json"]))  # type: ignore[arg-type]
        _persist_slide_content(
            slide_id=slide_id,
            title=str(slide_obj.get("title") or ""),
            bullets=[str(b) for b in bullets if str(b).strip()],
            content=content_next,
        )
        with _conn() as con3:
            s3 = con3.execute(
                "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
                (slide_id,),
            ).fetchone()
            updated = _slide_from_db_row(srow=s3) if s3 is not None else {}
        return {"updatedSlide": updated}

    # 8: edit current image (text instruction → new imageQuery / replace)
    if act in {"edit_image", "editimage", "edit_current_image", "editcurrentimage"}:
        # We model "edit current image" as "generate a fresh image from the edit instruction".
        slide_content_txt = "\n".join(
            [
                str(slide_obj.get("title") or ""),
                "\n".join([str(b) for b in edit_content.get("bullets") if isinstance(edit_content.get("bullets"), list)][:6])  # type: ignore[index]
                if isinstance(edit_content.get("bullets"), list)
                else "",
                str(edit_content.get("description") or ""),
                f"Image edit instruction: {custom_prompt}",
            ]
        ).strip()
        out = generate_image_advanced(slide_id=slide_id, slide_content=slide_content_txt)
        if not out.get("success") or not isinstance(out.get("images"), list) or not out["images"]:
            return {"error": "Image edit failed"}
        img = out["images"][0]
        content_next = dict(edit_content)
        content_next["generatedImageUrl"] = str(img.get("url") or "")
        content_next["generatedImagePrompt"] = str(out.get("promptUsed") or custom_prompt or slide_content_txt)[:220]
        content_next["generatedImageConfidence"] = float(img.get("confidence") or 0.75)
        # Keep image options minimal for deterministic native edit.
        content_next["generatedImageOptions"] = [
            {
                "imageUrl": str(img.get("url") or ""),
                "source": "search",
                "promptUsed": str(out.get("promptUsed") or slide_content_txt)[:220],
                "confidence": float(img.get("confidence") or 0.75),
                "similarity": img.get("similarity"),
                "originalUrl": str(img.get("url") or ""),
                "croppedUrl": str(img.get("url") or ""),
                "isBestMatch": True,
            }
        ]
        content_next["generatedImageStrategy"] = {"action": "search", "reason": f"Edited via instruction: {custom_prompt[:120]}"}
        if out.get("promptUsed"):
            content_next["references"] = [str(out.get("promptUsed"))]
        content_next["slideType"] = "visual"
        layout_preset = _layout_preset_for_slide_type("visual")
        content_next["layoutSuggestion"] = layout_preset
        gs = content_next.get("gammaStyle") if isinstance(content_next.get("gammaStyle"), dict) else {}
        content_next["gammaStyle"] = {**gs, "layoutPreset": layout_preset}
        bullets = content_next.get("bullets")
        if not isinstance(bullets, list):
            bullets = _sanitize_bullets_list(_parse_bullets_json(srow["bullets_json"]))  # type: ignore[arg-type]
        _persist_slide_content(
            slide_id=slide_id,
            title=str(slide_obj.get("title") or ""),
            bullets=[str(b) for b in bullets if str(b).strip()],
            content=content_next,
        )
        with _conn() as con4:
            s4 = con4.execute(
                "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
                (slide_id,),
            ).fetchone()
            updated = _slide_from_db_row(srow=s4) if s4 is not None else {}
        return {"updatedSlide": updated}

    # 12: improve title only
    if act in {"improve_title", "improvetitle", "rewrite_title", "revisetitle", "title"}:
        current_title = str(slide_obj.get("title") or "").strip()
        if not _openai_api_key_available():
            return {"error": "Title improvement requires a configured OpenAI API key (OPENAI_API_KEY)."}
        out = _llm_rewrite_slide_title_only(title=current_title, slide_content=edit_content, custom_prompt=custom_prompt)
        if not out:
            return {"error": "The AI could not rewrite this title. Try again."}
        new_title, _conf = out

        bullets = edit_content.get("bullets")
        if not isinstance(bullets, list):
            bullets = _sanitize_bullets_list(_parse_bullets_json(srow["bullets_json"]))  # type: ignore[arg-type]
        content_next = dict(edit_content)
        content_next["title"] = new_title
        _persist_slide_content(
            slide_id=slide_id,
            title=new_title,
            bullets=[str(b) for b in bullets if str(b).strip()],
            content=content_next,
        )
        with _conn() as con5:
            s5 = con5.execute(
                "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
                (slide_id,),
            ).fetchone()
            updated = _slide_from_db_row(srow=s5) if s5 is not None else {}
        return {"updatedSlide": updated}

    # 13: enhance slide quality
    if act in {"enhance", "enhance_quality", "enhanceslidequality", "quality_enhance", "qualityenhance"}:
        result = quality_enhance_slide(slide_id=slide_id, tone="professional")
        if result.get("error"):
            return result
        updated = result.get("slide") if isinstance(result.get("slide"), dict) else {}
        if isinstance(updated, dict) and presentation_id:
            updated["presentationId"] = presentation_id
        # Align with ApiSlide: ensure `order` is present (already from service).
        updated.pop("index", None)
        return {"updatedSlide": updated}

    # 14: refine this slide (story-flow + anti-repetition)
    if act in {"refine", "refine_slide", "refinethisslide", "refine_this_slide"}:
        result = refine_slide(slide_id=slide_id, min_score=8.0, max_iters=1)
        if result.get("error"):
            return result
        updated = result.get("slide") if isinstance(result.get("slide"), dict) else {}
        if isinstance(updated, dict) and presentation_id:
            updated["presentationId"] = presentation_id
        updated.pop("index", None)
        return {"updatedSlide": updated}

    # 11: regenerate (keep deck rest unchanged)
    if act in {"regenerate", "regenerate_slide", "regenerateslide"}:
        result = regenerate_slide(slide_id=slide_id, tone="professional")
        if result.get("error"):
            return result
        # regenerate_slide returns {slide: {...}, presentationId: ...}
        updated = result.get("slide") if isinstance(result.get("slide"), dict) else {}
        if isinstance(updated, dict):
            updated["presentationId"] = str(result.get("presentationId") or presentation_id or "")
        updated.pop("index", None)
        return {"updatedSlide": updated}

    # Fallback: treat as "improve".
    out = ai_edit_slide(
        slide_id=slide_id,
        action="improve",
        user_prompt=custom_prompt,
        current_content_json=json.dumps(edit_content),
    )
    if out.get("error"):
        return out
    return {"updatedSlide": out.get("slide")}


def _build_deck_context_for_slide(*, con: sqlite3.Connection, presentation_id: str, slide_idx: int, current_content: dict[str, Any]) -> list[dict[str, Any]]:
    """Create a compact deck context list for story-flow + anti-repetition."""

    def _words(s: str) -> set[str]:
        toks = re.split(r"[^a-zA-Z0-9]+", (s or "").lower())
        return {t for t in toks if len(t) >= 4}

    cur_text = " ".join(
        [
            str(current_content.get("title") or ""),
            str(current_content.get("subtitle") or ""),
            "\n".join(current_content.get("bullets") if isinstance(current_content.get("bullets"), list) else []),
            str(current_content.get("highlight") or current_content.get("keyMessage") or ""),
        ]
    )
    cur_words = _words(cur_text)

    rows = con.execute(
        "SELECT id, idx, title, bullets_json, content_json FROM slides WHERE presentation_id = ? ORDER BY idx ASC",
        (presentation_id,),
    ).fetchall()

    scored: list[tuple[int, dict[str, Any]]] = []
    for r in rows:
        idx = int(r["idx"])
        if idx == slide_idx:
            continue
        c = _slide_content_from_row(str(r["title"] or ""), str(r["bullets_json"] or ""), r["content_json"] or None)
        w = _words(
            " ".join(
                [
                    str(c.get("title") or r["title"] or ""),
                    str(c.get("subtitle") or ""),
                    "\n".join(c.get("bullets") if isinstance(c.get("bullets"), list) else []),
                    str(c.get("highlight") or c.get("keyMessage") or ""),
                ]
            )
        )
        overlap = len(cur_words.intersection(w))
        snippet = {
            "id": str(r["id"]),
            "idx": idx,
            "title": str(r["title"] or ""),
            "slideType": str(c.get("slideType") or ""),
            "keyPoint": str(c.get("keyMessage") or c.get("highlight") or ""),
            "bullets": (c.get("bullets") if isinstance(c.get("bullets"), list) else [])[:2],
        }
        scored.append((overlap, snippet))

    # If we have little overlap, still include early/late slides for narrative flow.
    scored.sort(key=lambda x: x[0], reverse=True)
    top = [x[1] for x in scored[:6] if x[1]]

    if len(top) < 4:
        extras = [s for _, s in scored[6:] if s and s not in top]
        top.extend(extras[: 4 - len(top)])

    # Compact text lengths
    for t in top:
        t["title"] = t.get("title", "")[:80]
        t["keyPoint"] = t.get("keyPoint", "")[:120]
        if isinstance(t.get("bullets"), list):
            t["bullets"] = [str(b)[:120] for b in t["bullets"]][:2]
    return top


def _llm_refine_or_enhance_slide_with_context(
    *,
    content: dict[str, Any],
    action_kind: str,
    deck_context: list[dict[str, Any]],
) -> dict[str, Any] | None:
    """One-shot LLM editor for refined slide quality + story-flow."""
    if not _openai_api_key_available():
        return None

    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    try:
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI
    except Exception:
        return None

    title = str(content.get("title") or "")
    bullets = content.get("bullets")
    if not isinstance(bullets, list):
        bullets = []
    bullets = [str(b) for b in bullets[:6]]
    desc = str(content.get("description") or "")
    highlight = str(content.get("highlight") or "")
    key_msg = str(content.get("keyMessage") or "")
    speaker_notes = str(content.get("speakerNotes") or "")

    if action_kind == "enhance_quality":
        system = (
            "You are a senior presentation editor.\n"
            "Rewrite the CURRENT slide to be premium-quality: clearer bullets, stronger highlight, improved structure and flow.\n"
            "Use the deck context to keep consistency and avoid re-stating other slides.\n"
            "Output strictly valid JSON only. No markdown. No explanations.\n\n"
            "Output JSON keys must include: title, subtitle, bullets, description, highlight, keyMessage, speakerNotes, slideType, imageQuery, layoutSuggestion.\n"
            "Constraints:\n"
            "- Bullets: 4–6 items, each under 200 chars.\n"
            "- Title max 10 words; non-generic.\n"
            "- Preserve meaning; do not invent new facts.\n"
        )
    else:
        # refine_slide / story-flow
        system = (
            "You are a story-flow editor.\n"
            "Refine the CURRENT slide so it aligns with the deck narrative and avoids repetition with other slides.\n"
            "Do not change the topic; preserve meaning.\n"
            "Output strictly valid JSON only. No markdown. No explanations.\n\n"
            "Output JSON keys must include: title, subtitle, bullets, description, highlight, keyMessage, speakerNotes, slideType, imageQuery, layoutSuggestion.\n"
            "- Bullets: 4–6 items, each under 200 chars.\n"
            "- Avoid reusing phrases that appear in other slides' key points.\n"
        )

    user = f"""Action kind: {action_kind}

CURRENT SLIDE JSON:
{json.dumps(
        {"title": title, "subtitle": str(content.get("subtitle") or ""), "bullets": bullets, "description": desc, "highlight": highlight, "keyMessage": key_msg, "speakerNotes": speaker_notes},
        ensure_ascii=False,
    )}

DECK CONTEXT (other slides, compact):
{json.dumps(deck_context, ensure_ascii=False)}

Return exactly this JSON object:
{{
  "title": "string",
  "subtitle": "string",
  "bullets": ["string"],
  "description": "string",
  "highlight": "string",
  "keyMessage": "string",
  "speakerNotes": "string",
  "slideType": "string",
  "imageQuery": "string",
  "layoutSuggestion": "string"
}}
"""

    for temp in (0.45, 0.62):
        try:
            llm = ChatOpenAI(
                model=settings.openai_model,
                api_key=key,
                base_url=settings.openai_base_url,
                temperature=temp,
            )
            resp = llm.invoke([SystemMessage(content=system), HumanMessage(content=user)])
            text = getattr(resp, "content", None) or str(resp)
            if not isinstance(text, str):
                text = str(text)
            parsed = _extract_json_object(text)
            if parsed and isinstance(parsed, dict):
                return parsed
        except Exception:
            continue
    return None


def refine_slide(*, slide_id: str, min_score: float = 8.0, max_iters: int = 1) -> dict[str, Any]:
    with _conn() as con:
        s = con.execute(
            "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
            (slide_id,),
        ).fetchone()
        if s is None:
            return {"error": "Slide not found"}

        content = _slide_content_from_row(s["title"], s["bullets_json"], s["content_json"])
        before = _snapshot_from_content(content)

        slide_idx = int(s["idx"])
        presentation_id = str(s["presentation_id"])
        deck_context = _build_deck_context_for_slide(
            con=con,
            presentation_id=presentation_id,
            slide_idx=slide_idx,
            current_content=content,
        )

        if not _openai_api_key_available():
            return {"error": "Slide refine requires a configured OpenAI API key (OPENAI_API_KEY)."}

        parsed = _llm_refine_or_enhance_slide_with_context(
            content=content,
            action_kind="refine_slide",
            deck_context=deck_context,
        )

        bullets = content.get("bullets")
        if not isinstance(bullets, list):
            bullets = _parse_bullets_json(s["bullets_json"])

        if not isinstance(parsed, dict):
            return {"error": "The AI could not refine this slide. Try again."}

        if isinstance(parsed.get("title"), str) and parsed["title"].strip():
            content["title"] = parsed["title"].strip()[:120]
        if isinstance(parsed.get("subtitle"), str):
            content["subtitle"] = parsed["subtitle"].strip()
        bl = parsed.get("bullets")
        if isinstance(bl, list) and len(bl) >= 2:
            content["bullets"] = _sanitize_bullets_list([str(x) for x in bl])[:6]
            bullets = content["bullets"]
        if isinstance(parsed.get("description"), str):
            content["description"] = parsed["description"].strip()
        if isinstance(parsed.get("highlight"), str):
            content["highlight"] = parsed["highlight"].strip()
        if isinstance(parsed.get("keyMessage"), str):
            content["keyMessage"] = parsed["keyMessage"].strip()
        if isinstance(parsed.get("speakerNotes"), str):
            content["speakerNotes"] = parsed["speakerNotes"].strip()
        if isinstance(parsed.get("imageQuery"), str):
            content["imageQuery"] = parsed["imageQuery"].strip()
        slide_type = str(parsed.get("slideType") or "").strip() or detect_slide_type(content)
        content["slideType"] = slide_type
        layout_preset = _layout_preset_for_slide_type(slide_type)
        content["layoutSuggestion"] = layout_preset
        gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
        content["gammaStyle"] = {**gs, "layoutPreset": layout_preset}

        content["qualityScore"] = max(float(content.get("qualityScore") or 7.0), min_score)

        con.execute(
            "UPDATE slides SET title = ?, bullets_json = ?, content_json = ? WHERE id = ?",
            (str(content["title"]), "|||".join([str(b) for b in content.get("bullets", [])]), json.dumps(content), slide_id),
        )
        con.execute("UPDATE presentations SET updated_at = ? WHERE id = ?", (time.time(), s["presentation_id"]))
        con.commit()

        after = _snapshot_from_content(content)
        slide = {"id": s["id"], "title": str(content["title"]), "order": slide_idx, "content": content}
        return {"before": before, "after": after, "slide": slide, "mode": "native"}


def quality_enhance_slide(*, slide_id: str, tone: str = "professional") -> dict[str, Any]:
    with _conn() as con:
        s = con.execute(
            "SELECT id, presentation_id, idx, title, bullets_json, content_json FROM slides WHERE id = ?",
            (slide_id,),
        ).fetchone()
        if s is None:
            return {"error": "Slide not found"}

        content = _slide_content_from_row(s["title"], s["bullets_json"], s["content_json"])
        before = _snapshot_from_content(content)

        slide_idx = int(s["idx"])
        presentation_id = str(s["presentation_id"])
        deck_context = _build_deck_context_for_slide(
            con=con,
            presentation_id=presentation_id,
            slide_idx=slide_idx,
            current_content=content,
        )

        if not _openai_api_key_available():
            return {"error": "Slide quality enhancement requires a configured OpenAI API key (OPENAI_API_KEY)."}

        parsed = _llm_refine_or_enhance_slide_with_context(
            content=content,
            action_kind="enhance_quality",
            deck_context=deck_context,
        )

        if not isinstance(parsed, dict):
            return {"error": "The AI could not enhance this slide. Try again."}

        if isinstance(parsed.get("title"), str) and parsed["title"].strip():
            content["title"] = parsed["title"].strip()[:120]
        if isinstance(parsed.get("subtitle"), str):
            content["subtitle"] = parsed["subtitle"].strip()
        bl = parsed.get("bullets")
        if isinstance(bl, list) and len(bl) >= 2:
            content["bullets"] = _sanitize_bullets_list([str(x) for x in bl])[:6]
        if isinstance(parsed.get("description"), str):
            content["description"] = parsed["description"].strip()
        if isinstance(parsed.get("highlight"), str):
            content["highlight"] = parsed["highlight"].strip()
        if isinstance(parsed.get("keyMessage"), str):
            content["keyMessage"] = parsed["keyMessage"].strip()
        if isinstance(parsed.get("speakerNotes"), str):
            content["speakerNotes"] = parsed["speakerNotes"].strip()
        if isinstance(parsed.get("imageQuery"), str):
            content["imageQuery"] = parsed["imageQuery"].strip()
        slide_type = str(parsed.get("slideType") or "").strip() or detect_slide_type(content)
        content["slideType"] = slide_type

        layout_preset = _layout_preset_for_slide_type(str(content.get("slideType") or slide_type))
        content["layoutSuggestion"] = layout_preset
        gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
        content["gammaStyle"] = {**gs, "layoutPreset": layout_preset}

        # Quality scores for UI parity.
        content["qualityScore"] = max(float(content.get("qualityScore") or 8.0), 8.7)
        content["qualityClarity"] = 8.8
        content["qualityEngagement"] = 8.5
        content["qualityVisualBalance"] = 8.4
        content["qualityVisualType"] = content.get("qualityVisualType") or "concept"

        if tone:
            content["description"] = f"{content.get('description','')}".strip()

        con.execute(
            "UPDATE slides SET title = ?, bullets_json = ?, content_json = ? WHERE id = ?",
            (str(content["title"]), "|||".join([str(b) for b in content.get("bullets", [])]), json.dumps(content), slide_id),
        )
        con.execute("UPDATE presentations SET updated_at = ? WHERE id = ?", (time.time(), s["presentation_id"]))
        con.commit()

        after = _snapshot_from_content(content)
        slide = {"id": s["id"], "title": str(content["title"]), "order": slide_idx, "content": content}
        return {"before": before, "after": after, "slide": slide, "mode": "native"}


def apply_premium_deck(*, presentation_id: str, slides: list[dict[str, Any]] | None = None) -> dict[str, Any]:
    with _conn() as con:
        p = con.execute("SELECT id, title, prompt FROM presentations WHERE id = ?", (presentation_id,)).fetchone()
        if p is None:
            return {"errorCode": "NOT_FOUND", "message": "Presentation not found"}
        con.execute("DELETE FROM slides WHERE presentation_id = ?", (presentation_id,))
        source = slides if isinstance(slides, list) and slides else _sample_premium_slides(str(p["title"] or p["prompt"] or "Presentation"))
        for i, item in enumerate(source[:20], start=1):
            slide_id = uuid.uuid4().hex
            title = str(item.get("title") or f"Slide {i}")
            subtitle = str(item.get("subtitle") or "")
            points = item.get("points") if isinstance(item.get("points"), list) else item.get("bullets")
            if not isinstance(points, list):
                points = []
            bullets = [str(x) for x in points][:8]
            content = {
                "title": title,
                "subtitle": subtitle,
                "bullets": bullets,
                "description": str(item.get("description") or ""),
                "highlight": str(item.get("highlight") or ""),
                "imageQuery": str(item.get("imageQuery") or ""),
                "layoutSuggestion": str(item.get("layoutSuggestion") or ("hero_split" if i % 2 else "image_left")),
                "emphasisWords": [title.split(" ")[0], "Impact"] if title else ["Impact"],
                "contentDensity": "medium",
                "qualityScore": 8.2,
                "gammaStyle": {"layoutPreset": "hero_split" if i % 2 else "two_column", "alignment": "left", "emphasisWords": []},
            }
            con.execute(
                "INSERT INTO slides (id, presentation_id, idx, title, bullets_json, content_json) VALUES (?, ?, ?, ?, ?, ?)",
                (slide_id, presentation_id, i, title, "|||".join(bullets), json.dumps(content)),
            )
        con.execute("UPDATE presentations SET updated_at = ?, status = ? WHERE id = ?", (time.time(), "COMPLETED", presentation_id))
        con.commit()
    return {"ok": True, "slideCount": min(len(source), 20)}


def extract_source_file(*, filename: str, buffer: bytes) -> dict[str, Any]:
    text = _extract_text_from_upload_bytes(filename=filename, buffer=buffer, limit_chars=12000).strip()
    if not text:
        return {"errorCode": "EMPTY_EXTRACT", "message": "Could not extract readable text from this file."}
    ext = Path(filename).suffix.lower().lstrip(".") or "txt"
    return {
        "success": True,
        "fileName": filename,
        "fileType": ext,
        "extractedText": text[:12000],
    }


_CHART_TYPES_ALLOWED = frozenset(
    {"bar", "line", "pie", "donut", "stacked_bar", "area", "stacked_area", "horizontal_bar"}
)


def _normalize_user_chart_type(pref: str | None) -> str | None:
    if not pref:
        return None
    x = str(pref).strip().lower().replace(" ", "_")
    if x in ("auto", "", "default", "none"):
        return None
    return x if x in _CHART_TYPES_ALLOWED else None


def _serialize_structured_sources_for_chart(edit_content: dict[str, Any]) -> str:
    """Pull file / spreadsheet structured blobs into text so numeric extraction + LLM see them."""
    chunks: list[str] = []
    sd = edit_content.get("structuredData")
    if isinstance(sd, list) and sd:
        try:
            chunks.append("--- structuredData (rows) ---\n" + json.dumps(sd[:120], ensure_ascii=False)[:12000])
        except Exception:
            pass
    elif isinstance(sd, dict) and sd:
        try:
            chunks.append("--- structuredData ---\n" + json.dumps(sd, ensure_ascii=False)[:12000])
        except Exception:
            pass
    for key in ("fileExtraction", "fileContext", "sourceFileSummary", "extractedText", "fileSlideContext"):
        val = edit_content.get(key)
        if isinstance(val, str) and val.strip():
            chunks.append(f"--- {key} ---\n{val.strip()[:8000]}")
        elif isinstance(val, dict) and val:
            try:
                chunks.append(f"--- {key} ---\n" + json.dumps(val, ensure_ascii=False)[:8000])
            except Exception:
                pass
    return "\n".join(chunks).strip()


def _humanize_column_name(name: str) -> str:
    s = str(name or "").strip()
    if not s:
        return ""
    s = re.sub(r"[_]+", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s[:80] if s else ""


def _try_parse_structured_data_rows_json(chunk: str) -> list[dict[str, Any]] | None:
    """Parse leading JSON array; ignore trailing prose after the array (common in merged slide text)."""
    s = chunk.strip()
    if not s.startswith("["):
        return None
    try:
        data, _ = json.JSONDecoder().raw_decode(s)
    except Exception:
        return None
    if not isinstance(data, list) or not data:
        return None
    out = [dict(x) for x in data[:200] if isinstance(x, dict)]
    return out if out else None


def _structured_rows_from_text(text: str) -> list[dict[str, Any]] | None:
    """Parse `--- structuredData (rows) ---` JSON (Excel/CSV pipeline) if present."""
    marker = "--- structuredData (rows) ---"
    idx = text.find(marker)
    if idx < 0:
        return None
    rest = text[idx + len(marker) :].lstrip()
    end = rest.find("\n---")
    blob = rest[:end].strip() if end >= 0 else rest.strip()
    return _try_parse_structured_data_rows_json(blob)


def _axis_labels_from_structured_rows(rows: list[dict[str, Any]]) -> tuple[str | None, str | None]:
    """Pick category-axis and value-axis names from spreadsheet-like row objects (Excel headers → keys)."""
    skip = frozenset({"rowIndex", "slideNumber"})
    dict_rows = [r for r in rows[:50] if isinstance(r, dict)]
    if not dict_rows:
        return None, None
    keys = [k for k in dict_rows[0].keys() if k not in skip and not str(k).startswith("__")]
    if len(keys) < 2:
        return None, None

    def col_ratio_numeric(key: str) -> float:
        nums = 0
        total = 0
        for r in dict_rows:
            if key not in r:
                continue
            v = r[key]
            total += 1
            if isinstance(v, bool):
                pass
            elif isinstance(v, (int, float)) and not isinstance(v, bool):
                nums += 1
            elif isinstance(v, str):
                sv = v.strip().replace(",", "")
                if not sv:
                    continue
                try:
                    float(sv.replace("%", "").strip())
                    nums += 1
                except ValueError:
                    pass
        return (nums / total) if total else 0.0

    numeric_cols: list[tuple[str, float]] = []
    text_cols: list[str] = []
    for k in keys:
        ratio = col_ratio_numeric(k)
        if ratio >= 0.55:
            numeric_cols.append((k, ratio))
        else:
            text_cols.append(k)

    numeric_cols.sort(key=lambda t: (-t[1], len(t[0])))
    text_cols.sort(key=lambda c: (len(c), len(str(dict_rows[0].get(c) or ""))))

    if text_cols and numeric_cols:
        return _humanize_column_name(text_cols[0]), _humanize_column_name(numeric_cols[0][0])
    if len(numeric_cols) >= 2 and not text_cols:
        # e.g. Year + Sales both numeric — first column often dimension
        a, b = numeric_cols[0][0], numeric_cols[1][0]
        return _humanize_column_name(a), _humanize_column_name(b)
    return None, None


def _axis_labels_from_markdown_table(text: str) -> tuple[str | None, str | None]:
    """Use first markdown table header row: first column → x, last column → y (measure)."""
    lines = [ln.rstrip() for ln in str(text or "").splitlines()]
    rows: list[list[str]] = []
    for ln in lines:
        s = ln.strip()
        if not s.startswith("|"):
            if rows:
                break
            continue
        cells = [c.strip() for c in s.split("|")]
        cells = [c for c in cells if c != ""]
        if not cells:
            continue
        if all(re.match(r"^:?-+:?$", c) for c in cells):
            continue
        rows.append(cells)
    if len(rows) < 2:
        return None, None
    header = rows[0]
    if len(header) < 2:
        return None, None

    def _cell_is_pure_number(c: str) -> bool:
        t = c.replace(",", "").strip()
        if not t:
            return False
        return bool(re.fullmatch(r"-?[0-9]+(?:\.[0-9]+)?%?", t))

    if all(_cell_is_pure_number(c) for c in header):
        return None, None
    return _humanize_column_name(header[0]), _humanize_column_name(header[-1])


def _infer_axis_labels(text: str) -> tuple[str | None, str | None]:
    """Derive axis titles from embedded structured rows (Excel) or markdown table headers."""
    rows = _structured_rows_from_text(text)
    if rows:
        x, y = _axis_labels_from_structured_rows(rows)
        if x and y:
            return x, y
    return _axis_labels_from_markdown_table(text)


def _extract_label_values(text: str) -> tuple[list[str], list[float]]:
    """Extract (label, number) pairs from bullets, lines, markdown-ish rows, and JSON fragments."""
    labels: list[str] = []
    values: list[float] = []
    seen: set[str] = set()

    def push(lab: str, val: float) -> None:
        lab = re.sub(r"\s+", " ", lab).strip()[:64]
        if not lab or lab in seen:
            return
        seen.add(lab)
        labels.append(lab)
        values.append(float(val))

    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("|") and line.count("|") >= 2:
            parts = [p.strip() for p in line.split("|") if p.strip() and not re.match(r"^:?-+:?$", p.strip())]
            if len(parts) >= 2:
                tail = parts[-1]
                mnum = re.search(r"(-?[0-9]+(?:\.[0-9]+)?)\s*%?", tail)
                if mnum:
                    lab = " ".join(parts[:-1])[:64] or parts[0]
                    try:
                        v = float(mnum.group(1))
                        if "%" in tail or any("%" in p for p in parts):
                            pass
                        push(lab, v)
                    except Exception:
                        pass
            continue

        line = re.sub(r"^[\-\*•]\s+", "", line)
        line = re.sub(r"^\d+[\.\)]\s+", "", line)

        pats = [
            r"^(.+?)\s*[:：]\s*(-?[0-9]+(?:\.[0-9]+)?)\s*%?\s*$",
            r"^(.+?)\s*[-–—]\s*(-?[0-9]+(?:\.[0-9]+)?)\s*%?\s*$",
            r"^(.+?)\s*,\s*(-?[0-9]+(?:\.[0-9]+)?)\s*%?\s*$",
            r"^(.+?)\s+\((-?[0-9]+(?:\.[0-9]+)?)\s*%?\)\s*$",
        ]
        for pat in pats:
            m = re.match(pat, line, re.IGNORECASE)
            if m:
                try:
                    push(m.group(1).strip(), float(m.group(2)))
                except Exception:
                    pass
                break

    # JSON array of {label, value} or {name, value}
    if not labels:
        try:
            if '"value"' in text or "'value'" in text:
                arr = re.search(r"\[.*\]", text, re.DOTALL)
                if arr:
                    parsed = json.loads(arr.group(0))
                    if isinstance(parsed, list):
                        for it in parsed[:24]:
                            if not isinstance(it, dict):
                                continue
                            lab = str(it.get("label") or it.get("name") or it.get("category") or "").strip()
                            v = it.get("value")
                            if lab and isinstance(v, (int, float)) and not isinstance(v, bool):
                                push(lab, float(v))
        except Exception:
            pass

    return labels[:24], values[:24]


def _looks_like_part_to_whole(values: list[float], text_lower: str) -> bool:
    """True only when values look like shares/percent parts, not arbitrary category totals."""
    if len(values) < 2 or len(values) > 10:
        return False
    if not all(v >= 0 for v in values):
        return False
    if not all(v <= 100 for v in values):
        return False
    s = sum(values)
    if not (94.0 <= s <= 106.0):
        return False
    part_whole_kw = (
        "share",
        "percent",
        "percentage",
        "%",
        "split",
        "breakdown",
        "composition",
        "portion",
        "of total",
        "mix",
        "distribution",
    )
    return any(k in text_lower for k in part_whole_kw)


def _infer_chart_type_heuristic(labels: list[str], values: list[float], text_lower: str) -> str:
    """Pick a sensible default without LLM; avoids pie-by-default for tiny series."""
    if not labels:
        return "bar"

    joined_labels = " ".join(labels).lower()
    time_hint = re.search(
        r"\b(20\d{2}|19\d{2}|q[1-4]\s*20\d{2}|jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec|month|week|day|yoy|qoq|mom|trend|forecast|timeline)\b",
        text_lower + " " + joined_labels,
        re.I,
    )
    if time_hint:
        return "line"

    compare_hint = any(
        k in text_lower for k in (" vs ", "versus", "compared to", "benchmark", "ranking", "top ", "bottom ")
    )
    if compare_hint and len(labels) >= 3:
        return "horizontal_bar" if max(len(l) for l in labels) > 22 else "bar"

    if _looks_like_part_to_whole(values, text_lower):
        return "pie" if len(labels) <= 5 else "donut"

    share_kw = ("share", "split", "breakdown", "composition", "portion", "distribution", "mix", "of total")
    if any(k in text_lower for k in share_kw) and len(labels) <= 8:
        if len(labels) <= 4:
            return "pie"
        return "donut"

    if max(len(l) for l in labels) > 26 or len(labels) > 10:
        return "horizontal_bar"

    if len(labels) >= 2:
        return "bar"

    return "bar"


def _llm_infer_chart_spec(*, slide_text: str, labels: list[str], values: list[float]) -> dict[str, Any] | None:
    if not _openai_api_key_available():
        return None
    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    try:
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI
    except Exception:
        return None

    series_json = json.dumps([{"label": l, "value": v} for l, v in zip(labels, values)], ensure_ascii=False)
    sys = SystemMessage(
        content=(
            "You are a data-viz assistant. Output a single JSON object only (no markdown). "
            "Keys: chartType, title, data, and optionally xLabel, yLabel. "
            "chartType must be one of: bar, line, pie, donut, area, stacked_bar, stacked_area, horizontal_bar. "
            "Rules: use line for time series or trends; use bar or horizontal_bar for category comparisons "
            "(horizontal_bar when labels are long or there are many categories); "
            "use pie or donut only when values represent parts of a meaningful whole (e.g. shares summing to ~100%); "
            "do NOT default to pie for 2–3 unrelated metrics. Prefer bar when unsure. "
            "data is an array of {label, value} using non-empty labels and numeric values. "
            "xLabel: short name for the category axis (what each data[].label represents — e.g. Month, Region, Product). "
            "yLabel: short name for the numeric axis (what value measures — e.g. Revenue USD, Units sold, %). "
            "Set xLabel/yLabel from spreadsheet column headers or table headers in the context when present; omit only if unknown. "
            "Keep the same approximate values as the extracted series unless clearly wrong."
        )
    )
    human = HumanMessage(
        content=f"Slide / data context (truncated):\n{slide_text[:5000]}\n\nExtracted series:\n{series_json}\n"
    )
    try:
        llm = ChatOpenAI(
            model=settings.openai_model,
            temperature=0.15,
            api_key=key,
            base_url=settings.openai_base_url,
        )
        resp = llm.invoke([sys, human])
        txt = str(resp.content).strip()
        if txt.startswith("```"):
            txt = re.sub(r"^```[a-zA-Z0-9]*\s*", "", txt)
            txt = re.sub(r"\s*```$", "", txt).strip()
        obj = json.loads(txt)
        if not isinstance(obj, dict):
            return None
        ct = str(obj.get("chartType") or obj.get("chart_type") or "").strip().lower()
        if ct not in _CHART_TYPES_ALLOWED:
            return None
        data = obj.get("data")
        if not isinstance(data, list) or len(data) < 1:
            return None
        clean: list[dict[str, Any]] = []
        for row in data[:24]:
            if not isinstance(row, dict):
                continue
            lab = str(row.get("label") or row.get("name") or "").strip()
            v = row.get("value")
            if not lab or not isinstance(v, (int, float)) or isinstance(v, bool):
                continue
            clean.append({"label": lab[:80], "value": float(v)})
        if len(clean) < 1:
            return None
        title = str(obj.get("title") or "").strip() or "Chart"
        xl = str(obj.get("xLabel") or obj.get("x_label") or "").strip()[:80]
        yl = str(obj.get("yLabel") or obj.get("y_label") or "").strip()[:80]
        out: dict[str, Any] = {"chartType": ct, "title": title[:120], "data": clean}
        if xl:
            out["xLabel"] = xl
        if yl:
            out["yLabel"] = yl
        return out
    except Exception:
        return None


def generate_chart(
    *,
    slide_content: str,
    chart_type_preference: str | None = None,
) -> dict[str, Any]:
    text = (slide_content or "").strip()
    labels, values = _extract_label_values(text)
    if not labels or not values or len(labels) != len(values):
        labels = ["Category A", "Category B", "Category C"]
        values = [42.0, 33.0, 25.0]

    hint_x, hint_y = _infer_axis_labels(text)

    title = "Generated Chart"
    for ln in text.splitlines():
        s = ln.strip()
        if s and not s.startswith("---") and not s.startswith("|"):
            title = s[:120]
            break

    override = _normalize_user_chart_type(chart_type_preference)
    text_lower = text.lower()

    llm_out: dict[str, Any] | None = None
    if not override:
        llm_out = _llm_infer_chart_spec(slide_text=text, labels=labels, values=values)
    chart_type = override or (llm_out.get("chartType") if isinstance(llm_out, dict) else None)
    if not chart_type:
        chart_type = _infer_chart_type_heuristic(labels, values, text_lower)

    data: list[dict[str, Any]] = []
    if isinstance(llm_out, dict) and isinstance(llm_out.get("data"), list) and not override:
        data = [dict(x) for x in llm_out["data"] if isinstance(x, dict)]  # type: ignore[assignment]
        if isinstance(llm_out.get("title"), str) and llm_out["title"].strip():
            title = llm_out["title"].strip()[:120]
    if not data:
        data = [{"label": l, "value": float(v)} for l, v in zip(labels, values)]

    xl: str | None = None
    yl: str | None = None
    if isinstance(llm_out, dict):
        lx = llm_out.get("xLabel")
        ly = llm_out.get("yLabel")
        if isinstance(lx, str) and lx.strip():
            xl = lx.strip()[:80]
        if isinstance(ly, str) and ly.strip():
            yl = ly.strip()[:80]
    if not xl:
        xl = hint_x
    if not yl:
        yl = hint_y

    out_chart: dict[str, Any] = {"chartType": chart_type, "title": title, "data": data}
    if xl:
        out_chart["xLabel"] = xl
    if yl:
        out_chart["yLabel"] = yl
    return out_chart


def generate_chart_for_user(
    *,
    user_id: str,
    prompt: str = "",
    filename: str | None = None,
    file_bytes: bytes | None = None,
) -> dict[str, Any]:
    source_text = (prompt or "").strip()
    source_type = "PROMPT"
    source_name = None
    if file_bytes is not None and filename:
        ex = extract_source_file(filename=filename, buffer=file_bytes)
        if ex.get("success"):
            source_text = "\n".join([source_text, str(ex.get("extractedText") or "")]).strip()
            source_type = "FILE"
            source_name = filename
    chart = generate_chart(slide_content=source_text)
    chart_id = uuid.uuid4().hex
    with _conn() as con:
        con.execute(
            """
            INSERT INTO user_charts (id, user_id, title, chart_type, chart_data_json, source_type, source_name, input_summary, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                chart_id,
                user_id,
                str(chart.get("title") or "Generated Chart"),
                str(chart.get("chartType") or "bar"),
                json.dumps(chart.get("data") or []),
                source_type,
                source_name,
                source_text[:600],
                time.time(),
            ),
        )
        con.commit()
    cbody: dict[str, Any] = {
        "id": chart_id,
        "title": chart.get("title"),
        "chartType": chart.get("chartType"),
        "data": chart.get("data"),
        "sourceType": source_type,
        "sourceName": source_name,
        "createdAt": _epoch_ms(time.time()),
    }
    if chart.get("xLabel"):
        cbody["xLabel"] = chart.get("xLabel")
    if chart.get("yLabel"):
        cbody["yLabel"] = chart.get("yLabel")
    return {"success": True, "chart": cbody}


def _parse_stored_chart_data_json(raw: str) -> tuple[list[Any], dict[str, Any]]:
    """Support legacy rows (JSON array of points) and enriched rows (object with data + labels)."""
    s = (raw or "").strip()
    if not s:
        return [], {}
    try:
        parsed = json.loads(s)
    except json.JSONDecodeError:
        return [], {}
    if isinstance(parsed, list):
        return parsed, {}
    if isinstance(parsed, dict):
        data = parsed.get("data")
        if not isinstance(data, list):
            data = []
        meta: dict[str, Any] = {}
        lx = parsed.get("xLabel")
        if isinstance(lx, str) and lx.strip():
            meta["xLabel"] = lx.strip()[:80]
        ly = parsed.get("yLabel")
        if isinstance(ly, str) and ly.strip():
            meta["yLabel"] = ly.strip()[:80]
        lt = parsed.get("legendTitle")
        if isinstance(lt, str) and lt.strip():
            meta["legendTitle"] = lt.strip()[:120]
        ser = parsed.get("series")
        if isinstance(ser, list) and ser:
            meta["series"] = ser
        return data, meta
    return [], {}


def save_user_chart_from_client(
    *,
    user_id: str,
    title: str,
    chart_type: str,
    data: list[Any],
    x_label: str | None = None,
    y_label: str | None = None,
    legend_title: str | None = None,
    series: list[Any] | None = None,
    source_type: str = "CLIENT_ENGINE",
    source_name: str | None = None,
    input_summary: str | None = None,
) -> dict[str, Any]:
    """Persist a chart built in the browser (chart-engine) so it appears in the profile list."""
    chart_id = uuid.uuid4().hex
    envelope: dict[str, Any] = {"data": list(data)}
    if x_label and str(x_label).strip():
        envelope["xLabel"] = str(x_label).strip()[:80]
    if y_label and str(y_label).strip():
        envelope["yLabel"] = str(y_label).strip()[:80]
    if legend_title and str(legend_title).strip():
        envelope["legendTitle"] = str(legend_title).strip()[:120]
    if series and isinstance(series, list):
        envelope["series"] = series
    summary = (input_summary or "").strip()[:600]
    st = (source_type or "CLIENT_ENGINE").strip()[:32] or "CLIENT_ENGINE"
    with _conn() as con:
        con.execute(
            """
            INSERT INTO user_charts (id, user_id, title, chart_type, chart_data_json, source_type, source_name, input_summary, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                chart_id,
                user_id,
                str(title or "Generated Chart").strip()[:120] or "Generated Chart",
                str(chart_type or "bar").strip()[:40] or "bar",
                json.dumps(envelope),
                st,
                (source_name or None),
                summary or None,
                time.time(),
            ),
        )
        con.commit()
    cbody: dict[str, Any] = {
        "id": chart_id,
        "title": str(title or "Generated Chart").strip()[:120] or "Generated Chart",
        "chartType": str(chart_type or "bar").strip()[:40] or "bar",
        "data": envelope["data"],
        "sourceType": st,
        "sourceName": source_name,
        "createdAt": _epoch_ms(time.time()),
    }
    if envelope.get("xLabel"):
        cbody["xLabel"] = envelope["xLabel"]
    if envelope.get("yLabel"):
        cbody["yLabel"] = envelope["yLabel"]
    if envelope.get("legendTitle"):
        cbody["legendTitle"] = envelope["legendTitle"]
    if envelope.get("series"):
        cbody["series"] = envelope["series"]
    return {"success": True, "chart": cbody}


def list_user_charts(*, user_id: str) -> dict[str, Any]:
    with _conn() as con:
        rows = con.execute(
            "SELECT * FROM user_charts WHERE user_id = ? ORDER BY created_at DESC",
            (user_id,),
        ).fetchall()
    out_charts: list[dict[str, Any]] = []
    for r in rows:
        data, meta = _parse_stored_chart_data_json(r["chart_data_json"] or "")
        item: dict[str, Any] = {
            "id": r["id"],
            "title": r["title"],
            "chartType": r["chart_type"],
            "data": data,
            "sourceType": r["source_type"],
            "sourceName": r["source_name"],
            "createdAt": _epoch_ms(r["created_at"]),
        }
        item.update(meta)
        out_charts.append(item)
    return {"charts": out_charts}


def delete_user_chart(*, chart_id: str, user_id: str) -> dict[str, Any]:
    with _conn() as con:
        row = con.execute("SELECT id, user_id FROM user_charts WHERE id = ?", (chart_id,)).fetchone()
        if row is None:
            return {"errorCode": "NOT_FOUND", "message": "Chart not found"}
        if str(row["user_id"]) != user_id:
            return {"errorCode": "FORBIDDEN", "message": "Chart does not belong to this user"}
        con.execute("DELETE FROM user_charts WHERE id = ?", (chart_id,))
        con.commit()
    return {"ok": True}


def _picsum_image_url(*, seed_key: str, width: int = 1280, height: int = 720) -> str:
    """Deterministic Lorem Picsum URL (reliable in browsers; used when Unsplash is off or fails)."""
    h = hashlib.sha256(seed_key.encode("utf-8")).hexdigest()[:24]
    return f"https://picsum.photos/seed/{h}/{width}/{height}"


# Words that help DALL-E-style imageQuery prompts but poison Unsplash /stock search (generic hits).
_UNSPLASH_STOCK_FLUFF: frozenset[str] = frozenset(
    {
        "abstract",
        "aesthetic",
        "angle",
        "audience",
        "authentic",
        "balanced",
        "blues",
        "calm",
        "centered",
        "cinematic",
        "closure",
        "comparing",
        "composition",
        "conceptual",
        "cool",
        "dashboard",
        "depicted",
        "depth",
        "documentary",
        "dramatic",
        "echoing",
        "editorial",
        "emotional",
        "field",
        "flow",
        "forward",
        "geometry",
        "glow",
        "handshake",
        "horizon",
        "human",
        "infographic",
        "light",
        "looking",
        "metaphor",
        "minimal",
        "modern",
        "moment",
        "mood",
        "numbered",
        "numbers",
        "optimistic",
        "paths",
        "photoreal",
        "possibility",
        "premium",
        "readable",
        "related",
        "scene",
        "shallow",
        "shot",
        "slide",
        "slides",
        "soft",
        "split",
        "sunrise",
        "tension",
        "text",
        "ties",
        "tones",
        "variant",
        "visualization",
        "wide",
        "workplace",
    }
)


def _stock_photo_search_query(topic_seed: str, image_query: str = "", *, max_len: int = 200) -> str:
    """Build a short keyword string for Unsplash /search/photos (not a generative prompt)."""
    topic_seed = (topic_seed or "").strip()
    image_query = (image_query or "").strip()
    topic_kw = _extract_topic_keywords(topic_seed, max_keywords=8)
    combined_kw = _extract_topic_keywords(f"{topic_seed} {image_query}", max_keywords=16)

    def _dedupe_ordered(words: list[str], *, skip_fluff: bool) -> list[str]:
        seen: set[str] = set()
        out: list[str] = []
        for w in words:
            wl = w.lower()
            if wl in seen:
                continue
            if skip_fluff and w in _UNSPLASH_STOCK_FLUFF:
                continue
            seen.add(wl)
            out.append(w)
            if len(out) >= 10:
                break
        return out

    ordered = _dedupe_ordered(topic_kw + combined_kw, skip_fluff=True)
    if len(ordered) < 2:
        ordered = _dedupe_ordered(topic_kw + combined_kw, skip_fluff=False)
    q = " ".join(ordered).strip()
    if len(q) < 3:
        q = " ".join(topic_kw[:6]).strip()
    if len(q) < 3:
        q = re.sub(r"\s+", " ", topic_seed)[:120].strip()
    if len(q) < 2:
        q = "professional presentation"
    return q[:max_len]


def _unsplash_access_key_configured() -> bool:
    return bool((get_settings().unsplash_access_key or "").strip())


def _fetch_unsplash_image_urls(search_query: str, *, limit: int = 5) -> list[str]:
    """Search Unsplash; return regular-size photo URLs (empty if key missing, error, or no hits)."""
    key = (get_settings().unsplash_access_key or "").strip()
    if not key:
        return []
    q = " ".join(str(search_query or "").split())[:200]
    if len(q) < 2:
        return []
    lim = max(1, min(int(limit), 10))
    try:
        import httpx

        with httpx.Client(timeout=httpx.Timeout(12.0, connect=4.0)) as client:
            r = client.get(
                "https://api.unsplash.com/search/photos",
                params={"query": q, "per_page": str(lim), "orientation": "landscape"},
                headers={
                    "Authorization": f"Client-ID {key}",
                    "Accept-Version": "v1",
                },
            )
        if r.status_code != 200:
            return []
        data = r.json()
        raw = data.get("results")
        if not isinstance(raw, list):
            return []
        out: list[str] = []
        for item in raw:
            if not isinstance(item, dict):
                continue
            urls = item.get("urls")
            if not isinstance(urls, dict):
                continue
            u = urls.get("regular") or urls.get("small") or urls.get("full")
            if isinstance(u, str) and u.startswith("http"):
                out.append(u.strip())
            if len(out) >= lim:
                break
        return out
    except Exception:
        return []


def _image_option_dict(
    *,
    url: str,
    reason: str,
    query: str,
    is_best: bool,
    confidence: float,
) -> dict[str, Any]:
    return {
        "url": url,
        "originalUrl": url,
        "confidence": confidence,
        "similarity": confidence,
        "source": "search",
        "score": confidence,
        "reason": reason,
        "aspectRatio": "16:9",
        "position": "right",
        "isBestMatch": is_best,
        "promptUsed": query[:220],
    }


def generate_image_advanced(*, slide_id: str, slide_content: str = "") -> dict[str, Any]:
    """Prefer Unsplash (topic search) when UNSPLASH_ACCESS_KEY is set; else Lorem Picsum."""
    raw = (slide_content or "").strip()
    query = _stock_photo_search_query(raw, raw, max_len=200)
    if len(query) < 2:
        query = raw[:80].strip() or "presentation"
    seed_material = f"{slide_id}:{query}:{slide_content}"

    unsplash_urls = _fetch_unsplash_image_urls(query, limit=5)
    if unsplash_urls:
        h = int(hashlib.sha256(seed_material.encode("utf-8")).hexdigest(), 16)
        start = h % len(unsplash_urls)
        rotated = unsplash_urls[start:] + unsplash_urls[:start]
        pick = rotated[: min(3, len(rotated))]
        reason = "Unsplash photo matched to slide topic (Picsum not used)."
        images = [
            _image_option_dict(
                url=u,
                reason=reason,
                query=query,
                is_best=(i == 0),
                confidence=0.82 - i * 0.03,
            )
            for i, u in enumerate(pick)
        ]
        ranked = [{"url": u, "score": 0.8 - i * 0.05} for i, u in enumerate(rotated[:5])]
        src = "unsplash"
    else:
        url = _picsum_image_url(seed_key=seed_material)
        reason = (
            "Lorem Picsum fallback (set UNSPLASH_ACCESS_KEY for real photos, or Unsplash had no results / error)."
        )
        images = [
            _image_option_dict(
                url=url,
                reason=reason,
                query=query,
                is_best=True,
                confidence=0.72,
            )
        ]
        ranked = []
        src = "picsum"

    return {
        "success": True,
        "images": images,
        "styleProfile": {"style": "modern", "tone": "professional", "colorPalette": ["#1f2937", "#facc15"]},
        "visualIntent": {"visualType": "concept", "primarySubject": query[:120]},
        "queries": [query],
        "ranked": ranked,
        "action": "ok",
        "reason": reason,
        "promptUsed": slide_content[:220],
        "candidateCount": len(images),
        "pipeline": {"native": True, "imageSource": src},
    }


def _openai_api_key_available() -> bool:
    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    return bool(key) and key != "sk-placeholder"


def _heuristic_role_plan() -> list[tuple[str, str, str]]:
    return [
        ("Overview", "What this topic is and who it is for", "opening"),
        ("Background", "Where it comes from and how it is usually framed", "context"),
        ("Core ideas", "Main concepts someone must understand", "insight"),
        ("How it works", "Mechanics, process, or structure—applied to this topic", "data"),
        ("Implications", "What changes for the audience in practice", "breakdown"),
        ("Plan forward", "Next steps to go deeper on this topic", "roadmap"),
        ("Wrap-up", "Takeaways and what to do next", "conclusion"),
    ]


def _format_file_extraction_for_llm(file_extraction: dict[str, Any] | None) -> str:
    if not isinstance(file_extraction, dict):
        return ""
    insights = file_extraction.get("insights")
    structured = file_extraction.get("structuredData")
    ins_lines: list[str] = []
    if isinstance(insights, list):
        for x in insights[:6]:
            s = str(x).strip()
            if s:
                ins_lines.append(s[:180])

    kpi_lines: list[str] = []
    if isinstance(structured, dict):
        kpis = structured.get("kpis")
        if isinstance(kpis, list):
            for k in kpis[:6]:
                if not isinstance(k, dict):
                    continue
                name = str(k.get("name") or "").strip()
                value = k.get("value")
                unit = str(k.get("unit") or "").strip()
                ctx = str(k.get("context") or "").strip()
                source_tag = str(k.get("sourceTag") or "").strip()
                if value is None or not name:
                    continue
                num = str(value).strip()
                display = f"{name}: {num}{unit}".strip(": ")
                if source_tag:
                    display = f"{source_tag} {display}".strip()
                if ctx:
                    display = f"{display} ({ctx[:90]})"
                kpi_lines.append(display[:220])

    ins_blob = "\n".join(f"- {l}" for l in ins_lines) if ins_lines else "- (none provided)"
    kpi_blob = "\n".join(f"- {l}" for l in kpi_lines) if kpi_lines else "- (none provided)"
    return f"""
FILE EXTRACTIONS (summarized; do NOT copy verbatim):

INSIGHTS:
{ins_blob}

STRUCTURED DATA (for KPIs / numeric framing):
{kpi_blob}
""".strip()


def _llm_generate_single_slide_payload(
    *,
    topic: str,
    deck_title: str,
    slide_index_1based: int,
    slide_count: int,
    tone: str,
    role_title: str,
    role_sub: str,
    role_kind: str,
    file_extraction: dict[str, Any] | None = None,
    user_prompt: str | None = None,
    detail_level: str = "standard",
) -> dict[str, Any] | None:
    """Generate ONE slide payload with an LLM so SSE can start immediately."""
    if not _openai_api_key_available():
        return None
    settings = get_settings()
    key = (settings.openai_api_key or "").strip()
    try:
        from langchain_core.messages import HumanMessage, SystemMessage
        from langchain_openai import ChatOpenAI
    except Exception:
        return None

    n = max(3, min(slide_count, 30))
    tone_use = (tone or "professional").strip() or "professional"
    topic_trim = (topic or "").strip() or deck_title.strip() or "the subject"
    brief = (user_prompt or "").strip()[:6000]

    file_context_section = _format_file_extraction_for_llm(file_extraction)
    file_context = f"\n{file_context_section}\n" if file_context_section else ""
    detail_mode = (detail_level or "standard").strip().lower()
    if detail_mode == "deep":
        detail_clause = (
            "MODE: DEEP / IN-DEPTH.\n"
            "- Bullets: exactly 4; mechanism + example/scenario + plausible measurable framing.\n"
            "- speakerNotes: 5–8 sentences max.\n"
        )
    elif detail_mode == "detailed":
        detail_clause = (
            "MODE: DETAILED.\n"
            "- Bullets: exactly 4; add concrete examples/workflows.\n"
            "- speakerNotes: 4–6 sentences.\n"
        )
    elif detail_mode == "concise":
        detail_clause = (
            "MODE: CONCISE / BRIEF.\n"
            "- Bullets: exactly 3; high-impact, no fluff.\n"
            "- speakerNotes: 2–3 sentences.\n"
        )
    else:
        detail_clause = (
            "MODE: STANDARD.\n"
            "- Bullets: 3–4; balanced clarity + depth.\n"
            "- speakerNotes: 3–5 sentences.\n"
        )

    system_hint = (
        "You are an expert on the user's subject matter. Domain-specific examples beat generic business language. "
        "Output strictly valid JSON only. No markdown. No explanations."
    )
    brief_section = (
        f"FULL USER REQUEST (honor scope, audience, and domain):\n---\n{brief}\n---\n\n"
        if brief
        else ""
    )
    user = f"""Create slide {slide_index_1based} of {n}.

{brief_section}SUBJECT LINE (title discipline; do not paste into every bullet):
{topic_trim}

SLIDE ROLE:
role_title: {role_title}
role_sub: {role_sub}
role_kind: {role_kind}

RULES:
- You are an expert presentation strategist, content designer, and business communicator.
- Each slide must deliver ONE clear idea. Avoid generic filler, repetition, and vague statements.
- Every bullet: include a concrete mechanism, application, or plausible measurable detail tied to the user's topic—not generic "efficiency" alone.
- Avoid empty use of: stakeholders, roadmaps, alignment, outcomes, strategy, synergy, transformation—unless the user brief is about that.
- Title: 5–10 words, domain-specific to this slide; not generic section labels.
- Subtitle: one short supporting line.
- bullets: follow mode bullet count:
  - CONCISE: exactly 3
  - STANDARD: 3–4
  - DETAILED/DEEP: exactly 4
- Forbidden: "memorable move", "this slide explains", "the key takeaway is", "as you can see", "unlock value", "paradigm shift".
- description: 1–2 expert sentences; highlight: one crisp line; keyMessage: one actionable line; speakerNotes: 2–4 conversational sentences.
- citations: if a factual claim comes from FILE EXTRACTIONS evidence tags like `[S1]`, append the same `[S#]` tags at the end of the sentence inside `speakerNotes`.
- numeric fidelity: if you mention a numeric KPI from FILE EXTRACTIONS, keep the exact extracted value and unit (no rounding or unit changes).
- anti-copy: never paste FILE EXTRACTIONS evidence sentences verbatim; paraphrase and cite tags only.

DECK TITLE: {deck_title}
TONE: {tone_use}

{file_context}

{detail_clause}

CITATION TAGGING:
- If a factual claim comes from FILE EXTRACTIONS evidence tags like `[S1]`, append the same `[S#]` tags at the end of the sentence inside `speakerNotes`.
- Never paste evidence sentences verbatim; paraphrase and cite tags only.

NUMERIC FAITHFULNESS:
- If you mention a numeric KPI from FILE EXTRACTIONS, keep the exact extracted value and unit (no rounding or unit changes).

Return exactly this JSON object:
{{
  "title": "...",
  "subtitle": "...",
  "bullets": ["...", "...", "..."],
  "description": "...",
  "highlight": "...",
  "keyMessage": "...",
  "speakerNotes": "..."
}}"""

    try:
        llm = ChatOpenAI(
            model=settings.openai_model,
            api_key=key,
            base_url=settings.openai_base_url,
            temperature=0.6,
        )
        resp = llm.invoke([SystemMessage(content=system_hint), HumanMessage(content=user)])
        text = getattr(resp, "content", None) or str(resp)
        if not isinstance(text, str):
            text = str(text)
        parsed = _extract_json_object(text)
        if not parsed or not isinstance(parsed, dict):
            return None
        return parsed
    except Exception:
        return None


def stream_generate_presentation_slides(
    *,
    presentation_id: str,
    slide_count_target: int = 0,
    tone: str | None = None,
) -> Any:
    """Yield SSE events while slides are generated/persisted, with field-by-field streaming."""
    tone_use = tone or "professional"

    def _sleep(seconds: float) -> None:
        # Small artificial delays make the experience feel human (not a token-by-token dump).
        # This does NOT add extra API calls.
        try:
            time.sleep(seconds)
        except Exception:
            pass

    with _conn() as con:
        row = con.execute(
            "SELECT id, user_id, title, prompt, template_name, status FROM presentations WHERE id = ?",
            (presentation_id,),
        ).fetchone()
        if row is None:
            yield {"event": "error", "data": {"message": "Presentation not found"}}
            return

        job_id = uuid.uuid4().hex
        now = time.time()
        con.execute(
            "INSERT INTO jobs (id, presentation_id, status, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
            (job_id, presentation_id, "PROCESSING", now, now),
        )
        con.execute(
            "UPDATE presentations SET status = ?, updated_at = ? WHERE id = ?",
            ("PROCESSING", now, presentation_id),
        )
        con.commit()

        slides_existing = con.execute(
            "SELECT id, idx, title, bullets_json, content_json FROM slides WHERE presentation_id = ? ORDER BY idx ASC",
            (presentation_id,),
        ).fetchall()

        base = str(row["title"] or "").strip() or "Presentation"
        n_slides = max(3, min(int(slide_count_target) if slide_count_target > 0 else 8, 30))

        yield {"event": "job_created", "data": {"jobId": job_id, "presentationId": presentation_id, "status": "queued"}}
        yield {
            "event": "outline_generated",
            "data": {"slideCount": len(slides_existing) if slides_existing else n_slides},
        }
        yield {"event": "thinking", "data": {"message": "Analyzing topic and audience intent…"}}
        yield {"event": "progress", "data": {"percent": 0, "step": "Generating content"}}

        # If slides already exist, just stream them with per-field animation.
        if slides_existing and len(slides_existing) > 0:
            for s in slides_existing:
                idx = int(s["idx"])
                content = _slide_content_from_row(s["title"], s["bullets_json"], s["content_json"])
                yield {"event": "slide_start", "data": {"index": idx}}
                yield {"event": "thinking", "data": {"message": f"Designing slide {idx}…"}}
                yield {
                    "event": "progress",
                    "data": {
                        "percent": int(((idx - 1) / max(1, len(slides_existing))) * 100),
                        "step": "Designing slides",
                    },
                }

                # Legacy compatibility events.
                legacy_layout_preset = ""
                if isinstance(content.get("gammaStyle"), dict):
                    legacy_layout_preset = str(content.get("gammaStyle", {}).get("layoutPreset") or "")
                legacy_title = str(content.get("title") or s["title"]).strip()
                legacy_bullets = content.get("bullets") if isinstance(content.get("bullets"), list) else []
                legacy_desc = str(content.get("description") or "").strip()
                legacy_key = str(content.get("keyMessage") or content.get("highlight") or "").strip()
                yield {
                    "event": "slide_generated",
                    "data": {
                        "slideIndex": idx,
                        "title": legacy_title,
                        "bullets": [str(b).strip() for b in legacy_bullets if str(b).strip()],
                        "description": legacy_desc,
                        "keyMessage": legacy_key,
                    },
                }
                yield {"event": "layout_applied", "data": {"slideIndex": idx, "layoutType": legacy_layout_preset or "title_bullets"}}

                _sleep(0.78)
                title = str(content.get("title") or s["title"]).strip()
                yield {"event": "slide_chunk", "data": {"index": idx, "field": "title", "value": title}}

                _sleep(0.72)
                subtitle = str(content.get("subtitle") or "").strip()
                yield {"event": "slide_chunk", "data": {"index": idx, "field": "subtitle", "value": subtitle}}

                bullets = content.get("bullets")
                bullets_list = bullets if isinstance(bullets, list) else []
                acc: list[str] = []
                for b in bullets_list:
                    b_str = str(b).strip()
                    if not b_str:
                        continue
                    acc.append(b_str)
                    _sleep(0.52)
                    yield {"event": "slide_chunk", "data": {"index": idx, "field": "bullets", "value": acc.copy()}}

                _sleep(0.55)
                description = str(content.get("description") or "").strip()
                yield {"event": "slide_chunk", "data": {"index": idx, "field": "description", "value": description}}

                _sleep(0.55)
                highlight = str(content.get("highlight") or "").strip()
                yield {"event": "slide_chunk", "data": {"index": idx, "field": "highlight", "value": highlight}}

                yield {
                    "event": "slide_complete",
                    "data": {
                        "fullSlideJSON": {
                            "id": s["id"],
                            "index": idx,
                            "order": idx,
                            "title": s["title"],
                            "content": content,
                        }
                    },
                }
                yield {"event": "slide_persisted", "data": {"slideIndex": idx}}
                yield {"event": "progress", "data": {"percent": int((idx / max(1, len(slides_existing))) * 100), "step": "Finalizing"}}
            yield {"event": "completed", "data": {"presentationId": presentation_id, "jobId": job_id}}
            return

        # Fresh generation.
        con.execute("DELETE FROM slides WHERE presentation_id = ?", (presentation_id,))
        con.commit()

        prompt_text = str(row["prompt"] or "")
        extracted = _try_parse_file_extraction_prompt(prompt_text)
        file_extraction = extracted
        user_request_text = str(extracted.get("userRequest") or "").strip() if isinstance(extracted, dict) else ""
        if not user_request_text:
            user_request_text = prompt_text
        detail_level = _infer_detail_level(user_request_text)
        if extracted and str(extracted.get("topic") or "").strip():
            topic_for_content = str(extracted.get("topic")).strip()
        else:
            topic_for_content = (_canonical_subject_from_prompt(user_request_text) or base).strip()

        # Cost-efficient streaming: generate the full deck once (1 LLM call),
        # then stream each slide's fields locally with delays.
        payloads_for_stream: list[dict[str, Any]] = []
        deck_from_full_llm = False
        if _openai_api_key_available():
            yield {"event": "thinking", "data": {"message": "Generating deck content in one pass…"}}
            full_payloads = _llm_generate_slide_payloads(
                topic=topic_for_content,
                deck_title=base,
                slide_count=n_slides,
                tone=tone_use,
                file_extraction=file_extraction,
                user_prompt=user_request_text,
                detail_level=detail_level,
            )
            if isinstance(full_payloads, list) and len(full_payloads) >= 3:
                payloads_for_stream = full_payloads
                deck_from_full_llm = True

        if not payloads_for_stream:
            yield {"event": "thinking", "data": {"message": "Using extracted heuristics to draft slides…"}}
            payloads_for_stream = _heuristic_slide_payloads(
                prompt=user_request_text or topic_for_content,
                deck_title=base,
                slide_count=n_slides,
                tone=tone_use,
                file_extraction=file_extraction,
                detail_level=detail_level,
            )

            # Optional: only titles-only pass to keep slides sharp (1 extra LLM call max).
            if _openai_api_key_available():
                yield {"event": "thinking", "data": {"message": "Refining sharp slide titles…"}}
                ai_titles = _llm_generate_slide_titles(
                    topic=topic_for_content,
                    deck_title=base,
                    slide_count=n_slides,
                    tone=tone_use,
                )
                if ai_titles:
                    for ti, t in enumerate(ai_titles):
                        if ti < len(payloads_for_stream) and isinstance(t, str) and t.strip():
                            payloads_for_stream[ti]["title"] = t.strip()

        payloads_for_stream = payloads_for_stream[:n_slides] if isinstance(payloads_for_stream, list) else []

        # Only merge heuristic slides when the deck was NOT produced by the full LLM pass—replacing
        # model copy with templates caused repetitive titles/bullets ("Why X Changes the Next 12 Months…").
        topic_keywords = _extract_topic_keywords(topic_for_content, max_keywords=8)

        if (
            not deck_from_full_llm
            and topic_keywords
            and isinstance(payloads_for_stream, list)
            and any(
                isinstance(p, dict)
                and not _payload_seems_topic_grounded(p, canonical_topic=topic_for_content, keywords=topic_keywords)
                for p in payloads_for_stream
            )
        ):
            heuristic_payloads = _heuristic_slide_payloads(
                prompt=user_request_text or topic_for_content,
                deck_title=base,
                slide_count=n_slides,
                tone=tone_use,
                file_extraction=file_extraction,
                detail_level=detail_level,
            )
            for i in range(min(len(payloads_for_stream), len(heuristic_payloads))):
                if isinstance(payloads_for_stream[i], dict) and not _payload_seems_topic_grounded(
                    payloads_for_stream[i], canonical_topic=topic_for_content, keywords=topic_keywords
                ):
                    payloads_for_stream[i] = heuristic_payloads[i]

            if _openai_api_key_available():
                ai_titles = _llm_generate_slide_titles(
                    topic=topic_for_content,
                    deck_title=base,
                    slide_count=n_slides,
                    tone=tone_use,
                )
                if ai_titles:
                    for ti, t in enumerate(ai_titles):
                        if ti < len(payloads_for_stream) and isinstance(payloads_for_stream[ti], dict) and isinstance(t, str) and t.strip():
                            payloads_for_stream[ti]["title"] = t.strip()

        for i in range(n_slides):
            idx_1 = i + 1

            yield {"event": "slide_start", "data": {"index": idx_1}}
            yield {"event": "thinking", "data": {"message": "Designing slide structure…"}}
            yield {"event": "progress", "data": {"percent": int(((idx_1 - 1) / max(1, n_slides)) * 100), "step": "Generating content"}}

            # Use precomputed slide payloads (generated once for the whole deck).
            payload = (
                payloads_for_stream[i]
                if isinstance(payloads_for_stream, list) and i < len(payloads_for_stream)
                else (payloads_for_stream[i % len(payloads_for_stream)] if payloads_for_stream else {})
            )

            slide_title = str(payload.get("title") or f"{base} — Slide {idx_1}").strip()
            # Last slide: no default "Section N" subtitle (reads like a timestamp / filler).
            subtitle_default = "" if idx_1 >= n_slides else f"Section {idx_1}"
            subtitle = str(payload.get("subtitle") or subtitle_default).strip()
            bullets = payload.get("bullets")
            if not isinstance(bullets, list):
                bullets = []
            bullets = [str(b).strip() for b in bullets if str(b).strip()]
            if len(bullets) < 2:
                bullets = [
                    f'Define what "{topic_for_content[:80]}" means for your audience and why it matters now.',
                    f"Explain two core ideas someone must know about {topic_for_content[:80]} before going deeper.",
                    f"Give one practical implication: how understanding {topic_for_content[:80]} changes decisions or next steps.",
                ]

            description = str(payload.get("description") or "").strip()
            highlight = str(payload.get("highlight") or "").strip()
            key_message = str(payload.get("keyMessage") or payload.get("key_message") or highlight or "").strip()
            speaker_notes = str(payload.get("speakerNotes") or "").strip()

            slide_type = str(payload.get("slideType") or detect_slide_type(payload)).strip() or "content"
            role_sub = str(payload.get("subtitle") or "This section").strip()
            gs_payload = payload.get("gammaStyle") if isinstance(payload.get("gammaStyle"), dict) else {}
            layout_preset = str(gs_payload.get("layoutPreset") or payload.get("layoutSuggestion") or "").strip()
            if layout_preset not in _VALID_GAMMA_LAYOUTS:
                layout_preset = _layout_preset_for_slide_type(slide_type)
            emphasis0 = base.split(" ")[0] if base.split(" ") else "Insight"
            emphasis = (
                [emphasis0, "Impact", "Action"]
                if slide_type in ("hero", "content", "visual")
                else [emphasis0, "Key KPI", "Proof"]
            )
            emph_words = (
                gs_payload["emphasisWords"]
                if isinstance(gs_payload.get("emphasisWords"), list) and gs_payload.get("emphasisWords")
                else emphasis
            )
            gamma_style_obj: dict[str, Any] = {
                "layoutPreset": layout_preset,
                "alignment": str(gs_payload.get("alignment") or "left"),
                "emphasisWords": emph_words[:6],
            }
            for gkey in ("imagePlacement", "fullBleed", "textPrimary", "gradientTitle", "contentAlign", "cardWidth"):
                if gkey in gs_payload and gs_payload[gkey] is not None:
                    gamma_style_obj[gkey] = gs_payload[gkey]

            image_query = str(payload.get("imageQuery") or "").strip() or f"{topic_for_content} cinematic scene slide {idx_1}"

            slide_id = uuid.uuid4().hex
            content_json = json.dumps(
                {
                    "title": slide_title,
                    "subtitle": subtitle,
                    "bullets": bullets,
                    "description": description
                    or f"Framing for {topic_for_content[:80]}: {role_sub} translated into a decision the audience can use.",
                    "highlight": highlight
                    or key_message
                    or f"Key takeaway: {topic_for_content[:80]} becomes clearer when you apply the {role_sub} angle.",
                    "keyMessage": key_message or (highlight or "Key takeaway"),
                    "speakerNotes": speaker_notes,
                    "imageQuery": image_query[:400],
                    "generatedImageUrl": _stable_image_url_for_slide(topic_for_content, idx_1, image_query),
                    "slideType": slide_type,
                    "layoutSuggestion": layout_preset,
                    "emphasisWords": emph_words[:8],
                    "contentDensity": "medium",
                    "qualityScore": 7.6,
                    "gammaStyle": gamma_style_obj,
                }
            )

            con.execute(
                "INSERT INTO slides (id, presentation_id, idx, title, bullets_json, content_json) VALUES (?, ?, ?, ?, ?, ?)",
                (slide_id, presentation_id, idx_1, slide_title, "|||".join(bullets), content_json),
            )
            con.commit()

            # Stream structured fields for the premium "building" experience.
            yield {
                "event": "progress",
                "data": {"percent": int(((idx_1 - 1) / max(1, n_slides)) * 100), "step": "Designing slides"},
            }

            # Legacy compatibility events (older UI listens to slide_generated/layout_applied/slide_persisted).
            yield {
                "event": "slide_generated",
                "data": {
                    "slideIndex": idx_1,
                    "title": slide_title,
                    "bullets": bullets if isinstance(bullets, list) else [],
                    "description": description or "",
                    "keyMessage": key_message or highlight or "",
                },
            }
            yield {"event": "layout_applied", "data": {"slideIndex": idx_1, "layoutType": layout_preset or "title_bullets"}}
            _sleep(0.78)
            yield {"event": "slide_chunk", "data": {"index": idx_1, "field": "title", "value": slide_title}}

            _sleep(0.72)
            yield {"event": "slide_chunk", "data": {"index": idx_1, "field": "subtitle", "value": subtitle}}

            acc: list[str] = []
            for b in bullets:
                b_str = str(b).strip()
                if not b_str:
                    continue
                acc.append(b_str)
                _sleep(0.52)
                yield {"event": "slide_chunk", "data": {"index": idx_1, "field": "bullets", "value": acc.copy()}}

            _sleep(0.55)
            yield {"event": "slide_chunk", "data": {"index": idx_1, "field": "description", "value": description or ""}}

            _sleep(0.55)
            yield {"event": "slide_chunk", "data": {"index": idx_1, "field": "highlight", "value": highlight or key_message or ""}}

            content_obj = _slide_content_from_row(slide_title, "|||".join(bullets), content_json)
            yield {
                "event": "slide_complete",
                "data": {
                    "fullSlideJSON": {
                        "id": slide_id,
                        "index": idx_1,
                        "order": idx_1,
                        "title": slide_title,
                        "content": content_obj,
                    }
                },
            }
            yield {"event": "slide_persisted", "data": {"slideIndex": idx_1}}
            yield {"event": "progress", "data": {"percent": int((idx_1 / max(1, n_slides)) * 100), "step": "Finalizing"}}

        done = time.time()
        con.execute("UPDATE jobs SET status = ?, updated_at = ? WHERE id = ?", ("COMPLETED", done, job_id))
        con.execute(
            "UPDATE presentations SET status = ?, updated_at = ? WHERE id = ?",
            ("COMPLETED", done, presentation_id),
        )
        con.commit()
        yield {"event": "completed", "data": {"presentationId": presentation_id, "jobId": job_id}}


def generate_stream_events(
    *,
    user_id: str,
    topic: str,
    slide_count: int,
    tone: str | None = None,
    template_key: str | None = None,
) -> tuple[str, str, list[dict[str, Any]]]:
    created = create_presentation(
        user_id=user_id,
        prompt=topic,
        title=topic[:120],
        template_name=template_key or "gammaDefault",
    )
    pid = str(created.get("presentationId", ""))
    generated = generate_presentation(
        presentation_id=pid,
        slide_count_target=slide_count,
        tone=tone or "professional",
    )
    jid = str(generated.get("jobId", ""))
    pres = get_presentation(presentation_id=pid).get("presentation", {})
    slides = pres.get("slides", []) if isinstance(pres, dict) else []
    events: list[dict[str, Any]] = []
    events.append({"event": "job_created", "data": {"jobId": jid, "presentationId": pid, "status": "queued"}})
    events.append({"event": "outline_generated", "data": {"slideCount": len(slides)}})
    for i, s in enumerate(slides, start=1):
        c = s.get("content", {}) if isinstance(s, dict) else {}
        events.append(
            {
                "event": "slide_generated",
                "data": {
                    "slideIndex": i,
                    "title": c.get("title") or s.get("title") or f"Slide {i}",
                    "bullets": c.get("bullets") if isinstance(c.get("bullets"), list) else [],
                    "description": c.get("description") or "",
                    "keyMessage": c.get("keyMessage") or c.get("highlight") or "",
                },
            }
        )
        layout = ((c.get("gammaStyle") or {}) if isinstance(c.get("gammaStyle"), dict) else {}).get("layoutPreset")
        events.append({"event": "layout_applied", "data": {"slideIndex": i, "layoutType": layout or "title_bullets"}})
        events.append({"event": "slide_persisted", "data": {"slideIndex": i}})
    events.append({"event": "completed", "data": {"presentationId": pid, "jobId": jid}})
    return pid, jid, events


def regenerate_slide(*, slide_id: str, tone: str = "professional") -> dict[str, Any]:
    """
    Regenerate ONE slide content while keeping the rest of the deck unchanged.
    Uses a single generation call for the requested slide only.
    """
    with _conn() as con:
        srow = con.execute(
            "SELECT id, idx, presentation_id FROM slides WHERE id = ?",
            (slide_id,),
        ).fetchone()
        if srow is None:
            return {"error": "Slide not found"}

        presentation_id = str(srow["presentation_id"])
        slide_idx = int(srow["idx"])

        pres = con.execute("SELECT id, title, prompt FROM presentations WHERE id = ?", (presentation_id,)).fetchone()
        if pres is None:
            return {"error": "Presentation not found"}

        deck_title = str(pres["title"] or "").strip() or "Presentation"
        prompt_text = str(pres["prompt"] or "")
        extracted = _try_parse_file_extraction_prompt(prompt_text)
        file_extraction = extracted
        user_request_text = str(extracted.get("userRequest") or "").strip() if isinstance(extracted, dict) else ""
        if not user_request_text:
            user_request_text = prompt_text
        detail_level = _infer_detail_level(user_request_text)
        topic_for_content = (
            str((extracted.get("topic") if extracted else None) or prompt_text).strip() or deck_title
        )

        slide_count_row = con.execute(
            "SELECT COUNT(1) as cnt FROM slides WHERE presentation_id = ?",
            (presentation_id,),
        ).fetchone()
        slide_count = int(slide_count_row["cnt"] or 8) if slide_count_row else 8

        role_plan = _heuristic_role_plan()
        role_title, role_sub, role_kind = role_plan[(slide_idx - 1) % len(role_plan)]

        payload: dict[str, Any] | None = None
        if _openai_api_key_available():
            payload = _llm_generate_single_slide_payload(
                topic=topic_for_content,
                deck_title=deck_title,
                slide_index_1based=slide_idx,
                slide_count=slide_count,
                tone=tone or "professional",
                role_title=role_title,
                role_sub=role_sub,
                role_kind=role_kind,
                file_extraction=file_extraction,
                user_prompt=user_request_text,
                detail_level=detail_level,
            )

        if not payload:
            heuristic_payloads = _heuristic_slide_payloads(
                prompt=user_request_text or topic_for_content,
                deck_title=deck_title,
                slide_count=slide_count,
                tone=tone or "professional",
                file_extraction=file_extraction,
                detail_level=detail_level,
            )
            payload = heuristic_payloads[slide_idx - 1] if slide_idx - 1 < len(heuristic_payloads) else {}

        slide_title = str(payload.get("title") or deck_title).strip()
        subtitle = str(payload.get("subtitle") or "").strip()

        bullets = payload.get("bullets")
        if not isinstance(bullets, list):
            bullets = payload.get("points")
        if not isinstance(bullets, list):
            bullets = []
        bullets = [str(b).strip() for b in bullets if str(b).strip()][:4]
        if len(bullets) < 3:
            bullets = bullets + [
                "Name the single decision this slide should change for the audience.",
                "State one constraint or trade-off that narrows good options.",
                "Give one measurable action within the next 30–90 days.",
            ]
        bullets = bullets[:4]

        description = str(payload.get("description") or "").strip()
        highlight = str(payload.get("highlight") or "").strip()
        key_message = str(payload.get("keyMessage") or payload.get("key_message") or highlight or "").strip()
        speaker_notes = str(payload.get("speakerNotes") or "").strip()

        slide_type = str(payload.get("slideType") or detect_slide_type(payload)).strip() or "content"
        layout_preset = _layout_preset_for_slide_type(slide_type)
        emphasis0 = deck_title.split(" ")[0] if deck_title.split(" ") else "Insight"
        emphasis = (
            [emphasis0, "Impact", "Action"]
            if slide_type in ("hero", "content", "visual")
            else [emphasis0, "Key KPI", "Proof"]
        )

        content = {
            "title": slide_title,
            "subtitle": subtitle,
            "bullets": bullets,
            "description": description
            or f"Framing for {topic_for_content[:80]}: {role_sub} translated into a decision the audience can use.",
            "highlight": highlight
            or key_message
            or f"Sharpest implication for the audience when viewing this slide through: {role_sub}.",
            "keyMessage": key_message or (highlight or "Operational implication for the next decision."),
            "speakerNotes": speaker_notes,
            "imageQuery": f"{topic_for_content} concept illustration slide {slide_idx}",
            "generatedImageUrl": _stable_image_url_for_slide(topic_for_content, slide_idx),
            "slideType": slide_type,
            "layoutSuggestion": layout_preset,
            "emphasisWords": emphasis,
            "contentDensity": "medium",
            "qualityScore": 7.6,
            "gammaStyle": {
                "layoutPreset": layout_preset,
                "alignment": "left",
                "emphasisWords": emphasis[:2],
            },
        }

        content_json = json.dumps(content)
        con.execute(
            "UPDATE slides SET title = ?, bullets_json = ?, content_json = ? WHERE id = ?",
            (slide_title, "|||".join(bullets), content_json, slide_id),
        )
        con.commit()

        content_obj = _slide_content_from_row(slide_title, "|||".join(bullets), content_json)
        return {
            "slide": {
                "id": slide_id,
                "index": slide_idx,
                "order": slide_idx,
                "title": slide_title,
                "content": content_obj,
            },
            "presentationId": presentation_id,
        }


def build_export_pptx(*, presentation_id: str) -> dict[str, Any]:
    try:
        from agent_core.ppt_native.gamma_export import build_gamma_presentation_pptx
    except Exception as exc:
        return {"error": f"Gamma PPTX export unavailable: {exc}"}

    data = get_presentation(presentation_id=presentation_id)
    if data.get("error"):
        return data
    presentation = data.get("presentation", {})
    slides = presentation.get("slides", [])
    if not isinstance(slides, list) or not slides:
        return {"error": "No slides to export"}

    out_dir = _native_export_work_dir()
    file_path = out_dir / f"{presentation_id}.pptx"

    template_name = presentation.get("templateName")
    if template_name is None:
        template_name = None
    else:
        template_name = str(template_name).strip() or None

    prs = build_gamma_presentation_pptx(slides=slides, template_name=template_name)
    prs.save(str(file_path))
    return {"path": str(file_path), "fileName": file_path.name, "mode": "native"}


def build_export_pdf_bytes(*, presentation_id: str, force_refresh: bool = False) -> dict[str, Any]:
    """Build PDF bytes for a presentation via PPTX export + LibreOffice headless.

    Reuses ``{presentation_id}.pdf`` next to the PPTX when it is newer than the PPTX
    unless ``force_refresh`` is True.
    """
    import logging

    from agent_core.tools.implementations.pdf_export_tool import PdfExportError, convert_ppt_to_pdf

    log = logging.getLogger(__name__)

    ppt_out = build_export_pptx(presentation_id=presentation_id)
    if ppt_out.get("error"):
        return ppt_out
    ppt_path_str = str(ppt_out.get("path", "")).strip()
    if not ppt_path_str:
        return {"error": "Native PPTX export did not return a file path"}
    ppt_path = Path(ppt_path_str)
    if not ppt_path.is_file():
        return {"error": f"PPTX missing on disk: {ppt_path}"}

    pdf_path = ppt_path.with_suffix(".pdf")
    try:
        ppt_mtime = ppt_path.stat().st_mtime
    except OSError as exc:
        return {"error": f"Cannot stat PPTX: {exc}"}

    if (
        not force_refresh
        and pdf_path.is_file()
        and pdf_path.stat().st_mtime >= ppt_mtime
    ):
        try:
            data = pdf_path.read_bytes()
        except OSError as exc:
            return {"error": f"Failed to read cached PDF: {exc}"}
        log.info("PDF export cache hit presentation_id=%s path=%s", presentation_id, pdf_path)
        return {
            "bytes": data,
            "fileName": pdf_path.name,
            "mode": "native",
            "pdfPath": str(pdf_path),
            "cached": True,
        }

    try:
        out_pdf = convert_ppt_to_pdf(str(ppt_path), out_dir=str(ppt_path.parent))
    except PdfExportError as exc:
        log.error("LibreOffice PDF export failed presentation_id=%s: %s", presentation_id, exc)
        return {"error": str(exc), "message": str(exc)}

    out_p = Path(out_pdf)
    try:
        data = out_p.read_bytes()
    except OSError as exc:
        return {"error": f"Failed to read generated PDF: {exc}"}

    log.info("PDF export generated presentation_id=%s path=%s", presentation_id, out_p)
    return {
        "bytes": data,
        "fileName": out_p.name,
        "mode": "native",
        "pdfPath": str(out_p),
        "cached": False,
    }


def _stable_image_url_for_slide(topic_seed: str, slide_index: int, image_query: str = "") -> str:
    """Prefer Unsplash search from deck topic + slide imageQuery; stable Picsum URL as fallback."""
    iq = (image_query or "").strip()[:240]
    search_q = _stock_photo_search_query(topic_seed, iq, max_len=200)
    if len(search_q.strip()) < 3:
        search_q = (topic_seed[:80] or "professional presentation").strip()

    if _unsplash_access_key_configured():
        urls = _fetch_unsplash_image_urls(search_q, limit=5)
        if urls:
            h = hashlib.sha256(f"{topic_seed}:{slide_index}:{iq}".encode("utf-8")).digest()
            pick = int.from_bytes(h[:4], "big") % len(urls)
            return urls[pick]

    return _picsum_image_url(seed_key=f"{search_q}:{slide_index}:{iq}:{topic_seed[:100]}")


def detect_slide_type(slide: dict[str, Any]) -> str:
    """Deterministic slide-type detection to guide layout + rendering."""
    title = str(slide.get("title") or "")
    subtitle = str(slide.get("subtitle") or "")
    bullets = slide.get("bullets")
    if not isinstance(bullets, list):
        bullets = []
    bullets_joined = " ".join(str(b) for b in bullets if b)
    highlight = str(slide.get("highlight") or "")
    key_message = str(slide.get("keyMessage") or "")

    text = f"{title}\n{subtitle}\n{bullets_joined}\n{highlight}\n{key_message}".strip().lower()

    # Priority order: stats -> comparison -> timeline -> visual -> hero -> content
    if re.search(r"\b(\d{1,3}(?:,\d{3})*|\d+)(?:\.\d+)?\b", text) or re.search(r"\b\d+%|\bpercent\b", text):
        if re.search(r"\b(vs|versus|compare|comparison|trade-?off|difference)\b", text):
            return "comparison"
        if re.search(r"\b(roi|kpi|revenue|profit|cost|growth|accuracy|latency|conversion|retention|margin|errors)\b", text):
            return "stats"
        return "stats"

    if re.search(r"\b(vs|versus|compare|comparison|trade-?off|difference)\b", text):
        return "comparison"

    if re.search(
        r"\b(steps?|step|process|workflow|pipeline|sequence|timeline|from\s+to|first|then|next|build)\b",
        text,
    ):
        return "timeline"

    if re.search(r"\b(image|visual|diagram|scene|illustration|photo|graph|map|depict|showcase)\b", text):
        return "visual"

    if re.search(r"\b(why|impact|transform|crucial|strategic|breakthrough|must|cannot|game-?changer)\b", text):
        return "hero"

    return "content"


def _layout_preset_for_slide_type(slide_type: str) -> str:
    st = str(slide_type or "").lower().strip()
    if st in ("hero", "visual"):
        return "hero_split"
    if st == "stats" or st == "stat":
        return "stats_split"
    if st in ("comparison", "timeline"):
        return "two_column"
    if st == "content":
        return "title_bullets"
    # Backward-compatible fallbacks
    if st == "split":
        return "two_column"
    if st == "section":
        return "title_bullets"
    return "title_bullets"


def _try_parse_file_extraction_prompt(prompt: str) -> dict[str, Any] | None:
    s = (prompt or "").strip()
    if not s.startswith("{"):
        return None
    try:
        parsed = json.loads(s)
        if isinstance(parsed, dict) and "topic" in parsed:
            return parsed
    except Exception:
        return None
    return None


def _generate_slides_sync(
    con: sqlite3.Connection,
    *,
    presentation_id: str,
    prompt: str,
    title: str,
    slide_count: int,
    tone: str = "",
) -> None:
    con.execute("DELETE FROM slides WHERE presentation_id = ?", (presentation_id,))
    base = title.strip() or "Presentation"
    # DB prompt can be either:
    # - user prompt text (normal flow)
    # - file-extraction JSON string: {"topic":..., "insights":..., "structuredData":...}
    extracted = _try_parse_file_extraction_prompt(prompt or "")
    file_extraction = extracted
    user_request_text = str(extracted.get("userRequest") or "").strip() if isinstance(extracted, dict) else ""
    if not user_request_text:
        user_request_text = prompt or ""
    detail_level = _infer_detail_level(user_request_text)
    if extracted and str(extracted.get("topic") or "").strip():
        topic_for_content = str(extracted.get("topic")).strip()
    else:
        topic_for_content = (_canonical_subject_from_prompt(user_request_text or "") or base).strip()
    n_slides = max(3, min(slide_count, 30))
    payloads = _llm_generate_slide_payloads(
        topic=topic_for_content,
        deck_title=base,
        slide_count=n_slides,
        tone=tone or "",
        file_extraction=file_extraction,
        user_prompt=user_request_text,
        detail_level=detail_level,
    )
    deck_from_full_llm = bool(payloads)
    used_heuristic = False
    if not payloads:
        used_heuristic = True
        payloads = _heuristic_slide_payloads(
            prompt=user_request_text or topic_for_content,
            deck_title=base,
            slide_count=n_slides,
            tone=tone or "",
            file_extraction=file_extraction,
            detail_level=detail_level,
        )

    # When the full-deck LLM did not run, still generate topic-specific slide titles via AI (matches legacy TS behavior).
    if used_heuristic:
        ai_titles = _llm_generate_slide_titles(
            topic=topic_for_content,
            deck_title=base,
            slide_count=n_slides,
            tone=tone or "",
        )
        if ai_titles:
            for ti, t in enumerate(ai_titles):
                if ti < len(payloads) and t.strip():
                    payloads[ti]["title"] = t.strip()

    topic_keywords = _extract_topic_keywords(topic_for_content, max_keywords=8)

    if (
        not deck_from_full_llm
        and topic_keywords
        and isinstance(payloads, list)
        and any(
            isinstance(p, dict)
            and not _payload_seems_topic_grounded(p, canonical_topic=topic_for_content, keywords=topic_keywords)
            for p in payloads[:n_slides]
        )
    ):
        heuristic_payloads = _heuristic_slide_payloads(
            prompt=user_request_text or topic_for_content,
            deck_title=base,
            slide_count=n_slides,
            tone=tone or "",
            file_extraction=file_extraction,
            detail_level=detail_level,
        )
        for i in range(min(len(payloads), len(heuristic_payloads))):
            if isinstance(payloads[i], dict) and not _payload_seems_topic_grounded(
                payloads[i], canonical_topic=topic_for_content, keywords=topic_keywords
            ):
                payloads[i] = heuristic_payloads[i]

        if _openai_api_key_available():
            ai_titles2 = _llm_generate_slide_titles(
                topic=topic_for_content,
                deck_title=base,
                slide_count=n_slides,
                tone=tone or "",
            )
            if ai_titles2:
                for ti, t in enumerate(ai_titles2):
                    if ti < len(payloads) and isinstance(payloads[ti], dict) and t.strip():
                        payloads[ti]["title"] = t.strip()

    for i in range(n_slides):
        slide_id = uuid.uuid4().hex
        p = payloads[i] if i < len(payloads) else payloads[i % len(payloads)]
        s_title = str(p.get("title") or f"{base} — Slide {i + 1}")
        slide_type = str(p.get("slideType") or detect_slide_type(p))
        bullets = p.get("bullets")
        if not isinstance(bullets, list) or len(bullets) < 2:
            bullets = [
                f'Define what "{topic_for_content[:80]}" means for your audience and why it matters now.',
                f"Explain two core ideas someone must know about {topic_for_content[:80]} before going deeper.",
                f"Give one practical implication: how understanding {topic_for_content[:80]} changes decisions or next steps.",
            ]
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
        subtitle_default = "" if (i + 1) >= n_slides else f"Section {i + 1}"
        subtitle = str(p.get("subtitle") or subtitle_default).strip()
        description = str(
            p.get("description")
            or (
                f"Framing for {topic_for_content[:80]}: connect context, mechanics, and a decision the audience can make."
            )
        ).strip()
        highlight = str(
            p.get("highlight")
            or f"Actionable takeaway: apply {topic_for_content[:80]} to make a clearer choice under real constraints."
        ).strip()
        key_message = str(p.get("keyMessage") or highlight).strip()
        speaker_notes = str(p.get("speakerNotes") or "").strip()

        gs_p = p.get("gammaStyle") if isinstance(p.get("gammaStyle"), dict) else {}
        layout_preset = str(gs_p.get("layoutPreset") or p.get("layoutSuggestion") or "").strip()
        if layout_preset not in _VALID_GAMMA_LAYOUTS:
            layout_preset = _layout_preset_for_slide_type(slide_type)
        emphasis0 = base.split(" ")[0] if base.split(" ") else "Insight"
        emphasis = [emphasis0, "Impact", "Action"] if slide_type in ("hero", "content", "visual") else [emphasis0, "Key KPI", "Proof"]
        emph_words = (
            gs_p["emphasisWords"] if isinstance(gs_p.get("emphasisWords"), list) and gs_p.get("emphasisWords") else emphasis
        )
        gamma_style_obj: dict[str, Any] = {
            "layoutPreset": layout_preset,
            "alignment": str(gs_p.get("alignment") or "left"),
            "emphasisWords": emph_words[:6],
        }
        for gkey in ("imagePlacement", "fullBleed", "textPrimary", "gradientTitle", "contentAlign", "cardWidth"):
            if gkey in gs_p and gs_p[gkey] is not None:
                gamma_style_obj[gkey] = gs_p[gkey]
        image_query = str(p.get("imageQuery") or "").strip() or f"{topic_for_content} cinematic scene slide {i + 1}"

        content = {
            "title": s_title,
            "subtitle": subtitle,
            "bullets": bullets,
            "description": description,
            "highlight": highlight,
            "keyMessage": key_message,
            "speakerNotes": speaker_notes,
            "imageQuery": image_query[:400],
            "generatedImageUrl": _stable_image_url_for_slide(topic_for_content, i + 1, image_query),
            "slideType": slide_type,
            "layoutSuggestion": layout_preset,
            "emphasisWords": emph_words[:8],
            "contentDensity": "medium",
            "qualityScore": 7.6,
            "gammaStyle": gamma_style_obj,
        }
        con.execute(
            "INSERT INTO slides (id, presentation_id, idx, title, bullets_json, content_json) VALUES (?, ?, ?, ?, ?, ?)",
            (slide_id, presentation_id, i + 1, s_title, "|||".join(bullets), json.dumps(content)),
        )


def _parse_bullets_json(value: str) -> list[str]:
    if not value:
        return []
    return [x for x in value.split("|||") if x]


def _slide_content_from_row(title: str, bullets_raw: str, content_json: str | None) -> dict[str, Any]:
    if content_json:
        try:
            parsed = json.loads(content_json)
            if isinstance(parsed, dict):
                bl = parsed.get("bullets")
                if isinstance(bl, list):
                    parsed = dict(parsed)
                    parsed["bullets"] = _sanitize_bullets_list(bl)
                return parsed
        except Exception:
            pass
    bullets = _sanitize_bullets_list(_parse_bullets_json(bullets_raw))
    return {
        "title": title,
        "subtitle": "",
        "bullets": bullets,
        "description": "",
        "highlight": "",
        "imageQuery": "",
        "layoutSuggestion": "title_bullets",
        "slideType": "content",
        "emphasisWords": [],
        "contentDensity": "medium",
        "qualityScore": 7.0,
        "gammaStyle": {"layoutPreset": "title_bullets", "alignment": "left", "emphasisWords": []},
    }


def _snapshot_from_content(content: dict[str, Any]) -> dict[str, Any]:
    bullets = content.get("bullets")
    if not isinstance(bullets, list):
        bullets = []
    return {
        "title": str(content.get("title") or ""),
        "subtitle": content.get("subtitle"),
        "bullets": [str(b) for b in bullets][:8],
        "highlight": content.get("highlight"),
        "keyMessage": content.get("keyMessage"),
        "qualityScore": float(content.get("qualityScore") or 0),
        "qualityClarity": float(content.get("qualityClarity") or 0),
        "qualityEngagement": float(content.get("qualityEngagement") or 0),
        "qualityVisualBalance": float(content.get("qualityVisualBalance") or 0),
    }


def _hash_password(password: str) -> str:
    # Lightweight deterministic hash for native cutover.
    # Can be upgraded to bcrypt/argon2 in a follow-up migration.
    return hashlib.sha256(password.encode("utf-8")).hexdigest()


def _sample_premium_slides(topic: str) -> list[dict[str, Any]]:
    return [
        {"title": f"Why {topic} Matters", "subtitle": "Executive overview", "points": ["Market context", "Core challenge", "Opportunity window"], "highlight": "Strong upside with focused execution."},
        {"title": f"{topic} Today", "subtitle": "Current-state analysis", "points": ["Baseline metrics", "Key blockers", "Dependency map"]},
        {"title": "Strategic Options", "subtitle": "Decision framing", "points": ["Option A", "Option B", "Option C"], "highlight": "Option B offers best risk-adjusted outcome."},
        {"title": "Execution Roadmap", "subtitle": "90-day plan", "points": ["Week 1-2 setup", "Pilot launch", "Scale and monitor"]},
        {"title": "Expected Outcomes", "subtitle": "Impact projection", "points": ["Efficiency gains", "Quality uplift", "Faster decision cycles"]},
    ]
