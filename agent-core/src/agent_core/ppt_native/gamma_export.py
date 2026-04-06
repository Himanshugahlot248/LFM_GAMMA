"""
Gamma-style PPTX export aligned with next-frontend `templates.ts` + backend `gamma-deck-theme.ts`
and `gamma-export-preset.ts`. Uses python-pptx shapes (background, card, typography) so exports
match deck preview theming instead of the default Office white template.
"""

from __future__ import annotations

import base64
import hashlib
import io
import math
import re
import urllib.request
from typing import Any, TypedDict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Palettes — keep keys in sync with `next-frontend/src/lib/templates.ts`
# ---------------------------------------------------------------------------

GAMMA_DEFAULT_KEY = "gammaDefault"

GAMMA_DECK_PREVIEW = {
    "pageBg": "05070C",
    "cardBg": "0C1118",
    "bodyText": "CBD5E1",
    "titleText": "FFFFFF",
    "keyMuted": "D4D4D8",
    "bulletAccent": "FB923C",
    "visualColumnBg": "0A0E16",
}

EXPORT_FONTS = {"heading": "Georgia", "body": "Calibri"}

# EMU per inch for OOXML / python-pptx (must match slide coordinate system).
_EMU_PER_INCH = 914400


class _TemplatePalette(TypedDict, total=False):
    background: str
    card_fill: str
    title: str
    body: str


TEMPLATE_PALETTES: dict[str, _TemplatePalette] = {
    "gammaDefault": {
        "background": "#05070c",
        "card_fill": "#0c1118",
        "title": "#FFFFFF",
        "body": "#cbd5e1",
    },
    "clementa": {
        "background": "#CBB89D",
        "card_fill": "#E8DDC8",
        "title": "#5A4634",
        "body": "#7A6A55",
    },
    "stratos": {
        "background": "#0B0F2A",
        "card_fill": "#11183C",
        "title": "#FFFFFF",
        "body": "#AAB0FF",
    },
    "nova": {
        "background": "#EAF0FF",
        "card_fill": "#FFFFFF",
        "title": "#4A5D73",
        "body": "#6B7C93",
    },
    "twilight": {
        "background": "#EAD9D2",
        "card_fill": "#F4EDE7",
        "title": "#5C4B44",
        "body": "#7A6A64",
    },
    "coralGlow": {
        "background": "#F7C6C7",
        "card_fill": "#FFFFFF",
        "title": "#A14D4E",
        "body": "#6E3C3D",
    },
    "mercury": {
        "background": "#DDE3EA",
        "card_fill": "#F5F7FA",
        "title": "#4B5563",
        "body": "#6B7280",
    },
    "ashrose": {
        "background": "#E5E5E5",
        "card_fill": "#F2F2F2",
        "title": "#6B6B6B",
        "body": "#8A8A8A",
    },
    "spectrum": {
        "background": "#DDEBFF",
        "card_fill": "#FFFFFF",
        "title": "#3B82F6",
        "body": "#6366F1",
    },
    "stardust": {
        "background": "#000000",
        "card_fill": "#0A0A0A",
        "title": "#FFFFFF",
        "body": "#FFA500",
    },
    "seafoam": {
        "background": "#CFE8E2",
        "card_fill": "#F1FAF7",
        "title": "#2C6E65",
        "body": "#3B8276",
    },
}


def _strip_hex(h: str) -> str:
    x = h.replace("#", "").strip().upper()
    if len(x) == 3:
        x = "".join(c * 2 for c in x)
    return x.zfill(6)[:6]


def _hex_to_rgb_tuple(h: str) -> tuple[int, int, int]:
    s = _strip_hex(h)
    return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)


def _rgb(h: str) -> RGBColor:
    r, g, b = _hex_to_rgb_tuple(h)
    return RGBColor(r, g, b)


def _blend_hex(a: str, b: str, weight_b: float) -> str:
    ra, ga, ba = _hex_to_rgb_tuple(a)
    rb, gb, bb = _hex_to_rgb_tuple(b)
    w = max(0.0, min(1.0, weight_b))
    u = 1.0 - w

    def c(x: int, y: int) -> int:
        return max(0, min(255, int(u * x + w * y)))

    return f"{c(ra, rb):02X}{c(ga, gb):02X}{c(ba, bb):02X}"


def _derive_visual_column_hex(page_bg: str, card_fill: str, accent: str) -> str:
    mixed = _blend_hex(page_bg, accent, 0.22)
    return _blend_hex(mixed, "000000", 0.12)


def _accent_from_palette(p: _TemplatePalette | None) -> str:
    if p and p.get("title"):
        return _strip_hex(p["title"])
    return "1B6EF3"


class ExportTheme(TypedDict):
    pageBg: str
    cardFillHex: str
    titleHex: str
    bodyHex: str
    accentHex: str
    bulletHex: str
    visualColHex: str
    fonts: dict[str, str]


def resolve_export_theme(template_name: str | None) -> ExportTheme:
    """Mirror `backend/src/lib/gamma-deck-theme.ts` `resolveGammaExportTheme`."""
    name = (template_name or "").strip()
    use_gamma = not name or name == GAMMA_DEFAULT_KEY
    fonts = dict(EXPORT_FONTS)

    if use_gamma:
        return {
            "pageBg": GAMMA_DECK_PREVIEW["pageBg"],
            "cardFillHex": GAMMA_DECK_PREVIEW["cardBg"],
            "titleHex": GAMMA_DECK_PREVIEW["titleText"],
            "bodyHex": GAMMA_DECK_PREVIEW["bodyText"],
            "accentHex": _accent_from_palette(TEMPLATE_PALETTES.get("gammaDefault")),
            "bulletHex": GAMMA_DECK_PREVIEW["bulletAccent"],
            "visualColHex": GAMMA_DECK_PREVIEW["visualColumnBg"],
            "fonts": fonts,
        }

    pal = TEMPLATE_PALETTES.get(name)
    page_bg = _strip_hex(pal["background"]) if pal and pal.get("background") else GAMMA_DECK_PREVIEW["pageBg"]
    card_fill = _strip_hex(pal["card_fill"]) if pal and pal.get("card_fill") else page_bg
    title_hex = _strip_hex(pal["title"]) if pal and pal.get("title") else GAMMA_DECK_PREVIEW["titleText"]
    body_hex = _strip_hex(pal["body"]) if pal and pal.get("body") else GAMMA_DECK_PREVIEW["bodyText"]
    accent_hex = _accent_from_palette(pal)
    bullet_hex = accent_hex
    visual_col = _derive_visual_column_hex(page_bg, card_fill, accent_hex)

    return {
        "pageBg": page_bg,
        "cardFillHex": card_fill,
        "titleHex": title_hex,
        "bodyHex": body_hex,
        "accentHex": accent_hex,
        "bulletHex": bullet_hex,
        "visualColHex": visual_col,
        "fonts": fonts,
    }


def resolve_gamma_preset_for_export(content: dict[str, Any]) -> str:
    """Mirror `backend/src/lib/gamma-export-preset.ts` `resolveGammaPresetForExport`."""
    gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
    preset = gs.get("layoutPreset") if isinstance(gs, dict) else None
    if preset in ("hero_split", "two_column", "three_cards", "stats_split", "title_bullets", "section_only"):
        return str(preset)

    bullets = content.get("bullets")
    blist = bullets if isinstance(bullets, list) else []
    st = str(content.get("slideType") or "").lower()

    if st in ("hero", "visual"):
        return "hero_split"
    if st in ("stats", "stat"):
        return "stats_split"
    if st in ("comparison", "timeline", "split"):
        return "two_column"
    if st == "section":
        return "section_only" if len(blist) == 0 else "title_bullets"
    if st == "content":
        return "title_bullets"

    lt = str(content.get("layoutType") or content.get("layoutSuggestion") or "title_bullets")
    if lt == "two_column":
        return "hero_split"
    if lt == "summary":
        return "stats_split"
    if lt == "section_break" and len(blist) == 0:
        return "section_only"

    return "title_bullets"


def _card_fill_from_content(content: dict[str, Any], default_hex: str) -> str:
    gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
    cc = gs.get("cardColor") if isinstance(gs, dict) else None
    if isinstance(cc, str) and len(cc.replace("#", "").strip()) >= 3:
        return _strip_hex(cc)
    return _strip_hex(default_hex)


def _gamma_typography_from_content(content: dict[str, Any], theme: ExportTheme) -> tuple[str, str, float, bool, float | None]:
    """Return (heading_font, body_font, body_pt, body_bold, legacy_line_height_pt_or_none).

    Keep this return shape stable because many render functions unpack it.
    New PowerPoint-style spacing/alignment controls are handled by
    `_gamma_paragraph_controls_from_content`.
    """
    gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
    ff = str(gs.get("fontFamily") or "").strip() if isinstance(gs, dict) else ""
    body_font = ff or theme["fonts"]["body"]
    heading_font = ff or theme["fonts"]["heading"]
    fs = gs.get("fontSize") if isinstance(gs, dict) else None
    body_pt = float(fs) if isinstance(fs, (int, float)) and fs > 0 else 16.0
    fw = gs.get("fontWeight") if isinstance(gs, dict) else None
    body_bold = bool(isinstance(fw, (int, float)) and fw >= 600)
    lh_pt = gs.get("lineHeightPt") if isinstance(gs, dict) else None
    lh_px = gs.get("lineHeightPx") if isinstance(gs, dict) else None
    legacy_lh_pt = float(lh_pt) if isinstance(lh_pt, (int, float)) and lh_pt > 0 else None
    if legacy_lh_pt is None and isinstance(lh_px, (int, float)) and lh_px > 0:
        legacy_lh_pt = float(lh_px) / 1.333
    return heading_font, body_font, body_pt, body_bold, legacy_lh_pt


def _gamma_paragraph_controls_from_content(content: dict[str, Any]) -> tuple[float, float, float, str]:
    """Return (line_spacing_mult, para_before_pt, para_after_pt, align)."""
    gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
    ls = gs.get("lineSpacing") if isinstance(gs, dict) else None
    line_spacing_mult = float(ls) if isinstance(ls, (int, float)) and ls > 0 else 1.25
    pb = gs.get("paraSpaceBeforePt") if isinstance(gs, dict) else None
    pa = gs.get("paraSpaceAfterPt") if isinstance(gs, dict) else None
    para_before_pt = float(pb) if isinstance(pb, (int, float)) and pb >= 0 else 0.0
    para_after_pt = float(pa) if isinstance(pa, (int, float)) and pa >= 0 else 8.0
    ta = str(gs.get("textAlign") or "").strip().lower() if isinstance(gs, dict) else ""
    align = ta if ta in ("left", "center", "right") else "left"

    # Legacy mapping: old “lineHeight” becomes approximate multiplier.
    lh_pt = gs.get("lineHeightPt") if isinstance(gs, dict) else None
    lh_px = gs.get("lineHeightPx") if isinstance(gs, dict) else None
    fs = gs.get("fontSize") if isinstance(gs, dict) else None
    body_pt = float(fs) if isinstance(fs, (int, float)) and fs > 0 else 16.0
    if isinstance(lh_pt, (int, float)) and lh_pt > 0:
        line_spacing_mult = max(1.0, min(2.2, float(lh_pt) / body_pt))
    elif isinstance(lh_px, (int, float)) and lh_px > 0:
        line_spacing_mult = max(1.0, min(2.2, (float(lh_px) / 1.333) / body_pt))

    return line_spacing_mult, para_before_pt, para_after_pt, align


def _bullet_marker_text_from_content(content: dict[str, Any]) -> str:
    gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
    bm = str(gs.get("bulletMarker") or "").strip().lower() if isinstance(gs, dict) else ""
    if bm == "square":
        return "■"
    if bm == "check":
        return "✓"
    if bm == "arrow":
        return "→"
    return "●"


def _add_full_bleed_rect(slide: Any, prs: Presentation, hex6: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex6)
    sp = shape.line
    sp.fill.background()


def _add_rounded_card(
    slide: Any,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill_hex: str,
    line_alpha: float = 0.88,
    *,
    show_outline: bool = True,
) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    line = shape.line
    if show_outline:
        line.color.rgb = RGBColor(255, 255, 255)
        line.transparency = line_alpha
    else:
        line.fill.background()
    return shape


def _bullet_row_from_content(content: dict[str, Any], index: int) -> list[dict[str, Any]] | None:
    raw = content.get("bulletRuns")
    if not isinstance(raw, list) or index >= len(raw):
        return None
    row = raw[index]
    if not isinstance(row, list):
        return None
    out: list[dict[str, Any]] = []
    for r in row:
        if isinstance(r, dict) and str(r.get("text", "")).strip() != "":
            out.append(
                {
                    "text": str(r.get("text") or ""),
                    "bold": bool(r.get("bold")),
                    "italic": bool(r.get("italic")),
                    "underline": bool(r.get("underline")),
                    "strike": bool(r.get("strike")),
                }
            )
    return out or None


def _runs_list_from_content(content: dict[str, Any], key: str) -> list[dict[str, Any]] | None:
    raw = content.get(key)
    if not isinstance(raw, list) or not raw:
        return None
    out: list[dict[str, Any]] = []
    for r in raw:
        if isinstance(r, dict) and str(r.get("text", "")).strip() != "":
            out.append(
                {
                    "text": str(r.get("text") or ""),
                    "bold": bool(r.get("bold")),
                    "italic": bool(r.get("italic")),
                    "underline": bool(r.get("underline")),
                    "strike": bool(r.get("strike")),
                }
            )
    return out or None


def _fill_paragraph_runs(
    paragraph: Any,
    runs: list[dict[str, Any]],
    *,
    font_name: str,
    size_pt: float,
    default_rgb: RGBColor,
) -> None:
    paragraph.clear()
    for r in runs:
        text = str(r.get("text") or "")
        if not text:
            continue
        run = paragraph.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.color.rgb = default_rgb
        run.font.bold = bool(r.get("bold"))
        run.font.italic = bool(r.get("italic"))
        if r.get("underline"):
            run.font.underline = True
        if r.get("strike"):
            run.font.strike = True


def _fill_bullet_paragraph(
    para: Any,
    plain_line: str,
    rich: list[dict[str, Any]] | None,
    *,
    marker_rgb: RGBColor,
    body_rgb: RGBColor,
    font_name: str,
    body_pt: float,
    marker_pt: float,
    marker_text: str = "●",
    default_bold: bool = False,
) -> None:
    para.clear()
    para.alignment = PP_ALIGN.LEFT
    m = para.add_run()
    m.text = f"{marker_text}  "
    m.font.name = font_name
    m.font.size = Pt(marker_pt)
    m.font.color.rgb = marker_rgb
    m.font.bold = bool(default_bold)
    if rich:
        for r in rich:
            text = str(r.get("text") or "")
            if not text:
                continue
            run = para.add_run()
            run.text = text
            run.font.name = font_name
            run.font.size = Pt(body_pt)
            run.font.color.rgb = body_rgb
            run.font.bold = bool(r.get("bold")) if r.get("bold") is not None else bool(default_bold)
            run.font.italic = bool(r.get("italic"))
            if r.get("underline"):
                run.font.underline = True
            if r.get("strike"):
                run.font.strike = True
    else:
        run = para.add_run()
        run.text = plain_line
        run.font.name = font_name
        run.font.size = Pt(body_pt)
        run.font.color.rgb = body_rgb
        run.font.bold = bool(default_bold)


def _ensure_text_frame_wrap(tf: Any) -> None:
    """Prevent clipped single-line overflow in PowerPoint (common export issue)."""
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.margin_left = Inches(0.03)
        tf.margin_right = Inches(0.03)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
    except Exception:
        pass


def _fetch_image_bytes(url: str, timeout: float = 20.0) -> bytes | None:
    if not isinstance(url, str) or not url.startswith("http"):
        return None
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (compatible; agent-core/1.0)"})
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            return resp.read()
    except Exception:
        return None


def _picsum_image_url(seed_key: str, *, width: int = 1280, height: int = 720) -> str:
    """Deterministic Picsum fallback for export.

    The browser preview can sometimes show an image even when export can't
    fetch it (different fetch environment). This ensures exported PPTs still
    contain a picture frame.
    """
    mat = (seed_key or "").strip()
    if not mat:
        mat = "lf-ai"
    h = hashlib.md5(mat.encode("utf-8")).hexdigest()[:10]
    return f"https://picsum.photos/seed/{h}/{width}/{height}"


def _safe_line_spacing_pt(*, font_pt: float, line_height_pt: float | None) -> float | None:
    """Clamp line spacing so it never becomes smaller than the font size.

    In PPT, if exact line spacing is smaller than the font size, wrapped lines can overlap.
    We convert px→pt (1px ≈ 0.75pt) and clamp to ~1.15× font size for readability.
    """
    if not isinstance(line_height_pt, (int, float)) or line_height_pt <= 0:
        return None
    req_pt = float(line_height_pt)
    min_pt = float(font_pt) * 1.15
    return max(req_pt, min_pt)


def _image_bytes_from_query(url_or_data: str) -> bytes | None:
    """HTTP(S) URL or ``data:image/...;base64,...`` (whitespace in base64 stripped)."""
    t = (url_or_data or "").strip()
    if not t:
        return None
    m = re.match(r"^data:image/(png|jpeg|webp|jpg);base64,([\s\S]+)$", t, re.IGNORECASE)
    if m:
        b64 = re.sub(r"\s+", "", m.group(2))
        try:
            return base64.b64decode(b64, validate=False)
        except Exception:
            return None
    low = t.lower()
    if low.startswith("http://") or low.startswith("https://"):
        return _fetch_image_bytes(t)
    return None


def _chart_snapshot_url(content: dict[str, Any]) -> str:
    u = content.get("chartSnapshotUrl")
    return str(u).strip() if isinstance(u, str) else ""


def _parse_chart_placement_pct(raw: Any) -> float | None:
    """Accept int/float or numeric strings (JSON may deserialize numbers as strings)."""
    try:
        if isinstance(raw, (int, float)) and math.isfinite(float(raw)):
            return float(raw)
        if isinstance(raw, str) and raw.strip():
            v = float(raw.strip())
            if math.isfinite(v):
                return v
    except Exception:
        return None
    return None


def _has_explicit_chart_placement(content: dict[str, Any]) -> bool:
    cp = content.get("chartPlacement")
    if not isinstance(cp, dict):
        return False
    for k in ("xPct", "yPct", "wPct", "hPct"):
        if _parse_chart_placement_pct(cp.get(k)) is None:
            return False
    return True


def _chart_placement_emu_box(prs: Presentation, content: dict[str, Any]) -> tuple[int, int, int, int] | None:
    if not _has_explicit_chart_placement(content):
        return None
    cp = content.get("chartPlacement")
    if not isinstance(cp, dict):
        return None
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    if sw <= 0 or sh <= 0:
        return None

    def pct_to_emu(axis: int, raw: Any) -> int:
        pv = _parse_chart_placement_pct(raw)
        if pv is None:
            return 0
        v = max(0.0, min(100.0, pv))
        return int(round(axis * v / 100.0))

    x = pct_to_emu(sw, cp.get("xPct"))
    y = pct_to_emu(sh, cp.get("yPct"))
    w = pct_to_emu(sw, cp.get("wPct"))
    h = pct_to_emu(sh, cp.get("hPct"))
    min_w = int(Inches(0.55))
    min_h = int(Inches(0.45))
    w = max(min_w, w)
    h = max(min_h, h)
    if x + w > sw:
        x = max(0, sw - w)
    if y + h > sh:
        y = max(0, sh - h)
    return (x, y, w, h)


def _allow_picture_stretch_to_frame(picture: Any) -> None:
    """Disable OOXML ``noChangeAspect`` on a picture.

    python-pptx sets ``a:picLocks/@noChangeAspect=1`` by default. Viewers then
    preserve aspect ratio using the bitmap's DPI-based native size; tall
    portrait images scaled to panel width overflow the frame vertically.
    LibreOffice PDF export is especially prone to this. Turning the flag off
    lets ``a:stretch`` fill ``p:spPr/a:xfrm`` exactly (we pre-rasterize to
    match aspect ratio so there is no visible stretch).
    """
    try:
        root = picture._element
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for pl in root.findall(f".//{{{ns}}}picLocks"):
            if pl.get("noChangeAspect") in ("1", "true", "True"):
                pl.set("noChangeAspect", "0")
    except Exception:
        pass


def _raster_to_picture_frame(
    raw: bytes,
    cx_emu: int,
    cy_emu: int,
    *,
    fit: str,
    pad_hex: str | None,
    jpeg_quality: int = 90,
    dpi: int = 150,
    encode: str = "jpeg",
) -> bytes | None:
    """Rasterize to ``cx_emu``×``cy_emu`` (pixel size from dpi). ``encode`` is ``jpeg`` or ``png``.

    ``fit`` is ``\"cover\"`` (center-crop then resize) or ``\"contain\"`` (letterbox on ``pad_hex``).
    Chart snapshots use PNG to avoid lossy JPEG artifacts vs deck preview.
    """
    if cx_emu <= 0 or cy_emu <= 0:
        return None
    try:
        from PIL import Image, ImageOps
    except ImportError:
        return None
    try:
        im = Image.open(io.BytesIO(raw))
        im = ImageOps.exif_transpose(im)
        target_ar = (cx_emu / cy_emu) if cy_emu else 1.0
        w_px = max(1, int(round(cx_emu * dpi / _EMU_PER_INCH)))
        h_px = max(1, int(round(cy_emu * dpi / _EMU_PER_INCH)))
        cap = 2400
        w_px = min(w_px, cap)
        h_px = min(h_px, cap)

        if fit == "cover":
            if im.mode != "RGB":
                im = im.convert("RGB")
            iw, ih = im.size
            if iw <= 1 or ih <= 1:
                return None
            src_ar = iw / ih
            if src_ar > target_ar:
                new_w = max(1, int(round(ih * target_ar)))
                x0 = max(0, (iw - new_w) // 2)
                im = im.crop((x0, 0, x0 + new_w, ih))
            else:
                new_h = max(1, int(round(iw / target_ar)))
                y0 = max(0, (ih - new_h) // 2)
                im = im.crop((0, y0, iw, y0 + new_h))
            im = im.resize((w_px, h_px), Image.Resampling.LANCZOS)
        else:
            pad = pad_hex or "0C1118"
            if im.mode != "RGBA":
                im = im.convert("RGBA")
            iw, ih = im.size
            if iw <= 1 or ih <= 1:
                return None
            scale = min(w_px / iw, h_px / ih)
            nw = max(1, int(iw * scale))
            nh = max(1, int(ih * scale))
            rim = im.resize((nw, nh), Image.Resampling.LANCZOS)
            br, bg, bb = _hex_to_rgb_tuple(pad)
            canvas = Image.new("RGBA", (w_px, h_px), (br, bg, bb, 255))
            ox = (w_px - nw) // 2
            oy = (h_px - nh) // 2
            canvas.paste(rim, (ox, oy), rim.split()[3])
            im = canvas.convert("RGB")

        buf = io.BytesIO()
        enc = (encode or "jpeg").strip().lower()
        if enc == "png":
            im.save(buf, format="PNG", optimize=True)
        else:
            im.save(buf, format="JPEG", quality=jpeg_quality, optimize=True, dpi=(dpi, dpi))
        return buf.getvalue()
    except Exception:
        return None


def _add_chart_snapshot_contain(
    slide: Any,
    content: dict[str, Any],
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    *,
    pad_hex: str,
) -> bool:
    snap = _chart_snapshot_url(content)
    if not snap:
        return False
    raw = _image_bytes_from_query(snap)
    if not raw:
        return False
    blob = _raster_to_picture_frame(
        raw,
        int(width),
        int(height),
        fit="contain",
        pad_hex=pad_hex,
        jpeg_quality=90,
        encode="png",
    )
    if not blob:
        return False
    try:
        pic = slide.shapes.add_picture(io.BytesIO(blob), left, top, width=width, height=height)
        _allow_picture_stretch_to_frame(pic)
        return True
    except Exception:
        return False


def _add_picture_from_url(
    slide: Any,
    url: str,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
) -> bool:
    raw = _image_bytes_from_query(url)
    if not raw:
        return False
    blob = _raster_to_picture_frame(
        raw,
        int(width),
        int(height),
        fit="cover",
        pad_hex=None,
        jpeg_quality=88,
    )
    if not blob:
        return False
    try:
        pic = slide.shapes.add_picture(io.BytesIO(blob), left, top, width=width, height=height)
        _allow_picture_stretch_to_frame(pic)
        return True
    except Exception:
        return False


def _set_textbox_style(
    box: Any,
    *,
    font_name: str,
    size_pt: float,
    color_hex: str,
    bold: bool = False,
    align: int | None = None,
    valign: int | None = None,
    line_height_px: float | None = None,
    line_height_pt: float | None = None,
    line_spacing_mult: float | None = None,
    space_before_pt: float | None = None,
    space_after_pt: float | None = None,
) -> None:
    tf = box.text_frame
    _ensure_text_frame_wrap(tf)
    if valign is not None:
        tf.vertical_anchor = valign
    for p in tf.paragraphs:
        p.font.name = font_name
        p.font.size = Pt(size_pt)
        p.font.bold = bold
        p.font.color.rgb = _rgb(color_hex)
        if align is not None:
            p.alignment = align
        if isinstance(line_spacing_mult, (int, float)) and line_spacing_mult > 0:
            # python-pptx supports float multiplier for line spacing
            p.line_spacing = float(line_spacing_mult)
        else:
            lh_pt = line_height_pt
            if lh_pt is None and isinstance(line_height_px, (int, float)) and line_height_px > 0:
                lh_pt = float(line_height_px) / 1.333
            sp = _safe_line_spacing_pt(font_pt=float(size_pt), line_height_pt=lh_pt)
            if sp is not None:
                p.line_spacing = Pt(sp)
        if isinstance(space_before_pt, (int, float)) and space_before_pt >= 0:
            p.space_before = Pt(float(space_before_pt))
        if isinstance(space_after_pt, (int, float)) and space_after_pt >= 0:
            p.space_after = Pt(float(space_after_pt))


def _render_section_only(
    slide: Any,
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    key_message: str,
    theme: ExportTheme,
    content: dict[str, Any],
) -> None:
    sw, sh = prs.slide_width, prs.slide_height
    margin = Inches(0.5)
    _add_rounded_card(slide, margin, margin, sw - 2 * margin, sh - 2 * margin, theme["cardFillHex"])
    heading_font, body_font, body_pt_base, body_bold, lh_px = _gamma_typography_from_content(content, theme)
    line_mult, para_before_pt, para_after_pt, align_txt = _gamma_paragraph_controls_from_content(content)
    line_mult, para_before_pt, para_after_pt, align_txt = _gamma_paragraph_controls_from_content(content)
    line_mult, para_before_pt, para_after_pt, align_txt = _gamma_paragraph_controls_from_content(content)

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), sw - Inches(2.0), Inches(2.0))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = heading_font
    p.font.size = Pt(max(22, int(round(body_pt_base + 20))))
    p.font.bold = True
    p.font.color.rgb = _rgb(theme["titleHex"])
    p.line_spacing = float(line_mult)
    p.space_before = Pt(float(para_before_pt))
    p.space_after = Pt(float(max(0.0, para_after_pt * 0.5)))
    p.alignment = PP_ALIGN.CENTER

    sub = subtitle.strip()
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.name = body_font
        p2.font.size = Pt(max(10, int(round(body_pt_base - 2))))
        p2.font.color.rgb = _rgb(theme["bodyHex"])
        p2.alignment = PP_ALIGN.CENTER
        p2.line_spacing = float(line_mult)
        p2.space_before = Pt(float(para_before_pt))
        p2.space_after = Pt(float(max(0.0, para_after_pt * 0.5)))

    if key_message.strip():
        km = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), sw - Inches(2.4), Inches(1.2))
        km.text_frame.text = key_message.strip()
        _set_textbox_style(
            km,
            font_name=body_font,
            size_pt=max(10, int(round(body_pt_base - 2))),
            color_hex=theme["bodyHex"],
            bold=body_bold,
            align=PP_ALIGN.CENTER,
            line_spacing_mult=line_mult,
            space_before_pt=para_before_pt,
            space_after_pt=para_after_pt,
        )

    if not _has_explicit_chart_placement(content):
        _add_chart_snapshot_contain(
            slide,
            content,
            Inches(1.2),
            Inches(2.0),
            sw - Inches(2.4),
            Inches(3.2),
            pad_hex=theme["cardFillHex"],
        )


def _render_hero_split(
    slide: Any,
    prs: Presentation,
    *,
    title: str,
    bullets: list[str],
    subtitle: str,
    key_message: str,
    content: dict[str, Any],
    theme: ExportTheme,
) -> None:
    sw, sh = prs.slide_width, prs.slide_height

    margin = Inches(0.45)
    gutter = Inches(0.25)
    visual_w = Inches(5.15)
    top_y = Inches(0.5)
    panel_h = sh - Inches(1.0)

    gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
    lh_px = None
    try:
        v = gs.get("lineHeightPx") if isinstance(gs, dict) else None
        if isinstance(v, (int, float)) and v > 0:
            lh_px = float(v)
    except Exception:
        lh_px = None
    image_on_right = gs.get("imagePlacement") == "right"
    if image_on_right:
        visual_x = sw - margin - visual_w
        text_x = margin
        text_w = sw - margin - gutter - visual_w - margin
    else:
        visual_x = margin
        text_x = margin + visual_w + gutter
        text_w = sw - (margin + visual_w + gutter) - margin

    vch = theme["visualColHex"]
    cf = _card_fill_from_content(content, theme["cardFillHex"])
    ah = theme["accentHex"]
    bu = theme["bulletHex"]
    th = theme["titleHex"]
    bh = theme["bodyHex"]
    heading_font, body_font, body_pt_base, body_bold, lh_px = _gamma_typography_from_content(content, theme)
    line_mult, para_before_pt, para_after_pt, align_txt = _gamma_paragraph_controls_from_content(content)
    marker_text = _bullet_marker_text_from_content(content)

    # Visual column (no stroke — thin light outlines often show as stray vertical lines in PDF export).
    vc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, visual_x, top_y, visual_w, panel_h)
    vc.fill.solid()
    vc.fill.fore_color.rgb = _rgb(vch)
    vc.line.fill.background()

    img_url = content.get("generatedImageUrl")
    pic_ok = False
    if isinstance(img_url, str) and img_url.strip():
        pic_ok = _add_picture_from_url(
            slide,
            img_url,
            visual_x + Inches(0.1),
            top_y + Inches(0.1),
            visual_w - Inches(0.2),
            panel_h - Inches(0.2),
        )
    if not pic_ok:
        # Export may fail to fetch `generatedImageUrl` depending on runtime/network.
        # Use deterministic Picsum fallback so PPTs still show an image.
        seed_key = ""
        for k in ("imageQuery", "generatedImagePrompt", "title"):
            v = content.get(k)
            if isinstance(v, str) and v.strip():
                seed_key = v.strip()
                break
        fallback = _picsum_image_url(seed_key, width=1280, height=720)
        pic_ok = _add_picture_from_url(
            slide,
            fallback,
            visual_x + Inches(0.1),
            top_y + Inches(0.1),
            visual_w - Inches(0.2),
            panel_h - Inches(0.2),
        )
    if not pic_ok:
        # Soft accent ellipses when no image (matches preview placeholder look)
        e1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, visual_x + Inches(0.35), top_y + Inches(0.45), Inches(2.4), Inches(2.4))
        e1.fill.solid()
        e1.fill.fore_color.rgb = _rgb(ah)
        e1.fill.transparency = 0.78
        e1.line.fill.background()

        e2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, visual_x + Inches(2.4), top_y + Inches(3.0), Inches(2.4), Inches(2.0))
        e2.fill.solid()
        e2.fill.fore_color.rgb = _rgb(bu)
        e2.fill.transparency = 0.72
        e2.line.fill.background()

    _add_rounded_card(slide, text_x, top_y, text_w, panel_h, cf, show_outline=False)

    title_box_w = text_w - Inches(0.7)
    fitted = max(22, min(36, int(round(body_pt_base + 12))))
    # Match deck preview (snug title + mt-2 to bullets): avoid oversized title box that adds a false gap.
    raw_title = (title or "").strip()
    nl = raw_title.count("\n") + 1 if raw_title else 1
    wrap_lines = max(1, (len(raw_title.replace("\n", " ")) + 34) // 35)
    est_title_lines = max(nl, wrap_lines)
    line_h_in = float(fitted) * 1.12 / 72.0
    title_h = Inches(min(2.1, 0.06 + line_h_in * est_title_lines))

    title_runs = _runs_list_from_content(content, "titleRuns")
    title_top = top_y + Inches(0.38)
    tt = slide.shapes.add_textbox(text_x + Inches(0.35), title_top, title_box_w, title_h)
    ttf = tt.text_frame
    ttf.clear()
    _ensure_text_frame_wrap(ttf)
    p0 = ttf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    if title_runs:
        _fill_paragraph_runs(
            p0,
            title_runs,
            font_name=heading_font,
            size_pt=float(fitted),
            default_rgb=_rgb(th),
        )
    else:
        p0.text = title
        p0.font.name = heading_font
        p0.font.size = Pt(fitted)
        p0.font.bold = True
        p0.font.color.rgb = _rgb(th)

    sub = subtitle.strip()
    gap_after_title = Inches(0.16)
    if sub:
        subtitle_y = title_top + title_h + gap_after_title
        bullets_y = subtitle_y + Inches(0.52)
        stb = slide.shapes.add_textbox(text_x + Inches(0.35), subtitle_y, title_box_w, Inches(0.48))
        stf = stb.text_frame
        stf.clear()
        _ensure_text_frame_wrap(stf)
        sp0 = stf.paragraphs[0]
        sp0.text = sub
        sp0.alignment = PP_ALIGN.LEFT
        sp0.font.name = body_font
        sp0.font.size = Pt(max(10, int(round(body_pt_base - 4))))
        sp0.font.color.rgb = _rgb(bh)
        sp0.line_spacing = float(line_mult)
        sp0.space_before = Pt(float(para_before_pt))
        sp0.space_after = Pt(float(max(0.0, para_after_pt * 0.5)))
    else:
        # Preview: `mt-2` below title ≈ 8px @ 96dpi — not gap_after_title + extra 0.14" (was ~0.30" total).
        bullets_y = title_top + title_h + Inches(8.0 / 96.0)

    chart_h = Inches(2.0)
    if not _has_explicit_chart_placement(content):
        if _add_chart_snapshot_contain(
            slide,
            content,
            text_x + Inches(0.4),
            bullets_y,
            text_w - Inches(0.85),
            chart_h,
            pad_hex=cf,
        ):
            bullets_y = bullets_y + chart_h + Inches(0.14)

    km = key_message.strip()
    key_h = Inches(1.05) if km else Inches(0.12)
    gap = Inches(0.12)
    bullets_h = panel_h - (bullets_y - top_y) - key_h - gap
    if bullets_h < Inches(1.1):
        bullets_h = Inches(1.1)

    bt = slide.shapes.add_textbox(text_x + Inches(0.4), bullets_y, text_w - Inches(0.8), bullets_h)
    btf = bt.text_frame
    btf.clear()
    _ensure_text_frame_wrap(btf)
    items = bullets[:6]
    body_pt = max(10, int(round(body_pt_base - (1 if len(items) >= 5 else 0))))
    if items:
        for i, line in enumerate(items):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            br = _bullet_row_from_content(content, i)
            _fill_bullet_paragraph(
                para,
                line,
                br,
                marker_rgb=_rgb(bu),
                body_rgb=_rgb(bh),
                font_name=body_font,
                body_pt=float(body_pt),
                marker_pt=float(max(11, body_pt - 1)),
                marker_text=marker_text,
                default_bold=body_bold,
            )
            para.line_spacing = float(line_mult)
            # First bullet: no extra space_before (spacing comes from title→bullets gap only; matches preview ul).
            para.space_before = Pt(float(0 if i == 0 else para_before_pt))
            para.space_after = Pt(float(para_after_pt))
            para.alignment = PP_ALIGN.LEFT if align_txt == "left" else PP_ALIGN.CENTER if align_txt == "center" else PP_ALIGN.RIGHT

    if km:
        km_top = top_y + panel_h - key_h
        # Text only (no thin accent bar — reads as a stray vertical line in exports).
        kb = slide.shapes.add_textbox(
            text_x + Inches(0.4),
            km_top,
            text_w - Inches(0.85),
            key_h,
        )
        ktf = kb.text_frame
        ktf.clear()
        _ensure_text_frame_wrap(ktf)
        kp = ktf.paragraphs[0]
        kp.text = km
        kp.alignment = PP_ALIGN.LEFT
        kp.font.name = body_font
        kp.font.size = Pt(max(9, int(round(body_pt_base - 5))))
        kp.font.bold = True
        kp.font.color.rgb = _rgb(GAMMA_DECK_PREVIEW["keyMuted"])
        kp.line_spacing = float(line_mult)
        kp.space_before = Pt(float(para_before_pt))
        kp.space_after = Pt(float(max(0.0, para_after_pt * 0.5)))


def _render_title_bullets(
    slide: Any,
    prs: Presentation,
    *,
    title: str,
    bullets: list[str],
    subtitle: str,
    key_message: str,
    content: dict[str, Any],
    theme: ExportTheme,
) -> None:
    sw, sh = prs.slide_width, prs.slide_height
    cf = _card_fill_from_content(content, theme["cardFillHex"])
    th = theme["titleHex"]
    bh = theme["bodyHex"]
    bu = theme["bulletHex"]
    vch = theme["visualColHex"]
    gs = content.get("gammaStyle") if isinstance(content.get("gammaStyle"), dict) else {}
    heading_font, body_font, body_pt_base, body_bold, lh_px = _gamma_typography_from_content(content, theme)
    line_mult, para_before_pt, para_after_pt, align_txt = _gamma_paragraph_controls_from_content(content)
    marker_text = _bullet_marker_text_from_content(content)

    card_left = Inches(0.5)
    card_top = Inches(0.55)
    card_w = sw - Inches(1.0)
    card_h = sh - Inches(1.1)
    card_right = card_left + card_w
    card_bottom = card_top + card_h
    _add_rounded_card(slide, card_left, card_top, card_w, card_h, cf)

    title_box_w = sw - Inches(2.0)
    fitted = max(20, min(34, int(round(body_pt_base + 10))))
    title_h = Inches(min(2.45, 0.42 + 0.42 * max(1, min(5, (len(title) // 44) + 1))))

    title_runs_tb = _runs_list_from_content(content, "titleRuns")
    title_top_y = Inches(0.72)
    tt = slide.shapes.add_textbox(Inches(1.0), title_top_y, title_box_w, title_h)
    ttf = tt.text_frame
    ttf.clear()
    _ensure_text_frame_wrap(ttf)
    p0 = ttf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    if title_runs_tb:
        _fill_paragraph_runs(
            p0,
            title_runs_tb,
            font_name=heading_font,
            size_pt=float(fitted),
            default_rgb=_rgb(th),
        )
    else:
        p0.text = title
        p0.font.name = heading_font
        p0.font.size = Pt(fitted)
        p0.font.bold = True
        p0.font.color.rgb = _rgb(th)

    list_y = title_top_y + title_h + Inches(0.22)
    sub = subtitle.strip()
    if sub:
        sub_y = title_top_y + title_h + Inches(0.14)
        sb = slide.shapes.add_textbox(Inches(1.0), sub_y, title_box_w, Inches(0.48))
        sb.text_frame.text = sub
        _set_textbox_style(
            sb,
            font_name=body_font,
            size_pt=max(10, int(round(body_pt_base - 3))),
            color_hex=bh,
            bold=body_bold,
            align=PP_ALIGN.CENTER,
            line_spacing_mult=line_mult,
            space_before_pt=para_before_pt,
            space_after_pt=max(0.0, para_after_pt * 0.5),
        )
        list_y = sub_y + Inches(0.52)

    # Content (bullets/chart) top — image (if present) should align with this top.
    content_top_y = list_y

    key_h = Inches(1.05) if key_message.strip() else Inches(0.4)
    # Keep content inside the rounded card area with consistent padding.
    bottom_pad = Inches(0.35)
    content_panel_h = (card_bottom - bottom_pad) - content_top_y - key_h
    if content_panel_h < Inches(1.2):
        content_panel_h = Inches(1.2)

    # Export must render any image shown in deck preview, even for title-only layouts.
    img_url = content.get("generatedImageUrl")
    has_image = isinstance(img_url, str) and img_url.strip()
    image_on_right = gs.get("imagePlacement") == "right"

    if has_image:
        gutter = Inches(0.28)
        inner_pad = Inches(0.25)
        left_edge = card_left + inner_pad
        right_edge = card_right - inner_pad
        # Slightly smaller image column to avoid crowded text.
        visual_w = int(card_w * 0.32)
        visual_w = max(Inches(3.5), min(visual_w, Inches(4.9)))
        # Ensure text column never collapses too far.
        min_text_w = Inches(6.0)
        if (right_edge - left_edge - visual_w - gutter) < min_text_w:
            visual_w = max(Inches(3.0), (right_edge - left_edge - gutter - min_text_w))

        if image_on_right:
            visual_x = right_edge - visual_w
            text_x = left_edge
            text_w = visual_x - gutter - text_x
        else:
            visual_x = left_edge
            text_x = visual_x + visual_w + gutter
            text_w = right_edge - text_x

        # Visual column frame.
        vc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, visual_x, content_top_y, visual_w, content_panel_h)
        vc.fill.solid()
        vc.fill.fore_color.rgb = _rgb(vch)
        vc.line.fill.background()

        # Image (try the URL first; fall back to deterministic Picsum for export).
        img_padding = Inches(0.1)
        pic_ok = False
        if isinstance(img_url, str) and img_url.strip():
            pic_ok = _add_picture_from_url(
                slide,
                img_url,
                visual_x + img_padding,
                content_top_y + img_padding,
                visual_w - Inches(0.2),
                content_panel_h - Inches(0.2),
            )
        if not pic_ok:
            seed_key = ""
            for k in ("imageQuery", "generatedImagePrompt", "title"):
                v = content.get(k)
                if isinstance(v, str) and v.strip():
                    seed_key = v.strip()
                    break
            fallback = _picsum_image_url(seed_key, width=1280, height=720)
            pic_ok = _add_picture_from_url(
                slide,
                fallback,
                visual_x + img_padding,
                content_top_y + img_padding,
                visual_w - Inches(0.2),
                content_panel_h - Inches(0.2),
            )

        if not pic_ok:
            # Soft accent blob when export cannot fetch/rasterize the image.
            ah = theme["accentHex"]
            blob_w = Inches(2.0)
            blob_h = Inches(1.6)
            e1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, visual_x + Inches(0.35), content_top_y + Inches(0.45), blob_w, blob_h)
            e1.fill.solid()
            e1.fill.fore_color.rgb = _rgb(ah)
            e1.fill.transparency = 0.78
            e1.line.fill.background()

        list_w_tb = text_w
    else:
        text_x = card_left + Inches(0.7)
        list_w_tb = card_w - Inches(1.4)

    chart_block_h = Inches(2.25)
    if not _has_explicit_chart_placement(content):
        if _add_chart_snapshot_contain(
            slide,
            content,
            text_x,
            list_y,
            list_w_tb,
            chart_block_h,
            pad_hex=cf,
        ):
            list_y = list_y + chart_block_h + Inches(0.14)

    list_h = (card_bottom - bottom_pad) - list_y - key_h
    bt = slide.shapes.add_textbox(text_x, list_y, list_w_tb, list_h)
    btf = bt.text_frame
    btf.clear()
    _ensure_text_frame_wrap(btf)

    items = bullets[:8]
    bp = max(10, int(round(body_pt_base - (1 if len(items) > 5 else 0))))
    if items:
        for i, line in enumerate(items):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            br = _bullet_row_from_content(content, i)
            _fill_bullet_paragraph(
                para,
                line,
                br,
                marker_rgb=_rgb(bu),
                body_rgb=_rgb(bh),
                font_name=body_font,
                body_pt=float(bp),
                marker_pt=float(max(10, bp - 3)),
                marker_text=marker_text,
                default_bold=body_bold,
            )
            para.line_spacing = float(line_mult)
            para.space_before = Pt(float(para_before_pt))
            para.space_after = Pt(float(para_after_pt))
            para.alignment = PP_ALIGN.LEFT if align_txt == "left" else PP_ALIGN.CENTER if align_txt == "center" else PP_ALIGN.RIGHT

    if key_message.strip():
        km = slide.shapes.add_textbox(card_left + Inches(0.5), card_bottom - Inches(0.95), card_w - Inches(1.0), Inches(0.85))
        km.text_frame.text = key_message.strip()
        _set_textbox_style(
            km,
            font_name=body_font,
            size_pt=max(9, int(round(body_pt_base - 4))),
            color_hex=bh,
            bold=True,
            align=PP_ALIGN.CENTER,
            line_spacing_mult=line_mult,
            space_before_pt=para_before_pt,
            space_after_pt=max(0.0, para_after_pt * 0.5),
        )


def _render_stats_split(
    slide: Any,
    prs: Presentation,
    *,
    title: str,
    bullets: list[str],
    key_message: str,
    content: dict[str, Any],
    theme: ExportTheme,
) -> None:
    """Simplified stats row — two columns of text when not numeric."""
    sw, sh = prs.slide_width, prs.slide_height
    cf = _card_fill_from_content(content, theme["cardFillHex"])
    th = theme["titleHex"]
    bh = theme["bodyHex"]
    heading_font, body_font, body_pt_base, body_bold, lh_px = _gamma_typography_from_content(content, theme)
    line_mult, para_before_pt, para_after_pt, align_txt = _gamma_paragraph_controls_from_content(content)

    _add_rounded_card(slide, Inches(0.5), Inches(0.5), sw - Inches(1.0), sh - Inches(1.0), cf)

    tt = slide.shapes.add_textbox(Inches(0.9), Inches(0.75), sw - Inches(1.8), Inches(0.9))
    ttf = tt.text_frame
    ttf.paragraphs[0].text = title
    ttf.paragraphs[0].font.name = heading_font
    ttf.paragraphs[0].font.size = Pt(max(18, int(round(body_pt_base + 12))))
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = _rgb(th)
    ttf.paragraphs[0].alignment = PP_ALIGN.CENTER

    shift_emu = 0
    ch = Inches(1.9)
    if not _has_explicit_chart_placement(content):
        if _add_chart_snapshot_contain(
            slide,
            content,
            Inches(0.9),
            Inches(1.38),
            sw - Inches(1.8),
            ch,
            pad_hex=cf,
        ):
            shift_emu = int(ch) + int(Inches(0.18))

    body_top = int(Inches(2.2)) + shift_emu

    left = bullets[0] if len(bullets) > 0 else "—"
    right = bullets[1] if len(bullets) > 1 else bullets[0] if bullets else "—"

    b1 = slide.shapes.add_textbox(Inches(0.8), body_top, Inches(5.8), Inches(2.8))
    b1.text_frame.text = left
    _set_textbox_style(
        b1,
        font_name=body_font,
        size_pt=max(10, int(round(body_pt_base - 2))),
        color_hex=bh,
        valign=MSO_ANCHOR.TOP,
        line_spacing_mult=line_mult,
        space_before_pt=para_before_pt,
        space_after_pt=para_after_pt,
        align=PP_ALIGN.LEFT if align_txt == "left" else PP_ALIGN.CENTER if align_txt == "center" else PP_ALIGN.RIGHT,
        bold=body_bold,
    )

    b2 = slide.shapes.add_textbox(Inches(6.75), body_top, Inches(5.8), Inches(2.8))
    b2.text_frame.text = right
    _set_textbox_style(
        b2,
        font_name=body_font,
        size_pt=max(10, int(round(body_pt_base - 2))),
        color_hex=bh,
        valign=MSO_ANCHOR.TOP,
        line_spacing_mult=line_mult,
        space_before_pt=para_before_pt,
        space_after_pt=para_after_pt,
        align=PP_ALIGN.LEFT if align_txt == "left" else PP_ALIGN.CENTER if align_txt == "center" else PP_ALIGN.RIGHT,
        bold=body_bold,
    )

    if key_message.strip():
        km = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), sw - Inches(2.0), Inches(1.2))
        km.text_frame.text = key_message.strip()
        _set_textbox_style(
            km,
            font_name=body_font,
            size_pt=max(9, int(round(body_pt_base - 3))),
            color_hex=bh,
            align=PP_ALIGN.LEFT if align_txt == "left" else PP_ALIGN.CENTER if align_txt == "center" else PP_ALIGN.RIGHT,
            line_spacing_mult=line_mult,
            space_before_pt=para_before_pt,
            space_after_pt=max(0.0, para_after_pt * 0.5),
            bold=body_bold,
        )


def _render_three_cards(
    slide: Any,
    prs: Presentation,
    *,
    title: str,
    bullets: list[str],
    key_message: str,
    content: dict[str, Any],
    theme: ExportTheme,
) -> None:
    sw, sh = prs.slide_width, prs.slide_height
    cf = _card_fill_from_content(content, theme["cardFillHex"])
    th = theme["titleHex"]
    bh = theme["bodyHex"]
    bg_inner = theme["pageBg"]
    heading_font, body_font, body_pt_base, body_bold, lh_px = _gamma_typography_from_content(content, theme)
    line_mult, para_before_pt, para_after_pt, align_txt = _gamma_paragraph_controls_from_content(content)

    _add_rounded_card(slide, Inches(0.5), Inches(0.5), sw - Inches(1.0), sh - Inches(1.0), cf, line_alpha=1.0)

    tt = slide.shapes.add_textbox(Inches(0.9), Inches(0.65), sw - Inches(1.8), Inches(0.85))
    ttf = tt.text_frame
    ttf.paragraphs[0].text = title
    ttf.paragraphs[0].font.name = heading_font
    ttf.paragraphs[0].font.size = Pt(max(18, int(round(body_pt_base + 10))))
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = _rgb(th)
    ttf.paragraphs[0].alignment = PP_ALIGN.CENTER

    items = bullets[:3] if bullets else []
    while len(items) < 3:
        items.append(items[-1] if items else "—")

    slide_w_in = float(int(sw)) / 914400.0
    card_w_in = (slide_w_in - 1.4) / 3.0
    start_x_in = 0.7
    gutter_in = 0.15
    card_y_in = 1.65
    card_h_in = 4.2
    if not _has_explicit_chart_placement(content):
        if _add_chart_snapshot_contain(
            slide,
            content,
            Inches(0.9),
            Inches(1.4),
            sw - Inches(1.8),
            Inches(1.95),
            pad_hex=cf,
        ):
            card_y_in = 3.45
            card_h_in = min(2.75, 5.85 - card_y_in)

    for i in range(3):
        line = items[i] if i < len(items) else "—"
        if ":" in line:
            title_part, rest = line.split(":", 1)
            title_part = title_part.strip()
            body_part = rest.strip()
        else:
            title_part = line[:50]
            body_part = line[50:].strip() if len(line) > 50 else ""

        cx_in = start_x_in + i * (card_w_in + gutter_in)
        cx = Inches(cx_in)
        cy = Inches(card_y_in)
        cw = Inches(card_w_in)
        ch_box = Inches(card_h_in)

        inner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch_box)
        inner.fill.solid()
        inner.fill.fore_color.rgb = _rgb(bg_inner)
        inner.fill.transparency = 0.25
        inner.line.color.rgb = _rgb(th)
        inner.line.transparency = 0.4

        tbox = slide.shapes.add_textbox(
            cx + Inches(0.2),
            cy + Inches(0.25),
            cw - Inches(0.4),
            Inches(0.9),
        )
        tbox.text_frame.text = title_part[:80]
        _set_textbox_style(
            tbox,
            font_name=heading_font,
            size_pt=max(10, int(round(body_pt_base - 3))),
            color_hex=th,
            bold=True,
            valign=MSO_ANCHOR.TOP,
            line_spacing_mult=line_mult,
            space_before_pt=para_before_pt,
            space_after_pt=max(0.0, para_after_pt * 0.5),
        )

        if body_part:
            body_h_emu = max(int(Inches(0.6)), int(ch_box) - int(Inches(1.35)))
            bbox = slide.shapes.add_textbox(
                cx + Inches(0.2),
                cy + Inches(1.15),
                cw - Inches(0.4),
                body_h_emu,
            )
            bbox.text_frame.text = body_part
            _set_textbox_style(
                bbox,
                font_name=body_font,
                size_pt=max(9, int(round(body_pt_base - 5))),
                color_hex=bh,
                valign=MSO_ANCHOR.TOP,
                line_spacing_mult=line_mult,
                space_before_pt=para_before_pt,
                space_after_pt=para_after_pt,
            )

    if key_message.strip():
        km = slide.shapes.add_textbox(Inches(1.0), Inches(6.15), sw - Inches(2.0), Inches(0.9))
        km.text_frame.text = key_message.strip()
        _set_textbox_style(
            km,
            font_name=body_font,
            size_pt=max(9, int(round(body_pt_base - 5))),
            color_hex=bh,
            align=PP_ALIGN.LEFT if align_txt == "left" else PP_ALIGN.CENTER if align_txt == "center" else PP_ALIGN.RIGHT,
            line_spacing_mult=line_mult,
            space_before_pt=para_before_pt,
            space_after_pt=max(0.0, para_after_pt * 0.5),
            bold=body_bold,
        )


def render_gamma_slide(
    prs: Presentation,
    *,
    slide_content: dict[str, Any],
    slide_title: str,
    theme: ExportTheme,
    speaker_notes: str,
) -> None:
    """Append one slide to `prs` using Gamma-style layout + theme."""
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)

    _add_full_bleed_rect(slide, prs, theme["pageBg"])

    content = slide_content if isinstance(slide_content, dict) else {}
    title = str(content.get("title") or slide_title or "Slide")
    bullets_raw = content.get("bullets")
    bullets = [str(b).strip() for b in bullets_raw] if isinstance(bullets_raw, list) else []
    subtitle = str(content.get("subtitle") or "")
    km = str(content.get("keyMessage") or "").strip()
    hl = str(content.get("highlight") or "").strip()
    key_message = km or hl

    preset = resolve_gamma_preset_for_export(content)

    if preset == "section_only":
        _render_section_only(
            slide,
            prs,
            title=title,
            subtitle=subtitle,
            key_message=key_message,
            theme=theme,
            content=content,
        )
    elif preset in ("hero_split", "two_column"):
        # Match deck preview: no subtitle row under the title or footer strip — more room for bullets + chart.
        _render_hero_split(
            slide,
            prs,
            title=title,
            bullets=bullets,
            subtitle="",
            key_message="",
            content=content,
            theme=theme,
        )
    elif preset == "stats_split":
        _render_stats_split(
            slide,
            prs,
            title=title,
            bullets=bullets,
            key_message=key_message,
            content=content,
            theme=theme,
        )
    elif preset == "three_cards":
        _render_three_cards(
            slide,
            prs,
            title=title,
            bullets=bullets,
            key_message=key_message,
            content=content,
            theme=theme,
        )
    else:
        # title_bullets or unknown
        _render_title_bullets(
            slide,
            prs,
            title=title,
            bullets=bullets,
            subtitle=subtitle,
            key_message=key_message,
            content=content,
            theme=theme,
        )

    cf = _card_fill_from_content(content, theme["cardFillHex"])
    place_box = _chart_placement_emu_box(prs, content)
    if place_box is not None and _chart_snapshot_url(content):
        lx, ty, ww, hh = place_box
        _add_chart_snapshot_contain(slide, content, lx, ty, ww, hh, pad_hex=cf)

    notes = str(content.get("speakerNotes") or speaker_notes or "").strip()
    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass


def build_gamma_presentation_pptx(
    *,
    slides: list[dict[str, Any]],
    template_name: str | None,
) -> Presentation:
    """Build a full `Presentation` with themed Gamma-style slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = resolve_export_theme(template_name)

    for s in slides:
        c = s.get("content") if isinstance(s.get("content"), dict) else {}
        stitle = str(s.get("title") or "")
        render_gamma_slide(prs, slide_content=c, slide_title=stitle, theme=theme, speaker_notes="")

    return prs
